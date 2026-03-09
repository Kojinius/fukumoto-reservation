"""
福元鍼灸整骨院 オンライン予約システム 詳細設計書
PowerPoint 生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラー定義 ──
C_BROWN  = RGBColor(0x73, 0x57, 0x63)
C_ORANGE = RGBColor(0xF7, 0x93, 0x21)
C_BEIGE  = RGBColor(0xFC, 0xF0, 0xDE)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x4A, 0x30, 0x40)
C_MUTED  = RGBColor(0x9B, 0x7E, 0x8E)
C_BEIGE2 = RGBColor(0xF5, 0xE8, 0xD8)
C_BEIGE3 = RGBColor(0xF5, 0xF0, 0xED)
C_LIGHT  = RGBColor(0xF0, 0xEB, 0xF0)
C_GREEN  = RGBColor(0x3A, 0x8A, 0x5C)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def add_shape(slide, x, y, w, h, fill=None, border=None, border_w=1):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    return s


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True,
             line_sp=None, font="Noto Sans JP"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        if color:
            run.font.color.rgb = color
    return tb


def header_bar(slide, title, badge=None, badge_bg=C_ORANGE):
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.9), fill=C_BROWN)
    add_text(slide, title, Inches(0.4), Inches(0.08),
             Inches(10.5), Inches(0.75), size=26, bold=True, color=C_WHITE)
    if badge:
        b = add_shape(slide, Inches(11.5), Inches(0.2),
                      Inches(1.6), Inches(0.5), fill=badge_bg)
        tf = b.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = badge
        r.font.size = Pt(12); r.font.bold = True
        r.font.name = "Noto Sans JP"; r.font.color.rgb = C_WHITE


def section_label(slide, text, x, y, w=Inches(3.0), bg=C_BROWN, fg=C_WHITE):
    s = add_shape(slide, x, y, w, Inches(0.38), fill=bg)
    tf = s.text_frame; p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True
    r.font.name = "Noto Sans JP"; r.font.color.rgb = fg


def bullet_block(slide, items, x, y, w, size=13, color=C_DARK, gap=0.3):
    cy = y
    for item in items:
        add_text(slide, item, x, cy, w, Inches(gap), size=size, color=color)
        cy += Inches(gap)
    return cy


def divider(slide, y, x=Inches(0.4), w=None):
    if w is None: w = W - Inches(0.8)
    add_shape(slide, x, y, w, Inches(0.02), fill=C_BEIGE2)


# ═══════════════════════════════════════════════
# Slide 1 : 表紙
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_BROWN)
add_shape(s, Inches(0), Inches(5.8), W, Inches(1.7), fill=RGBColor(0x5a, 0x40, 0x50))
add_shape(s, Inches(1.0), Inches(2.1), Inches(0.8), Inches(0.07), fill=C_ORANGE)
add_text(s, "福元鍼灸整骨院\nオンライン予約システム",
         Inches(1.0), Inches(2.2), Inches(11), Inches(2.3),
         size=42, bold=True, color=C_WHITE, line_sp=Pt(52))
add_text(s, "詳細設計書",
         Inches(1.0), Inches(4.6), Inches(9), Inches(0.7),
         size=26, color=C_BEIGE)
add_text(s, "システム構成 / DB設計 / セキュリティ / API仕様 / バリデーション",
         Inches(1.0), Inches(5.25), Inches(10), Inches(0.45),
         size=14, color=C_ORANGE)
add_text(s, "ver 1.0　|　2026-03-09",
         Inches(1.0), Inches(6.1), Inches(8), Inches(0.4),
         size=13, color=C_MUTED)


# ═══════════════════════════════════════════════
# Slide 2 : 目次
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_BEIGE)
header_bar(s, "目次　— 詳細設計書")

toc_items = [
    ("1", "システム概要・構成図",       "技術スタック・構成要素"),
    ("2", "Firestore データ設計",        "コレクション・フィールド定義"),
    ("3", "認証・認可設計",              "Firebase Auth・admin クレーム"),
    ("4", "パスワード・アカウントポリシー", "パスワードルール・セッション管理"),
    ("5", "Firestore セキュリティルール", "アクセス制御・バリデーション"),
    ("6", "入力バリデーション一覧",      "クライアント側・Firestore ルール"),
    ("7", "Cloud Functions API 仕様",   "エンドポイント・パラメータ定義"),
    ("8", "メール通知設計",              "送信フロー・メールテンプレート"),
    ("9", "ディレクトリ構成",            "ファイル構成・役割"),
]

col_w = Inches(6.0)
col_h = Inches(5.9)
add_shape(s, Inches(0.2), Inches(1.0), col_w, col_h, fill=C_WHITE)
add_shape(s, Inches(6.8), Inches(1.0), Inches(6.3), col_h, fill=C_BEIGE2)

items_left  = toc_items[:5]
items_right = toc_items[5:]

for i, (num, title, sub) in enumerate(items_left):
    y = Inches(1.15) + i * Inches(1.1)
    add_shape(s, Inches(0.35), y, Inches(0.45), Inches(0.45), fill=C_BROWN)
    tf = s.shapes[-1].text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(16)
    r.font.bold = True; r.font.name = "Noto Sans JP"; r.font.color.rgb = C_WHITE
    add_text(s, title, Inches(0.9), y + Inches(0.01),
             Inches(5.0), Inches(0.4), size=15, bold=True, color=C_DARK)
    add_text(s, sub, Inches(0.9), y + Inches(0.42),
             Inches(5.0), Inches(0.3), size=11, color=C_MUTED)

for i, (num, title, sub) in enumerate(items_right):
    y = Inches(1.15) + i * Inches(1.3)
    add_shape(s, Inches(6.95), y, Inches(0.45), Inches(0.45), fill=C_ORANGE)
    tf = s.shapes[-1].text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(16)
    r.font.bold = True; r.font.name = "Noto Sans JP"; r.font.color.rgb = C_WHITE
    add_text(s, title, Inches(7.5), y + Inches(0.01),
             Inches(5.3), Inches(0.4), size=15, bold=True, color=C_DARK)
    add_text(s, sub, Inches(7.5), y + Inches(0.42),
             Inches(5.3), Inches(0.3), size=11, color=C_MUTED)


# ═══════════════════════════════════════════════
# Slide 3 : システム概要・構成図
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "1　システム概要・構成図", "技術スタック")

# 技術スタック左カラム
add_shape(s, Inches(0.2), Inches(1.0), Inches(4.0), Inches(6.1), fill=C_BEIGE)
add_text(s, "技術スタック", Inches(0.4), Inches(1.1),
         Inches(3.6), Inches(0.4), size=15, bold=True, color=C_BROWN)
stack = [
    ("フロントエンド", "HTML / CSS / Vanilla JS (ES Modules)"),
    ("データベース",   "Cloud Firestore（NoSQL）"),
    ("認証",           "Firebase Authentication"),
    ("ホスティング",   "Firebase Hosting（kojinius.jp）"),
    ("サーバーレス",   "Cloud Functions for Firebase（Node.js）"),
    ("メール配信",     "Resend API（SPF/DKIM 設定済み）"),
    ("PDFエンジン",    "pdf-lib + fontkit（クライアント生成）"),
    ("シークレット管理", "Firebase Secret Manager（APIキー）"),
    ("バージョン管理", "Git / GitHub（Private）"),
]
sy = Inches(1.55)
for label, val in stack:
    add_text(s, label, Inches(0.4), sy, Inches(1.5), Inches(0.28),
             size=11, bold=True, color=C_BROWN)
    add_text(s, val, Inches(1.95), sy, Inches(2.1), Inches(0.28),
             size=11, color=C_DARK)
    sy += Inches(0.6)

# 構成図（中央）
add_shape(s, Inches(4.4), Inches(1.0), Inches(8.7), Inches(6.1), fill=C_BEIGE2)
add_text(s, "システム構成図", Inches(4.6), Inches(1.1),
         Inches(8.0), Inches(0.4), size=15, bold=True, color=C_BROWN)

# ブラウザ
add_shape(s, Inches(4.6), Inches(1.6), Inches(2.5), Inches(1.8), fill=C_WHITE, border=C_MUTED)
add_text(s, "🌐 ブラウザ\n（患者・スタッフ）",
         Inches(4.7), Inches(1.7), Inches(2.3), Inches(0.8), size=12, bold=True, color=C_DARK)
add_text(s, "HTML / CSS / JS",
         Inches(4.7), Inches(2.5), Inches(2.3), Inches(0.3), size=10, color=C_MUTED)

# 矢印
add_text(s, "↔", Inches(7.3), Inches(2.2), Inches(0.5), Inches(0.4),
         size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Firebase Hosting
add_shape(s, Inches(7.9), Inches(1.6), Inches(2.2), Inches(1.8), fill=C_BROWN, border=C_BROWN)
add_text(s, "🔥 Firebase\nHosting",
         Inches(8.0), Inches(1.7), Inches(2.0), Inches(0.8), size=12, bold=True, color=C_ORANGE)
add_text(s, "kojinius.jp",
         Inches(8.0), Inches(2.5), Inches(2.0), Inches(0.3), size=10, color=C_BEIGE)

# 矢印下
add_text(s, "↕", Inches(8.9), Inches(3.5), Inches(0.4), Inches(0.4),
         size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Firestore
add_shape(s, Inches(7.9), Inches(3.9), Inches(2.2), Inches(1.2), fill=C_WHITE, border=C_BROWN)
add_text(s, "🗄️ Firestore",
         Inches(8.0), Inches(4.0), Inches(2.0), Inches(0.4), size=12, bold=True, color=C_BROWN)
add_text(s, "reservations\nslots / settings",
         Inches(8.0), Inches(4.42), Inches(2.0), Inches(0.55), size=10, color=C_DARK)

# Auth
add_shape(s, Inches(7.9), Inches(5.3), Inches(2.2), Inches(0.9), fill=C_WHITE, border=C_BROWN)
add_text(s, "🔐 Auth",
         Inches(8.0), Inches(5.35), Inches(2.0), Inches(0.35), size=12, bold=True, color=C_BROWN)
add_text(s, "admin クレーム",
         Inches(8.0), Inches(5.7), Inches(2.0), Inches(0.3), size=10, color=C_DARK)

# Functions → Resend
add_shape(s, Inches(4.6), Inches(3.9), Inches(2.5), Inches(1.2), fill=C_WHITE, border=C_MUTED)
add_text(s, "⚡ Cloud\nFunctions",
         Inches(4.7), Inches(4.0), Inches(2.3), Inches(0.65), size=12, bold=True, color=C_DARK)
add_text(s, "sendMail / scheduler",
         Inches(4.7), Inches(4.65), Inches(2.3), Inches(0.3), size=9, color=C_MUTED)

add_text(s, "→", Inches(7.3), Inches(4.3), Inches(0.5), Inches(0.4),
         size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

add_shape(s, Inches(10.3), Inches(3.9), Inches(2.5), Inches(1.2), fill=C_WHITE, border=C_MUTED)
add_text(s, "📧 Resend\nAPI",
         Inches(10.4), Inches(4.0), Inches(2.3), Inches(0.65), size=12, bold=True, color=C_DARK)
add_text(s, "kojinius.jp 送信ドメイン",
         Inches(10.4), Inches(4.65), Inches(2.3), Inches(0.3), size=9, color=C_MUTED)

add_text(s, "↕", Inches(5.7), Inches(3.5), Inches(0.4), Inches(0.4),
         size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Slide 4 : Firestore データ設計
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "2　Firestore データ設計", "DB Schema")

cols_def = [
    ("reservations/{uuid}", C_BROWN, C_ORANGE, [
        ("id",            "string(36)",  "UUID v4 予約番号"),
        ("name",          "string≤50",   "患者氏名（必須）"),
        ("furigana",      "string≤50",   "ふりがな（必須）"),
        ("phone",         "string≤20",   "電話番号（必須・半角）"),
        ("email",         "string≤200",  "メールアドレス（任意）"),
        ("birthdate",     "string≤10",   "生年月日 YYYY-MM-DD（任意）"),
        ("gender",        "string≤20",   "性別（任意）"),
        ("zip",           "string≤8",    "郵便番号（任意）"),
        ("address",       "string≤200",  "住所（任意）"),
        ("date",          "string=10",   "予約日 YYYY-MM-DD（必須）"),
        ("time",          "string≤10",   "予約時刻 H:MM（必須）"),
        ("visitType",     "string≤50",   "初診/再診（任意）"),
        ("insurance",     "string≤50",   "保険有無（任意）"),
        ("symptoms",      "string≤1000", "症状（任意）"),
        ("notes",         "string≤500",  "伝達事項（任意）"),
        ("contactMethod", "string≤50",   "連絡方法（任意）"),
        ("status",        "enum",        "pending / confirmed / cancelled"),
        ("createdAt",     "timestamp",   "予約作成日時"),
    ]),
    ("slots/{date_time}", C_BLUE, C_WHITE, [
        ("date",   "string=10", "予約日 YYYY-MM-DD"),
        ("time",   "string≤10", "予約時刻 H:MM"),
        ("status", "enum",      "pending / cancelled"),
    ]),
    ("settings/clinic", C_GREEN, C_WHITE, [
        ("clinicName",    "string",  "院名"),
        ("phone",         "string",  "電話番号"),
        ("clinicAddress", "string",  "住所"),
        ("clinicZip",     "string",  "郵便番号"),
        ("clinicUrl",     "string",  "HP URL"),
        ("clinicLogo",    "string",  "ロゴ（base64 PNG）"),
        ("colorTheme",    "string",  "カラーテーマ"),
        ("businessHours", "map",     "営業時間設定"),
        ("holidays",      "array",   "休日日付リスト"),
        ("announcement",  "map",     "お知らせバナー設定"),
        ("maintenance",   "map",     "メンテナンス期間設定"),
        ("bookingCutoffMinutes", "number", "予約締め切り分前"),
        ("cancelCutoffMinutes",  "number", "キャンセル締め切り分前"),
        ("privacyPolicy", "string",  "プライバシーポリシー本文"),
    ]),
]

x_positions = [Inches(0.2), Inches(4.6), Inches(9.0)]
col_w = Inches(4.2)

for ci, (coll_name, bg_color, title_color, fields) in enumerate(cols_def):
    cx = x_positions[ci]
    add_shape(s, cx, Inches(1.0), col_w, Inches(6.1), fill=C_BEIGE)
    # コレクション名
    b = add_shape(s, cx, Inches(1.0), col_w, Inches(0.5), fill=bg_color)
    tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = coll_name
    r.font.size = Pt(13); r.font.bold = True
    r.font.name = "Noto Sans JP"; r.font.color.rgb = title_color

    fy = Inches(1.55)
    for field, ftype, desc in fields:
        if fy > Inches(6.8): break
        add_text(s, field, cx + Inches(0.12), fy,
                 Inches(1.3), Inches(0.27), size=10, bold=True, color=bg_color)
        add_text(s, ftype, cx + Inches(1.44), fy,
                 Inches(1.0), Inches(0.27), size=9, color=C_MUTED)
        add_text(s, desc, cx + Inches(2.46), fy,
                 Inches(1.6), Inches(0.27), size=9, color=C_DARK)
        fy += Inches(0.34)


# ═══════════════════════════════════════════════
# Slide 5 : 認証・認可設計
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "3　認証・認可設計", "Firebase Auth")

# 左：認証フロー
add_shape(s, Inches(0.2), Inches(1.0), Inches(6.0), Inches(6.1), fill=C_BEIGE)
add_text(s, "ログインフロー", Inches(0.4), Inches(1.1),
         Inches(5.6), Inches(0.4), size=15, bold=True, color=C_BROWN)

flow_steps = [
    ("1", "メール・パスワード入力",  "スタッフがログインフォームに入力"),
    ("2", "setPersistence(browser)", "セッション限定（タブを閉じるとログアウト）"),
    ("3", "signInWithEmailAndPassword", "Firebase Auth で認証"),
    ("4", "getIdTokenResult()",      "カスタムクレームを取得"),
    ("5", "claims.admin == true ?",  "admin クレームを確認"),
    ("6", "YES → 管理画面表示",       "admin.html をロード"),
    ("6", "NO  → ログアウト",         "signOut() → login.html"),
]

fy = Inches(1.6)
for num, title, sub in flow_steps:
    circle = add_shape(s, Inches(0.4), fy, Inches(0.42), Inches(0.42), fill=C_BROWN)
    tf = circle.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(13)
    r.font.bold = True; r.font.name = "Noto Sans JP"; r.font.color.rgb = C_WHITE
    add_text(s, title, Inches(0.95), fy,
             Inches(5.0), Inches(0.3), size=13, bold=True, color=C_DARK)
    add_text(s, sub, Inches(0.95), fy + Inches(0.28),
             Inches(5.0), Inches(0.25), size=10, color=C_MUTED)
    fy += Inches(0.77)

# 右：認可設計
add_shape(s, Inches(6.5), Inches(1.0), Inches(6.6), Inches(6.1), fill=C_BEIGE2)
add_text(s, "認可マトリクス", Inches(6.7), Inches(1.1),
         Inches(6.2), Inches(0.4), size=15, bold=True, color=C_BROWN)

# テーブルヘッダー
add_shape(s, Inches(6.5), Inches(1.55), Inches(6.6), Inches(0.4), fill=C_BROWN)
headers = ["操作", "未認証（患者）", "認証済み管理者"]
hx = [Inches(6.55), Inches(8.5), Inches(10.9)]
for htext, hpos in zip(headers, hx):
    add_text(s, htext, hpos, Inches(1.6), Inches(2.3), Inches(0.3),
             size=11, bold=True, color=C_WHITE)

rows = [
    ("reservations 新規作成",   "✅ 可",  "✅ 可"),
    ("reservations 一覧取得",   "❌ 不可","✅ 可"),
    ("reservations 1件取得",    "✅ 可",  "✅ 可"),
    ("reservations キャンセル", "✅ 可*", "✅ 可"),
    ("reservations 削除",       "❌ 不可","✅ 可"),
    ("slots 読み取り/作成",     "✅ 可",  "✅ 可"),
    ("slots ステータス更新",    "✅ 可*", "✅ 可"),
    ("settings 読み取り",       "✅ 可",  "✅ 可"),
    ("settings 書き込み",       "❌ 不可","✅ 可"),
]

ry = Inches(2.0)
for i, (op, anon, admin) in enumerate(rows):
    bg = C_WHITE if i % 2 == 0 else C_BEIGE
    add_shape(s, Inches(6.5), ry, Inches(6.6), Inches(0.42), fill=bg)
    add_text(s, op, Inches(6.55), ry + Inches(0.06),
             Inches(1.9), Inches(0.3), size=10, color=C_DARK)
    gc = C_GREEN if "✅" in anon else C_RED
    add_text(s, anon, Inches(8.5), ry + Inches(0.06),
             Inches(2.3), Inches(0.3), size=10, bold=True, color=gc)
    add_text(s, admin, Inches(10.9), ry + Inches(0.06),
             Inches(2.0), Inches(0.3), size=10, bold=True, color=C_GREEN)
    ry += Inches(0.42)

add_text(s, "* 条件付き（電話番号一致 / _cancelVerify フィールド検証）",
         Inches(6.55), ry + Inches(0.05), Inches(6.3), Inches(0.35),
         size=10, color=C_MUTED)


# ═══════════════════════════════════════════════
# Slide 6 : パスワード・アカウントポリシー
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "4　パスワード・アカウントポリシー", "Security")

# 左上：パスワードルール
add_shape(s, Inches(0.2), Inches(1.0), Inches(6.3), Inches(2.8), fill=C_BEIGE)
section_label(s, "🔑 パスワードルール", Inches(0.4), Inches(1.1), Inches(3.2))
pw_rules = [
    ("最低文字数",     "8文字以上（フロントエンドバリデーション）"),
    ("文字種制限",     "Firebase Auth デフォルト（制限なし）"),
    ("確認入力",       "「新しいパスワード」「確認用」の2入力で一致チェック"),
    ("変更時再認証",   "現在のパスワードによる reauthenticateWithCredential 必須"),
    ("変更API",        "Firebase updatePassword()（セキュアな変更）"),
]
ry = Inches(1.56)
for label, val in pw_rules:
    add_text(s, label, Inches(0.4), ry, Inches(1.8), Inches(0.3),
             size=11, bold=True, color=C_BROWN)
    add_text(s, val, Inches(2.3), ry, Inches(4.1), Inches(0.3), size=11, color=C_DARK)
    ry += Inches(0.42)

# 右上：セッション管理
add_shape(s, Inches(6.7), Inches(1.0), Inches(6.4), Inches(2.8), fill=C_BEIGE2)
section_label(s, "🔐 セッション・ログイン管理", Inches(6.9), Inches(1.1), Inches(3.6))
sess_rules = [
    ("セッション種別",  "browserSessionPersistence（タブ閉じで自動ログアウト）"),
    ("ログイン失敗",    "Firebase Auth が自動的にレートリミット（too-many-requests）"),
    ("エラーメッセージ", "「メールアドレスまたはパスワードが正しくありません」\n（ユーザー列挙攻撃対策のため同一メッセージ）"),
    ("管理者確認",      "ログイン後に claims.admin == true を必須チェック"),
    ("ガード方式",      "全ページで requireAdmin() を先頭に配置、未認証は即リダイレクト"),
]
ry = Inches(1.56)
for label, val in sess_rules:
    add_text(s, label, Inches(6.9), ry, Inches(1.9), Inches(0.3),
             size=11, bold=True, color=C_BROWN)
    add_text(s, val, Inches(8.9), ry, Inches(4.0), Inches(0.35), size=11, color=C_DARK)
    ry += Inches(0.44)

# 下：メールアドレス変更フロー
add_shape(s, Inches(0.2), Inches(4.0), Inches(12.9), Inches(3.1), fill=C_BEIGE)
section_label(s, "📧 メールアドレス変更フロー（Firebase verifyBeforeUpdateEmail）",
              Inches(0.4), Inches(4.1), Inches(7.5))

steps_email = [
    ("1. 再認証",        "現在のパスワードで\nreauthenticate()"),
    ("2. 確認メール送信", "新しいアドレスに\nverifyBeforeUpdateEmail()"),
    ("3. リンクをクリック", "メール内のリンクを\nクリック"),
    ("4. 変更完了",       "Firebase が新アドレスに\n自動更新"),
    ("5. 旧アドレス無効", "旧アドレスでは\nログイン不可になる"),
]
sx = Inches(0.4)
for i, (title, desc) in enumerate(steps_email):
    cx = sx + i * Inches(2.5)
    add_shape(s, cx, Inches(4.55), Inches(2.2), Inches(2.2), fill=C_WHITE, border=C_MUTED)
    add_shape(s, cx, Inches(4.55), Inches(2.2), Inches(0.5), fill=C_BROWN)
    tf = s.shapes[-1].text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(11)
    r.font.bold = True; r.font.name = "Noto Sans JP"; r.font.color.rgb = C_WHITE
    add_text(s, desc, cx + Inches(0.1), Inches(5.1),
             Inches(2.0), Inches(0.9), size=11, color=C_DARK, line_sp=Pt(15))
    if i < len(steps_email) - 1:
        add_text(s, "→", cx + Inches(2.25), Inches(5.0),
                 Inches(0.3), Inches(0.5), size=18, bold=True,
                 color=C_ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Slide 7 : Firestore セキュリティルール
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "5　Firestore セキュリティルール", "firestore.rules")

add_shape(s, Inches(0.2), Inches(1.0), Inches(12.9), Inches(6.1), fill=RGBColor(0x28, 0x2C, 0x34))

code_lines = [
    ("// デフォルト拒否",                                     C_MUTED),
    ("match /{document=**} { allow read, write: if false; }","#abb2bf"),
    ("",                                                      "#abb2bf"),
    ("// 管理者判定",                                          C_MUTED),
    ("function isAdmin() {",                                  "#abb2bf"),
    ("  return request.auth != null && request.auth.token.admin == true;", "#abb2bf"),
    ("}",                                                     "#abb2bf"),
    ("",                                                      "#abb2bf"),
    ("// reservations",                                       C_MUTED),
    ("allow create: if request.auth == null && isValidReservation(data);", "#e06c75"),
    ("allow get:    if request.auth == null;                 // ID=トークン", "#56b6c2"),
    ("allow update: if request.auth == null                 // キャンセルのみ", "#56b6c2"),
    ("             && data.status != 'cancelled'",           "#56b6c2"),
    ("             && newData.status == 'cancelled'",        "#56b6c2"),
    ("             && newData._cancelVerify == data.phone;  // 電話番号検証", "#e5c07b"),
    ("allow read, write: if isAdmin();",                     "#98c379"),
    ("",                                                      "#abb2bf"),
    ("// slots: 二重予約防止用",                               C_MUTED),
    ("allow read, create: if request.auth == null;",         "#abb2bf"),
    ("allow update: if request.auth == null                 // キャンセル⇄再予約", "#56b6c2"),
    ("",                                                      "#abb2bf"),
    ("// settings: 未認証でも読み取り可（院名・テーマ表示用）", C_MUTED),
    ("allow read: if true;  allow write: if isAdmin();",     "#98c379"),
]

def hex_to_rgb(hex_str):
    if isinstance(hex_str, RGBColor): return hex_str
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

ly = Inches(1.1)
for line_text, line_color in code_lines:
    add_text(s, line_text, Inches(0.4), ly, Inches(12.5), Inches(0.27),
             size=10, color=hex_to_rgb(line_color), font="Courier New")
    ly += Inches(0.27)


# ═══════════════════════════════════════════════
# Slide 8 : 入力バリデーション一覧
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "6　入力バリデーション一覧", "Validation")

# テーブルヘッダー
add_shape(s, Inches(0.2), Inches(1.0), Inches(12.9), Inches(0.48), fill=C_BROWN)
headers = ["フィールド", "必須", "形式・制約", "クライアント", "Firestoreルール"]
hxw = [(Inches(0.25), Inches(2.0)), (Inches(2.3), Inches(0.7)),
       (Inches(3.1), Inches(3.3)), (Inches(6.5), Inches(3.0)),
       (Inches(9.6), Inches(3.4))]
for (hx, hw), ht in zip(hxw, headers):
    add_text(s, ht, hx, Inches(1.06), hw, Inches(0.34),
             size=11, bold=True, color=C_WHITE)

val_rows = [
    ("氏名（name）",       "必須", "最大50文字",          "未入力チェック",         "string & ≤50"),
    ("ふりがな",           "必須", "最大50文字",          "未入力チェック",         "string & ≤50"),
    ("電話番号",           "必須", "最大20文字・半角正規化","未入力＋toHankaku()",    "string & ≤20"),
    ("メールアドレス",     "任意", "最大200文字・形式チェック","形式チェック（正規表現）","string & ≤200（任意）"),
    ("生年月日",           "任意", "YYYY-MM-DD 形式",      "日付フォーマット確認",   "string & ≤10（任意）"),
    ("性別",               "任意", "最大20文字",          "ラジオボタン選択",       "string & ≤20（任意）"),
    ("郵便番号",           "必須", "最大8文字・半角",      "zipcloud API 入力補助",  "string & ≤8（任意）"),
    ("住所",               "必須", "最大200文字",         "zipcloud 自動補完",      "string & ≤200（任意）"),
    ("予約日",             "必須", "YYYY-MM-DD（=10文字）", "カレンダー選択のみ",     "string & =10"),
    ("予約時刻",           "必須", "H:MM（≤10文字）",      "スロット選択のみ",       "string & ≤10"),
    ("症状",               "任意", "最大1000文字",        "文字数カウント",         "string & ≤1000（任意）"),
    ("伝達事項",           "任意", "最大500文字",         "文字数カウント",         "string & ≤500（任意）"),
    ("ステータス",         "必須", "enum 3値",            "システム管理",           "in [pending,confirmed,cancelled]"),
    ("個人情報同意",       "必須", "チェック必須",         "checkboxチェック確認",   "（クライアント側のみ）"),
    ("HP URL（設定画面）", "任意", "URL形式",              "new URL() でパース検証",  "（設定はadminのみ）"),
]

ry = Inches(1.48)
for i, (field, req, rule, client, firestore) in enumerate(val_rows):
    bg = C_WHITE if i % 2 == 0 else C_BEIGE
    add_shape(s, Inches(0.2), ry, Inches(12.9), Inches(0.36), fill=bg)
    row_vals = [field, req, rule, client, firestore]
    for (hx, hw), val in zip(hxw, row_vals):
        col = C_RED if val == "必須" else (C_GREEN if val == "任意" else C_DARK)
        add_text(s, val, hx, ry + Inches(0.04), hw, Inches(0.3),
                 size=10, color=col, bold=(val in ("必須","任意")))
    ry += Inches(0.36)


# ═══════════════════════════════════════════════
# Slide 9 : Cloud Functions API 仕様
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "7　Cloud Functions API 仕様", "HTTP Functions")

apis = [
    {
        "name": "sendReservationEmail",
        "method": "POST",
        "color": C_BLUE,
        "endpoint": "https://sendreservationemail-po3aztuimq-uc.a.run.app",
        "trigger": "患者が予約確定時にクライアントから直接呼び出し",
        "params": [
            ("to",    "必須", "患者メールアドレス"),
            ("name",  "必須", "患者氏名"),
            ("date",  "必須", "予約日 YYYY-MM-DD"),
            ("time",  "必須", "予約時刻"),
            ("menu",  "必須", "診療メニュー"),
            ("id",    "任意", "予約番号（UUID）"),
        ],
        "cors": "https://kojinius.jp のみ",
        "secret": "RESEND_API_KEY（Secret Manager）",
    },
    {
        "name": "notifyAdminOnReservation",
        "method": "POST",
        "color": C_ORANGE,
        "endpoint": "https://notifyadminonreservation-po3aztuimq-uc.a.run.app",
        "trigger": "患者が予約確定時にクライアントから呼び出し（非同期・fire-and-forget）",
        "params": [
            ("name",      "必須", "患者氏名"),
            ("date",      "必須", "予約日"),
            ("time",      "必須", "予約時刻"),
            ("phone",     "必須", "電話番号"),
            ("visitType", "任意", "初診/再診"),
            ("symptoms",  "任意", "症状"),
        ],
        "cors": "https://kojinius.jp のみ",
        "secret": "RESEND_API_KEY（Secret Manager）",
    },
    {
        "name": "sendDailyReminders",
        "method": "Scheduler",
        "color": C_GREEN,
        "endpoint": "毎日 15:00 JST 自動実行（Cloud Scheduler）",
        "trigger": "翌日の予約者全員にリマインダーメールを自動送信",
        "params": [
            ("（パラメータなし）", "—", "Firestoreから翌日予約を自動取得"),
            ("対象",              "—", "status = pending or confirmed"),
            ("送信先",            "—", "各予約の email フィールド（空の場合はスキップ）"),
        ],
        "cors": "不要（Scheduler トリガー）",
        "secret": "RESEND_API_KEY（Secret Manager）",
    },
]

col_w2 = Inches(4.2)
for ai, api in enumerate(apis):
    cx = Inches(0.2 + ai * 4.4)
    add_shape(s, cx, Inches(1.0), col_w2, Inches(6.1), fill=C_BEIGE)
    add_shape(s, cx, Inches(1.0), col_w2, Inches(0.55), fill=api["color"])

    # メソッドバッジ
    mb = add_shape(s, cx + Inches(0.1), Inches(1.08), Inches(0.9), Inches(0.35),
                   fill=C_WHITE)
    tf = mb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = api["method"]; r.font.size = Pt(11)
    r.font.bold = True; r.font.name = "Noto Sans JP"; r.font.color.rgb = api["color"]

    add_text(s, api["name"], cx + Inches(1.1), Inches(1.1),
             col_w2 - Inches(1.2), Inches(0.38), size=12, bold=True, color=C_WHITE)

    fy2 = Inches(1.65)
    add_text(s, "エンドポイント", cx + Inches(0.1), fy2, col_w2 - Inches(0.2), Inches(0.27),
             size=9, bold=True, color=C_BROWN)
    fy2 += Inches(0.27)
    add_text(s, api["endpoint"], cx + Inches(0.1), fy2, col_w2 - Inches(0.2), Inches(0.27),
             size=9, color=C_DARK)
    fy2 += Inches(0.35)
    add_text(s, "トリガー", cx + Inches(0.1), fy2, col_w2 - Inches(0.2), Inches(0.27),
             size=9, bold=True, color=C_BROWN)
    fy2 += Inches(0.27)
    add_text(s, api["trigger"], cx + Inches(0.1), fy2, col_w2 - Inches(0.2), Inches(0.35),
             size=9, color=C_DARK, line_sp=Pt(13))
    fy2 += Inches(0.45)

    add_text(s, "パラメータ", cx + Inches(0.1), fy2, col_w2 - Inches(0.2), Inches(0.27),
             size=9, bold=True, color=C_BROWN)
    fy2 += Inches(0.3)
    for pname, preq, pdesc in api["params"]:
        add_text(s, f"• {pname} ({preq}) : {pdesc}",
                 cx + Inches(0.1), fy2, col_w2 - Inches(0.2), Inches(0.28),
                 size=9, color=C_DARK)
        fy2 += Inches(0.28)

    fy2 += Inches(0.1)
    add_text(s, f"CORS: {api['cors']}", cx + Inches(0.1), fy2,
             col_w2 - Inches(0.2), Inches(0.27), size=9, color=C_MUTED)
    fy2 += Inches(0.27)
    add_text(s, f"Secret: {api['secret']}", cx + Inches(0.1), fy2,
             col_w2 - Inches(0.2), Inches(0.27), size=9, color=C_MUTED)


# ═══════════════════════════════════════════════
# Slide 10 : メール通知設計
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "8　メール通知設計", "Email Flow")

# 送信フロー図
add_shape(s, Inches(0.2), Inches(1.0), Inches(12.9), Inches(2.5), fill=C_BEIGE)
add_text(s, "メール送信フロー", Inches(0.4), Inches(1.1),
         Inches(8), Inches(0.4), size=14, bold=True, color=C_BROWN)

flow_nodes = [
    ("患者が予約確定",  Inches(0.4)),
    ("app.js が\nFunctions呼び出し", Inches(2.7)),
    ("Functions が\nFirestore参照",  Inches(5.2)),
    ("Resend API\n経由で送信",       Inches(7.7)),
    ("患者メール\n受信",             Inches(10.3)),
]
for label, nx in flow_nodes:
    add_shape(s, nx, Inches(1.55), Inches(2.1), Inches(0.9), fill=C_BROWN)
    tf = s.shapes[-1].text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(10)
    r.font.bold = True; r.font.name = "Noto Sans JP"; r.font.color.rgb = C_WHITE

for i in range(len(flow_nodes) - 1):
    nx = flow_nodes[i][1] + Inches(2.15)
    add_text(s, "→", nx, Inches(1.85), Inches(0.45), Inches(0.4),
             size=18, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# メール種別カード
mail_types = [
    ("📩 予約確認メール", C_BLUE, [
        "送信先：患者（予約フォームのメールアドレス）",
        "送信タイミング：予約確定直後",
        "送信元：noreply@kojinius.jp",
        "件名：【院名】ご予約確認",
        "本文：予約番号・日時・診療メニュー・キャンセル方法",
        "追記：院の住所・電話番号・GoogleマップURL",
    ]),
    ("🔔 管理者通知メール", C_ORANGE, [
        "送信先：admin@kojinius.jp（固定）",
        "送信タイミング：新規予約確定時",
        "送信元：noreply@kojinius.jp",
        "件名：【新規予約】患者名 日時",
        "本文：患者名・日時・電話番号・初再診・症状",
    ]),
    ("⏰ 前日リマインダー", C_GREEN, [
        "送信先：翌日の予約者（email あるもの）",
        "送信タイミング：毎日 15:00 JST（Cloud Scheduler）",
        "送信元：noreply@kojinius.jp",
        "件名：【明日のご予約】患者名 様",
        "本文：日時・院名・電話番号",
        "対象：status = pending or confirmed",
    ]),
]

my = Inches(3.65)
mw = Inches(4.1)
for i, (title, color, items) in enumerate(mail_types):
    cx = Inches(0.2 + i * 4.35)
    add_shape(s, cx, my, mw, Inches(3.5), fill=C_BEIGE2)
    add_shape(s, cx, my, mw, Inches(0.5), fill=color)
    add_text(s, title, cx + Inches(0.15), my + Inches(0.07),
             mw - Inches(0.2), Inches(0.38), size=13, bold=True, color=C_WHITE)
    iy = my + Inches(0.6)
    for item in items:
        add_text(s, "• " + item, cx + Inches(0.15), iy,
                 mw - Inches(0.3), Inches(0.3), size=10, color=C_DARK)
        iy += Inches(0.44)


# ═══════════════════════════════════════════════
# Slide 11 : ディレクトリ構成
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_WHITE)
header_bar(s, "9　ディレクトリ構成", "File Structure")

add_shape(s, Inches(0.2), Inches(1.0), Inches(5.8), Inches(6.1), fill=RGBColor(0x28, 0x2C, 0x34))

dir_lines = [
    ("fukumoto-reservation/",                "",        "#abb2bf"),
    ("├── apps/OnlineAppointSystem/",         "",        "#e5c07b"),
    ("│   ├── index.html",                    "# 予約トップ",   "#98c379"),
    ("│   ├── admin.html",                    "# 管理ダッシュボード", "#98c379"),
    ("│   ├── login.html",                    "# ログイン",      "#98c379"),
    ("│   ├── cancel.html",                   "# 予約キャンセル", "#98c379"),
    ("│   ├── privacy-policy.html",           "# プライバシーポリシー", "#98c379"),
    ("│   ├── maintenance.html",              "# メンテナンス画面",  "#98c379"),
    ("│   ├── auth-action.html",              "# 認証アクションハンドラー", "#98c379"),
    ("│   ├── js/",                           "",        "#e5c07b"),
    ("│   │   ├── firebase.js",               "# Firebase初期化",  "#56b6c2"),
    ("│   │   ├── auth.js",                   "# 認証ロジック",    "#56b6c2"),
    ("│   │   ├── app.js",                    "# 患者側予約ロジック", "#56b6c2"),
    ("│   │   ├── admin.js",                  "# 管理側ロジック",  "#56b6c2"),
    ("│   │   ├── cancel.js",                 "# キャンセルロジック", "#56b6c2"),
    ("│   │   ├── config.js",                 "# Firebase設定値",  "#56b6c2"),
    ("│   │   └── utils.js",                  "# 共通ユーティリティ", "#56b6c2"),
    ("│   ├── css/style.css",                 "# スタイルシート",  "#c678dd"),
    ("│   └── fonts/",                        "# NotoSansJP (PDF用)", "#e5c07b"),
    ("├── functions/",                        "",        "#e5c07b"),
    ("│   └── index.js",                      "# Cloud Functions 3本", "#56b6c2"),
    ("├── firestore.rules",                   "# セキュリティルール", "#e06c75"),
    ("├── firestore.indexes.json",            "# 複合インデックス", "#e06c75"),
    ("└── firebase.json",                     "# Hosting設定",    "#e06c75"),
]

ly = Inches(1.1)
for path, comment, color in dir_lines:
    line = path + ("  " + comment if comment else "")
    add_text(s, line, Inches(0.35), ly, Inches(5.4), Inches(0.26),
             size=9, color=hex_to_rgb(color), font="Courier New")
    ly += Inches(0.26)

# 右側：ファイル役割説明
add_shape(s, Inches(6.2), Inches(1.0), Inches(6.9), Inches(6.1), fill=C_BEIGE2)
add_text(s, "主要ファイルの役割", Inches(6.4), Inches(1.1),
         Inches(6.5), Inches(0.4), size=14, bold=True, color=C_BROWN)

file_roles = [
    ("firebase.js",    "Firebase初期化・エミュレーター対応。db / auth をエクスポート。"),
    ("auth.js",        "requireAdmin()・login()・logout()・パスワードリセット関数。"),
    ("app.js",         "患者側：カレンダー・スロット・予約トランザクション・PDF生成。"),
    ("admin.js",       "管理側：onSnapshot リアルタイム同期・設定モーダル・CSV出力。"),
    ("cancel.js",      "予約キャンセル：予約検索・電話番号照合・writeBatch 更新。"),
    ("utils.js",       "DAY_NAMES・formatDate・THEMES・applyTheme・toHankaku・esc。"),
    ("functions/\nindex.js", "sendReservationEmail・notifyAdminOnReservation・\nsendDailyReminders の3関数。"),
    ("firestore\n.rules", "デフォルト拒否・isAdmin()・isValidReservation()・\n各コレクション権限定義。"),
]

ry2 = Inches(1.55)
for fname, desc in file_roles:
    add_text(s, fname, Inches(6.4), ry2, Inches(1.7), Inches(0.55),
             size=11, bold=True, color=C_BROWN, line_sp=Pt(14))
    add_text(s, desc, Inches(8.2), ry2, Inches(4.7), Inches(0.55),
             size=10, color=C_DARK, line_sp=Pt(14))
    ry2 += Inches(0.68)


# ═══════════════════════════════════════════════
# Slide 12 : クロージング
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_shape(s, Inches(0), Inches(0), W, H, fill=C_BROWN)
add_shape(s, Inches(0), Inches(5.5), W, Inches(2.0), fill=RGBColor(0x5a, 0x40, 0x50))
add_shape(s, Inches(1.0), Inches(2.0), Inches(0.8), Inches(0.07), fill=C_ORANGE)
add_text(s, "詳細設計書",
         Inches(1.0), Inches(2.1), Inches(11), Inches(1.3),
         size=50, bold=True, color=C_WHITE)
add_text(s, "福元鍼灸整骨院 オンライン予約システム",
         Inches(1.0), Inches(3.45), Inches(10), Inches(0.6),
         size=24, color=C_BEIGE)
add_shape(s, Inches(1.0), Inches(4.2), Inches(1.8), Inches(0.04), fill=C_ORANGE)

summary_points = [
    "Firebase（Firestore / Auth / Hosting / Functions）",
    "Firestoreルールによるデフォルト拒否 + admin クレーム認可",
    "パスワード 8文字以上 + 再認証必須 + セッション限定",
    "全フィールドバリデーション（クライアント + Firestoreルール二重）",
    "メール通知 3種（Resend API + Secret Manager）",
]
sy2 = Inches(4.35)
for pt in summary_points:
    add_text(s, "• " + pt, Inches(1.0), sy2, Inches(11), Inches(0.35),
             size=14, color=C_BEIGE)
    sy2 += Inches(0.35)


# ── 保存 ──
out = r"C:\Users\SM7B\Workspace\apps\fukumoto-reservation\documents\OnlineAppointSystem_詳細設計書.pptx"
prs.save(out)
print(f"保存完了: {out}")

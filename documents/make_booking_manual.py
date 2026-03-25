"""
OAS ユーザーマニュアル 予約編 PPTX 生成スクリプト
実行: PYTHONUTF8=1 python documents/make_booking_manual.py
     (fukumoto-reservation ディレクトリから実行)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

# ─── スライドサイズ ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ─── ブランドカラー（Bold Navy × Gold）──────────────────
C_NAVY  = RGBColor(0x1A, 0x2B, 0x4A)
C_GOLD  = RGBColor(0xD4, 0xAF, 0x37)
C_BASE  = RGBColor(0xF8, 0xF5, 0xEE)   # クリーム背景
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK  = RGBColor(0x1A, 0x2B, 0x4A)
C_MUTED = RGBColor(0x6B, 0x7C, 0x9E)
C_WARN  = RGBColor(0xB4, 0x5A, 0x00)

FONT_PRIMARY = "Noto Sans JP"

# ─── パス設定 ─────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent.parent
DOC_DIR     = BASE_DIR / "documents"
SS_DIR      = DOC_DIR / "screenshots"
OUTPUT_DIR  = DOC_DIR / "manuals"
OUTPUT_FILE = OUTPUT_DIR / "OAS_ユーザーマニュアル_予約編.pptx"


# ─── 部品関数（make_user_manual.py から継承）─────────────

def add_shape(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, line_spacing=None, font_name=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    fn = font_name or FONT_PRIMARY
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = fn
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox


def header_bar(slide, title_text, badge_text=None, badge_color=None):
    if badge_color is None:
        badge_color = C_GOLD
    bar_h = Inches(0.65)
    add_shape(slide, 0, 0, W, bar_h, fill_color=C_NAVY)
    add_text(slide, title_text, Inches(0.3), Inches(0.08),
             Inches(9.5), bar_h, font_size=22, bold=True, color=C_WHITE)
    if badge_text:
        badge_w = Inches(2.4)
        badge_x = W - badge_w - Inches(0.2)
        add_shape(slide, badge_x, Inches(0.1), badge_w, Inches(0.45),
                  fill_color=badge_color)
        add_text(slide, badge_text, badge_x, Inches(0.1),
                 badge_w, Inches(0.45), font_size=13, bold=True,
                 color=C_NAVY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, bg=None):
    return add_shape(slide, x, y, w, h, fill_color=bg or C_BASE)


def add_image_fit(slide, img_path, x, y, max_w, max_h):
    if not img_path or not Path(img_path).exists():
        box = add_shape(slide, x, y, max_w, max_h,
                        fill_color=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(slide, "（スクリーンショット準備中）",
                 x, y + max_h // 2 - Inches(0.3),
                 max_w, Inches(0.6), font_size=13,
                 color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)
        return
    try:
        with PILImage.open(img_path) as img:
            iw, ih = img.size
    except Exception:
        add_shape(slide, x, y, max_w, max_h, fill_color=RGBColor(0xCC, 0xCC, 0xCC))
        return
    scale = min(max_w / iw, max_h / ih)
    draw_w = int(iw * scale)
    draw_h = int(ih * scale)
    cx = x + (max_w - draw_w) // 2
    cy = y + (max_h - draw_h) // 2
    slide.shapes.add_picture(str(img_path), cx, cy, draw_w, draw_h)


def screen_slide(prs, blank_layout, title, subtitle, desc_lines, img_path,
                 label_text="患者向け", label_color=None, left_w=Inches(4.3)):
    if label_color is None:
        label_color = C_GOLD
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, title, badge_text=label_text, badge_color=label_color)

    bar_h = Inches(0.65)
    content_y = bar_h + Inches(0.1)
    content_h = H - content_y - Inches(0.15)

    card(slide, 0, content_y, left_w, content_h, bg=C_BASE)

    if subtitle:
        add_text(slide, subtitle, Inches(0.25), content_y + Inches(0.15),
                 left_w - Inches(0.3), Inches(0.5),
                 font_size=12, color=C_MUTED)

    text_y = content_y + (Inches(0.55) if subtitle else Inches(0.2))
    text_h = content_h - (Inches(0.55) if subtitle else Inches(0.2)) - Inches(0.1)
    txBox = slide.shapes.add_textbox(Inches(0.25), text_y, left_w - Inches(0.3), text_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in desc_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if line.startswith("## "):
            run = p.add_run()
            run.text = line[3:]
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = C_NAVY
        elif line.startswith("・"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK
        elif line.startswith("⚠"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(12)
            run.font.color.rgb = C_WARN
        elif line == "":
            run = p.add_run()
            run.text = ""
            run.font.size = Pt(5)
        else:
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK

    right_x = left_w + Inches(0.25)
    right_w = W - right_x - Inches(0.2)
    right_y = content_y + Inches(0.2)
    right_h = content_h - Inches(0.4)
    add_image_fit(slide, img_path, right_x, right_y, right_w, right_h)
    return slide


def screen_slide_two_images(prs, blank_layout, title, subtitle, desc_lines,
                             img_left, img_right, label_text="患者向け"):
    """左テキスト＋右に2枚並べる画像スライド"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, title, badge_text=label_text, badge_color=C_GOLD)

    bar_h = Inches(0.65)
    content_y = bar_h + Inches(0.1)
    content_h = H - content_y - Inches(0.15)
    left_w = Inches(4.3)

    card(slide, 0, content_y, left_w, content_h, bg=C_BASE)

    if subtitle:
        add_text(slide, subtitle, Inches(0.25), content_y + Inches(0.15),
                 left_w - Inches(0.3), Inches(0.5),
                 font_size=12, color=C_MUTED)

    text_y = content_y + (Inches(0.55) if subtitle else Inches(0.2))
    text_h = content_h - (Inches(0.55) if subtitle else Inches(0.2)) - Inches(0.1)
    txBox = slide.shapes.add_textbox(Inches(0.25), text_y, left_w - Inches(0.3), text_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in desc_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if line.startswith("## "):
            run = p.add_run()
            run.text = line[3:]
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = C_NAVY
        elif line.startswith("・"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK
        elif line.startswith("⚠"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(12)
            run.font.color.rgb = C_WARN
        elif line == "":
            run = p.add_run()
            run.text = ""
            run.font.size = Pt(5)
        else:
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK

    # 右エリアを2分割して画像2枚を並べる
    right_x = left_w + Inches(0.2)
    right_w_total = W - right_x - Inches(0.2)
    right_y = content_y + Inches(0.2)
    right_h = content_h - Inches(0.4)
    half_w = right_w_total / 2 - Inches(0.1)

    add_image_fit(slide, img_left,  right_x,                     right_y, half_w, right_h)
    add_image_fit(slide, img_right, right_x + half_w + Inches(0.2), right_y, half_w, right_h)
    return slide


def divider_slide(prs, blank_layout, section_title, items, bg=None, accent=None):
    if bg is None:
        bg = C_NAVY
    if accent is None:
        accent = C_GOLD
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=bg)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=accent)
    add_text(slide, section_title, Inches(0.5), Inches(2.5),
             Inches(8), Inches(1.2),
             font_size=42, bold=True, color=C_WHITE)
    add_shape(slide, Inches(0.5), Inches(3.6), Inches(3.5), Inches(0.06),
              fill_color=accent)
    y_start = Inches(4.0)
    for i, item in enumerate(items):
        add_text(slide, f"▸  {item}", Inches(0.8),
                 y_start + Inches(0.45) * i, Inches(10), Inches(0.45),
                 font_size=16, color=RGBColor(0xFF, 0xE0, 0x80))
    return slide


# ─── カバースライド（予約編専用）────────────────────────────

def cover_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    # 背景
    add_shape(slide, 0, 0, W, Inches(4.5), fill_color=C_NAVY)
    add_shape(slide, 0, Inches(4.5), W, Inches(3.0), fill_color=C_BASE)
    # 左ゴールドバー
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=C_GOLD)
    # タイトル
    add_text(slide, "オンライン予約システム",
             Inches(0.6), Inches(1.1), Inches(11), Inches(1.0),
             font_size=40, bold=True, color=C_WHITE)
    add_text(slide, "ユーザーマニュアル — 予約編",
             Inches(0.6), Inches(2.1), Inches(10), Inches(0.85),
             font_size=30, color=C_GOLD)
    # 区切り線
    add_shape(slide, Inches(0.6), Inches(3.0), Inches(5.0), Inches(0.05), fill_color=C_GOLD)
    # サブテキスト
    add_text(slide, "2026年3月  |  患者向け操作ガイド",
             Inches(0.6), Inches(3.2), Inches(5), Inches(0.6),
             font_size=16, color=C_MUTED)
    # 右下アクセント
    add_text(slide, "OAS", Inches(10.5), Inches(5.3), Inches(2.5), Inches(1.5),
             font_size=72, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    return slide


# ─── 目次スライド（予約編）──────────────────────────────────

def toc_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, "目次", badge_text="患者向け", badge_color=C_GOLD)

    bar_h = Inches(0.65)
    col_y = bar_h + Inches(0.4)
    left_x = Inches(0.6)
    col_w = Inches(5.5)

    add_text(slide, "予約編 — 目次",
             left_x, col_y, col_w, Inches(0.55),
             font_size=18, bold=True, color=C_NAVY)
    add_shape(slide, left_x, col_y + Inches(0.55),
              Inches(3.0), Inches(0.04), fill_color=C_GOLD)

    toc_items = [
        ("1", "予約の流れ（概要）"),
        ("2", "日時の選択 — カレンダーで日付を選ぶ"),
        ("3", "日時の選択 — 時間枠を選ぶ"),
        ("4", "患者情報の入力"),
        ("5", "確認・完了"),
        ("6", "キャンセル方法"),
        ("7", "問診票について"),
        ("8", "プライバシーポリシー"),
    ]
    for i, (num, label) in enumerate(toc_items):
        item_y = col_y + Inches(0.8) + Inches(0.52) * i
        # 番号バッジ
        add_shape(slide, left_x, item_y + Inches(0.05),
                  Inches(0.45), Inches(0.38), fill_color=C_NAVY)
        add_text(slide, num, left_x, item_y + Inches(0.02),
                 Inches(0.45), Inches(0.42),
                 font_size=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, label, left_x + Inches(0.55), item_y,
                 col_w - Inches(0.55), Inches(0.45),
                 font_size=14, color=C_DARK)

    # 右側：注意事項ボックス
    note_x = Inches(7.2)
    note_w = Inches(5.7)
    note_y = col_y
    note_h = Inches(5.6)
    add_shape(slide, note_x, note_y, note_w, note_h, fill_color=C_BASE)
    add_text(slide, "ご利用の流れ", note_x + Inches(0.25), note_y + Inches(0.2),
             note_w - Inches(0.5), Inches(0.45),
             font_size=15, bold=True, color=C_NAVY)
    add_shape(slide, note_x + Inches(0.25), note_y + Inches(0.65),
              Inches(2.0), Inches(0.04), fill_color=C_GOLD)

    flow_lines = [
        "ブラウザで URL を開く",
        "  ↓",
        "カレンダーで希望日を選ぶ",
        "  ↓",
        "時間枠を選ぶ",
        "  ↓",
        "患者情報を入力する",
        "  ↓",
        "内容を確認して予約を確定",
        "  ↓",
        "予約番号を控えておく",
    ]
    for i, fl in enumerate(flow_lines):
        color = C_MUTED if fl.strip() == "↓" else C_DARK
        size = 11 if fl.strip() == "↓" else 13
        add_text(slide, fl, note_x + Inches(0.3),
                 note_y + Inches(0.9) + Inches(0.42) * i,
                 note_w - Inches(0.5), Inches(0.42),
                 font_size=size, color=color)
    return slide


# ─── 予約フロー概要スライド ──────────────────────────────────

def flow_overview_slide(prs, blank_layout):
    """3ステップ図解スライド"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, "予約の流れ（概要）", badge_text="患者向け", badge_color=C_GOLD)

    bar_h = Inches(0.65)
    content_y = bar_h + Inches(0.3)

    steps = [
        ("Step 1", "日時の選択", "カレンダーで希望日を\nクリック後、時間枠を選ぶ"),
        ("Step 2", "情報の入力", "氏名・電話番号などの\n患者情報を入力する"),
        ("Step 3", "確認・完了", "内容を確認して\n「予約を確定する」"),
    ]

    step_w = Inches(3.5)
    step_h = Inches(4.5)
    gap = Inches(0.55)
    total_w = step_w * 3 + gap * 2
    start_x = (W - total_w) / 2

    for i, (num, title, desc) in enumerate(steps):
        sx = start_x + (step_w + gap) * i
        sy = content_y + Inches(0.3)

        # カードベース
        add_shape(slide, sx, sy, step_w, step_h, fill_color=C_BASE)
        # ネイビートップバー
        add_shape(slide, sx, sy, step_w, Inches(0.6), fill_color=C_NAVY)
        # ステップ番号
        add_text(slide, num, sx, sy + Inches(0.05), step_w, Inches(0.5),
                 font_size=15, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        # タイトル
        add_text(slide, title, sx, sy + Inches(0.75), step_w, Inches(0.7),
                 font_size=22, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        # ゴールドライン
        add_shape(slide, sx + Inches(0.5), sy + Inches(1.45),
                  step_w - Inches(1.0), Inches(0.04), fill_color=C_GOLD)
        # 説明テキスト
        add_text(slide, desc, sx + Inches(0.2), sy + Inches(1.65),
                 step_w - Inches(0.4), Inches(2.0),
                 font_size=15, color=C_DARK, align=PP_ALIGN.CENTER)

        # 矢印（最後のステップ以外）
        if i < 2:
            arrow_x = sx + step_w + gap * 0.15
            arrow_y = sy + step_h / 2 - Inches(0.15)
            add_text(slide, "▶", arrow_x, arrow_y, gap * 0.7, Inches(0.4),
                     font_size=22, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 下部補足
    note_y = content_y + Inches(5.2)
    add_text(slide,
             "⚠ 予約完了後に表示される「予約番号」はキャンセル時に必要です。スクリーンショットで保存することをお勧めします。",
             Inches(0.8), note_y, W - Inches(1.6), Inches(0.5),
             font_size=12, color=C_WARN)
    return slide


# ─── クロージングスライド ────────────────────────────────────

def closing_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_NAVY)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=C_GOLD)
    add_text(slide, "ご予約ありがとうございます",
             Inches(0.6), Inches(2.4), Inches(12), Inches(1.4),
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(3.5), Inches(3.8), Inches(6.3), Inches(0.06), fill_color=C_GOLD)
    add_text(slide, "ご不明な点はお気軽に院内スタッフへお問い合わせください",
             Inches(0.6), Inches(4.1), Inches(12), Inches(0.7),
             font_size=18, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "Online Appointment System — 予約編",
             Inches(0.6), Inches(5.0), Inches(12), Inches(0.7),
             font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    return slide


# ─── メイン処理 ───────────────────────────────────────────

def build_pptx():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank_layout = prs.slide_layouts[6]

    # ── 1. カバー ──────────────────────────────────────────
    cover_slide(prs, blank_layout)

    # ── 2. 目次 ────────────────────────────────────────────
    toc_slide(prs, blank_layout)

    # ── 3. セクション区切り「1. 予約の流れ（概要）」 ─────────
    divider_slide(prs, blank_layout,
                  "1. 予約の流れ（概要）",
                  [
                      "ブラウザでアクセスするだけ — アプリ不要",
                      "Step 1: 日時選択  →  Step 2: 情報入力  →  Step 3: 確認",
                      "予約番号で後からキャンセル可能",
                  ])

    # ── 4. 予約フロー概要（図解）─────────────────────────────
    flow_overview_slide(prs, blank_layout)

    # ── 5. セクション区切り「2. 日時の選択」 ──────────────────
    divider_slide(prs, blank_layout,
                  "2. 日時の選択",
                  [
                      "月カレンダーで希望日をクリック",
                      "予約可能日は通常カラー表示 / 休診日はグレー",
                      "時間枠から希望の時間を選んで「次へ」",
                  ])

    # ── 6. カレンダーで日付を選ぶ ─────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="カレンダーで日付を選ぶ",
        subtitle="Step 1 — 日時選択 ① 日付を選ぶ",
        desc_lines=[
            "## アクセス方法",
            "・ブラウザで URL を開くだけ",
            "・アプリインストールは不要",
            "",
            "## カレンダーの見方",
            "・グレー = 予約不可（休業日・過去日）",
            "・通常カラー = 予約可能な日",
            "・希望日をクリックすると時間枠が表示",
            "",
            "## 月の切り替え",
            "・「＜」「＞」で前後の月を表示",
            "・過去の日付は選択できません",
        ],
        img_path=SS_DIR / "01_calendar_top.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 7. 時間枠を選ぶ（2枚画像）────────────────────────────
    screen_slide_two_images(
        prs, blank_layout,
        title="時間枠を選ぶ",
        subtitle="Step 1 — 日時選択 ② 時間を選ぶ",
        desc_lines=[
            "## 時間枠の見方",
            "・日付をクリックすると時間枠が表示",
            "・「予約済み」枠は選択できません",
            "・残り枠数が少ない場合は表示されます",
            "",
            "## 選択後の操作",
            "・枠をクリックするとハイライト表示",
            "・「次へ」ボタンが有効になります",
            "・クリックで情報入力フォームへ",
            "",
            "⚠ 選択中に他の患者が予約した場合",
            "　確定時にエラーが表示されます",
        ],
        img_left=SS_DIR / "02_time_slots.png",
        img_right=SS_DIR / "02b_slot_selected.png",
        label_text="患者向け",
    )

    # ── 8. セクション区切り「3. 患者情報の入力」 ──────────────
    divider_slide(prs, blank_layout,
                  "3. 患者情報の入力",
                  [
                      "お名前・ふりがな・電話番号は必須",
                      "郵便番号を入力すると住所が自動補完",
                      "症状・お悩みを入力してスムーズな診察へ",
                  ])

    # ── 9. 情報入力フォーム ────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="患者情報入力フォーム（Step 2）",
        subtitle="必要事項をご入力ください",
        desc_lines=[
            "## 基本情報（必須）",
            "・氏名・ふりがな",
            "・生年月日",
            "・電話番号",
            "",
            "## 住所（必須）",
            "・郵便番号（7桁）→ 住所が自動補完",
            "⚠ 郵便番号は半角数字で入力してください",
            "",
            "## 診療情報",
            "・症状・お悩み（必須）",
            "・初診 / 再診・保険証の有無",
            "・メールアドレス（任意）",
            "",
            "⚠ 健康情報の取り扱い同意が必要です",
        ],
        img_path=SS_DIR / "03_form_step2.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 10. セクション区切り「4. 確認・完了」 ─────────────────
    divider_slide(prs, blank_layout,
                  "4. 確認・完了",
                  [
                      "入力内容を最終確認してから予約を確定",
                      "「戻る」で修正も可能",
                      "完了後に予約番号が発行 — 必ず控えておく",
                  ])

    # ── 11. 予約確認画面 ────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="予約確認画面（Step 3）",
        subtitle="予約内容を最終確認してください",
        desc_lines=[
            "## 確認内容",
            "・予約日時",
            "・氏名・ふりがな・生年月日",
            "・電話番号・住所・メール",
            "・初診/再診・保険証・症状",
            "",
            "## 修正したい場合",
            "・「戻る」で入力フォームに戻れます",
            "",
            "## 予約の確定",
            "・「予約を確定する」をクリック",
            "・サーバー側で空き確認後に登録",
            "",
            "⚠ 同時刻に他の予約が入った場合は",
            "　別の時間帯を選び直してください",
        ],
        img_path=SS_DIR / "05_confirm_step3.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 12. セクション区切り「5. その他の機能」 ───────────────
    divider_slide(prs, blank_layout,
                  "5. その他の機能",
                  [
                      "予約のキャンセル（予約番号 + 電話番号で本人確認）",
                      "問診票の事前記入でスムーズな来院",
                      "プライバシーポリシーで個人情報の取り扱いを確認",
                  ])

    # ── 13. キャンセル方法 ────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="キャンセル方法",
        subtitle="予約番号と電話番号で本人確認",
        desc_lines=[
            "## キャンセルの手順",
            "・予約番号（完了時に発行）を入力",
            "・電話番号（予約時と同じ）を入力",
            "・「予約を確認する」をクリック",
            "・内容を確認して「キャンセルする」",
            "",
            "## キャンセル理由",
            "・理由の入力は任意です",
            "・法令上、理由なくキャンセル可能",
            "",
            "⚠ 受付期限を過ぎた場合は",
            "　お電話でご連絡ください",
            "⚠ 前日までのキャンセルを推奨します",
        ],
        img_path=SS_DIR / "06_cancel_page.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 14. 問診票について ────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="問診票について",
        subtitle="来院前に記入でスムーズな施術へ",
        desc_lines=[
            "## 問診票の目的",
            "・来院前に症状の詳細を伝えられます",
            "・院側が事前に準備できスムーズに",
            "",
            "## アクセス方法",
            "・予約完了画面の",
            "　「問診票を記入する」ボタンから",
            "・URL: /questionnaire",
            "",
            "## 記入のタイミング",
            "・来院前日〜当日の記入を推奨",
            "・記入は任意ですが、ご協力ください",
            "",
            "⚠ 記入情報は院のみが閲覧します",
        ],
        img_path=SS_DIR / "08_questionnaire.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 15. プライバシーポリシー ──────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="プライバシーポリシー",
        subtitle="個人情報の取り扱いについて",
        desc_lines=[
            "## 個人情報の利用目的",
            "・予約管理・施術準備のみに使用",
            "・第三者への提供は行いません",
            "",
            "## 要配慮個人情報",
            "・症状・健康情報は個人情報保護法上",
            "　の「要配慮個人情報」に該当",
            "・入力前に同意確認を行います",
            "",
            "## お問い合わせ",
            "・個人情報に関するご要望は",
            "　院内窓口へお申し出ください",
            "",
            "⚠ プライバシーポリシーへの同意が",
            "　予約完了に必要です",
        ],
        img_path=SS_DIR / "07_privacy_policy.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 16. クロージング ──────────────────────────────────────
    closing_slide(prs, blank_layout)

    prs.save(str(OUTPUT_FILE))
    print(f"✅ 生成完了: {OUTPUT_FILE}")
    print(f"   スライド枚数: {len(prs.slides)}")


if __name__ == "__main__":
    build_pptx()

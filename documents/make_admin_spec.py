# -*- coding: utf-8 -*-
"""
OAS 管理画面 UI 仕様書 PPTX 生成スクリプト v2.0
実行: PYTHONUTF8=1 python documents/make_admin_spec.py
     （fukumoto-reservation ディレクトリから）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from PIL import Image as PILImage

# ─────────────────────────────────────────────
# カラー定義（make_design_doc.py と完全同一）
# ─────────────────────────────────────────────
C_NAVY    = RGBColor(0x1A, 0x2B, 0x4A)   # ダークネイビー（管理者ヘッダー）
C_GOLD    = RGBColor(0xD4, 0xAF, 0x37)   # ゴールド（アクセント）
C_BASE    = RGBColor(0xF8, 0xF5, 0xEE)   # クリーム（背景）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x2B, 0x4A)   # ネイビー（テキスト）
C_MUTED   = RGBColor(0x6B, 0x7C, 0x9E)   # ミュートブルー
C_RED     = RGBColor(0xDC, 0x26, 0x26)   # エラー・警告

# ─────────────────────────────────────────────
# パス設定
# ─────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent.parent
SS_DIR      = BASE_DIR / "documents" / "screenshots"
OUTPUT_PATH = BASE_DIR / "documents" / "specs" / "OAS_管理画面仕様書.pptx"

W = Inches(13.33)
H = Inches(7.5)


# ─────────────────────────────────────────────
# 共通ヘルパー関数（make_design_doc.py と同一実装）
# ─────────────────────────────────────────────

def add_shape(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, line_spacing=None, font_name="Noto Sans JP"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def header_bar(slide, title_text, badge_text=None, badge_color=None):
    """ヘッダーバー — ナビーグレー背景＋ゴールドバッジ"""
    if badge_color is None:
        badge_color = C_GOLD
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.9), fill_color=C_NAVY)
    add_text(slide, title_text, Inches(0.4), Inches(0.1), Inches(10), Inches(0.75),
             font_size=26, bold=True, color=C_WHITE)
    if badge_text:
        s = add_shape(slide, Inches(11.3), Inches(0.2), Inches(1.8), Inches(0.5),
                      fill_color=badge_color)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge_text
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.name = "Noto Sans JP"
        run.font.color.rgb = C_NAVY


def card(slide, x, y, w, h, bg=None):
    if bg is None:
        bg = C_BASE
    return add_shape(slide, x, y, w, h, fill_color=bg)


def add_image_fit(slide, img_path, x, y, max_w, max_h):
    """画像をアスペクト比を保ってフィット表示。失敗時はスキップ（テキストのみにする）"""
    try:
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        scale = min(max_w / iw, max_h / ih)
        disp_w = int(iw * scale)
        disp_h = int(ih * scale)
        off_x = (max_w - disp_w) // 2
        off_y = (max_h - disp_h) // 2
        slide.shapes.add_picture(str(img_path), x + off_x, y + off_y, disp_w, disp_h)
    except Exception:
        # 画像が存在しない場合はグレーのプレースホルダーを表示
        add_shape(slide, x, y, max_w, max_h, fill_color=C_MUTED)
        add_text(slide, "（スクリーンショット未取得）",
                 x + Inches(0.2), y + Inches(0.2), max_w - Inches(0.4), Inches(0.5),
                 font_size=12, color=C_WHITE)


def divider_slide(prs, blank_layout, section_title, items):
    """セクション区切りスライド（ネイビー背景）"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.2), W, Inches(1.2), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(2.0), Inches(12.33), Inches(0.06), fill_color=C_GOLD)
    add_text(slide, section_title, Inches(0.5), Inches(1.0), Inches(12), Inches(1.0),
             font_size=40, bold=True, color=C_WHITE)
    y = Inches(2.3)
    for item in items:
        add_text(slide, "  " + item, Inches(0.7), y, Inches(12), Inches(0.5),
                 font_size=16, color=C_BASE)
        y += Inches(0.5)


def three_col_cards(slide, cards_data, top=Inches(1.1)):
    """3カラムカード"""
    col_w = Inches(3.9)
    gap   = Inches(0.3)
    card_h = Inches(5.8)
    starts = [Inches(0.4), Inches(0.4) + col_w + gap,
              Inches(0.4) + (col_w + gap) * 2]
    for idx, (title, lines) in enumerate(cards_data):
        cx = starts[idx]
        card(slide, cx, top, col_w, card_h)
        add_shape(slide, cx, top, col_w, Inches(0.55), fill_color=C_NAVY)
        add_text(slide, title, cx + Inches(0.1), top + Inches(0.05),
                 col_w - Inches(0.2), Inches(0.5),
                 font_size=15, bold=True, color=C_WHITE)
        body_text = "\n".join(lines)
        add_text(slide, body_text, cx + Inches(0.15), top + Inches(0.65),
                 col_w - Inches(0.3), card_h - Inches(0.75),
                 font_size=13, color=C_DARK, line_spacing=Pt(20))


def two_col_cards(slide, left_title, left_lines, right_title, right_lines,
                  top=Inches(1.1)):
    """2カラムカード"""
    col_w = Inches(5.9)
    gap   = Inches(0.4)
    card_h = Inches(5.8)
    left_x  = Inches(0.5)
    right_x = left_x + col_w + gap
    for cx, title, lines in [(left_x, left_title, left_lines),
                              (right_x, right_title, right_lines)]:
        card(slide, cx, top, col_w, card_h)
        add_shape(slide, cx, top, col_w, Inches(0.55), fill_color=C_NAVY)
        add_text(slide, title, cx + Inches(0.1), top + Inches(0.05),
                 col_w - Inches(0.2), Inches(0.5),
                 font_size=15, bold=True, color=C_WHITE)
        body_text = "\n".join(lines)
        add_text(slide, body_text, cx + Inches(0.15), top + Inches(0.65),
                 col_w - Inches(0.3), card_h - Inches(0.75),
                 font_size=13, color=C_DARK, line_spacing=Pt(20))


def image_with_notes(slide, img_path, notes_lines, top=Inches(1.1)):
    """左側に画像、右側にノート（make_design_doc.py と同一）"""
    img_x, img_y = Inches(0.4), top
    img_max_w, img_max_h = Inches(8.0), Inches(5.8)
    note_x = Inches(8.7)
    note_w = Inches(4.3)
    if img_path and Path(img_path).exists():
        add_image_fit(slide, img_path, img_x, img_y, img_max_w, img_max_h)
    else:
        card(slide, img_x, img_y, img_max_w, img_max_h, bg=C_MUTED)
        add_text(slide, "（スクリーンショット未取得）",
                 img_x + Inches(0.2), img_y + Inches(0.2),
                 img_max_w - Inches(0.4), Inches(0.5),
                 font_size=12, color=C_WHITE)
    card(slide, note_x, top, note_w, Inches(5.8))
    body_text = "\n".join(notes_lines)
    add_text(slide, body_text, note_x + Inches(0.15), top + Inches(0.2),
             note_w - Inches(0.3), Inches(5.4),
             font_size=13, color=C_DARK, line_spacing=Pt(22))


# ─────────────────────────────────────────────
# スライド生成関数群
# ─────────────────────────────────────────────

def slide_01_cover(prs, blank_layout):
    """表紙"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.5), W, Inches(1.5), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(2.8), Inches(12.33), Inches(0.08), fill_color=C_GOLD)
    add_text(slide, "Online Appointment System",
             Inches(0.5), Inches(1.2), Inches(12), Inches(1.1),
             font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "管理画面 UI 仕様書  v2.0",
             Inches(0.5), Inches(3.0), Inches(12), Inches(0.9),
             font_size=36, bold=True, color=C_BASE, align=PP_ALIGN.CENTER)
    add_text(slide, "2026年3月  |  管理者向け各画面コンポーネント・バリデーション仕様",
             Inches(0.5), Inches(3.9), Inches(12), Inches(0.6),
             font_size=18, color=C_BASE, align=PP_ALIGN.CENTER)


def slide_02_toc(prs, blank_layout):
    """目次"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "目次", badge_text="管理画面")

    toc_items = [
        ("1",  "認証・アカウント",        "ログイン / パスワード変更強制フロー"),
        ("2",  "ダッシュボード",           "KPI カード / 検索・フィルター / 予約テーブル / 詳細モーダル"),
        ("3",  "診察履歴",                 "履歴ページ / CSV出力 / 訂正機能"),
        ("4",  "設定 — 基本情報",          "院名・電話・URL・住所・Google Maps"),
        ("5",  "設定 — 営業日設定",        "タイムラインバー・AM/PM・休日・締切"),
        ("6",  "設定 — お知らせ",          "バナー設定・メンテナンスモード"),
        ("7",  "設定 — ポリシー・規約",    "PP自動生成・要配慮同意・患者権利"),
        ("8",  "設定 — アカウント",        "管理者作成・メール変更・PW変更"),
    ]

    y = Inches(1.15)
    for num, section, desc in toc_items:
        badge = add_shape(slide, Inches(0.5), y + Inches(0.05),
                          Inches(0.5), Inches(0.4), fill_color=C_NAVY)
        tf = badge.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.name = "Noto Sans JP"
        run.font.color.rgb = C_WHITE
        add_text(slide, section, Inches(1.1), y,
                 Inches(3.5), Inches(0.5), font_size=14, bold=True, color=C_DARK)
        add_text(slide, desc, Inches(4.8), y,
                 Inches(8.2), Inches(0.5), font_size=13, color=C_MUTED)
        add_shape(slide, Inches(0.5), y + Inches(0.5), Inches(12.33), Inches(0.02), fill_color=C_MUTED)
        y += Inches(0.52)


# ── Section 1: 認証・アカウント ──────────────────────────

def slide_04_login(prs, blank_layout):
    """ログイン画面仕様"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "ログイン画面仕様（/login）", badge_text="管理画面")

    two_col_cards(slide,
        "UIコンポーネント",
        [
            "## ページ構成",
            "・<LoginLayout> ── 中央寄せ白カード",
            "・ロゴ（院名頭文字アイコン）",
            "・タイトル「管理者ログイン」",
            "",
            "## 入力フィールド",
            "・メールアドレス",
            "  type=email / 必須 / オートコンプリート",
            "・パスワード",
            "  <PasswordInput> / 表示トグル付き",
            "  必須 / minLength=8",
            "",
            "## ボタン",
            "・ログイン（primary / full-width）",
            "  loading 中はスピナー表示",
            "",
            "## リンク",
            "・パスワードをお忘れですか？",
            "  → sendPasswordResetEmail()",
        ],
        "バリデーション・エラー",
        [
            "## クライアント側",
            "・メール形式チェック（isValidEmail）",
            "  連続ドット・末尾ドット拒否",
            "・パスワード空欄チェック",
            "",
            "## Firebase Auth エラーハンドリング",
            "・auth/wrong-password",
            "  → 「メール・パスワードが違います」",
            "・auth/user-not-found",
            "  → 「メール・パスワードが違います」",
            "  （列挙攻撃防止のため同一文言）",
            "・auth/too-many-requests",
            "  → 「しばらく後にお試しください」",
            "",
            "## セッション",
            "・browserSessionPersistence",
            "  → タブ閉じでセッション消滅",
            "",
            "## 同意モーダル",
            "・未同意時: <ConsentModal> でブロッキング表示",
            "  TOS / PP 全文表示 → 同意ボタン",
        ],
        top=Inches(1.1),
    )


def slide_05_change_password(prs, blank_layout):
    """パスワード変更強制フロー"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "パスワード変更強制フロー（/admin/change-password）", badge_text="管理画面")

    img_path = SS_DIR / "19_change_password.png"
    notes = [
        "## トリガー条件",
        "・users/{uid}.mustChangePassword",
        "  == true のとき",
        "・AdminLayout が",
        "  ?forced=1 付きで強制リダイレクト",
        "",
        "## UIコンポーネント",
        "・現在のパスワード入力",
        "  <PasswordInput> / 必須",
        "・新しいパスワード入力",
        "  <PasswordInput> / 必須",
        "・確認パスワード入力",
        "  <PasswordInput> / 必須",
        "",
        "## バリデーション",
        "・isStrongPassword()",
        "  大文字・小文字・数字・記号",
        "  各1文字以上 / 8〜128文字",
        "・新旧パスワード一致チェック",
        "",
        "## 完了後",
        "・mustChangePassword = false",
        "・/admin へリダイレクト",
    ]
    image_with_notes(slide, img_path, notes)


# ── Section 2: ダッシュボード ──────────────────────────

def slide_07_kpi(prs, blank_layout):
    """KPI カード仕様"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "KPI カード仕様", badge_text="管理画面")

    three_col_cards(slide, [
        ("KPI カード 4種", [
            "① 本日の予約",
            "  カレンダーアイコン / 今日の件数",
            "  クリック → 日付フィルター適用",
            "",
            "② 今月の予約",
            "  グラフアイコン / 当月合計件数",
            "  クリック → 月次フィルター適用",
            "",
            "③ 新規患者",
            "  人物+プラスアイコン",
            "  visitType=='first' の件数",
            "  クリック → 初診フィルター適用",
            "",
            "④ 未確認",
            "  時計アイコン / status=='pending' 数",
            "  クリック → 未確認フィルター適用",
            "  バッジ数がヘッダーにも表示",
        ]),
        ("表示仕様", [
            "・<Card> コンポーネント使用",
            "  bg=white / shadow-sm",
            "",
            "・数値フォント: 32pt bold",
            "・ラベル: 12pt muted",
            "",
            "・アクティブ状態:",
            "  border-gold + shadow-gold",
            "  （選択中KPIを強調）",
            "",
            "・再クリックでフィルター解除",
            "",
            "・モバイル: 2×2 グリッド",
            "  (sm以上で横1行並び)",
            "",
            "・ローディング中:",
            "  スケルトンアニメーション",
        ]),
        ("データソース", [
            "## useKpis() hook",
            "・reservations 配列を引数に取る",
            "・todayCount: 今日の件数",
            "・monthCount: 今月の件数",
            "・newCount: 初診件数",
            "・pendingCount: 未確認件数",
            "",
            "## usePendingCount() hook",
            "・Firestore リアルタイム購読",
            "・status=='pending' をリアルタイム監視",
            "・ヘッダーバッジと連動",
            "  99件超は '99+' 表示",
            "  animate-pulse（アニメーション）",
            "",
            "## useReservations() hook",
            "・onSnapshot でリアルタイム同期",
            "・管理者ロール確認済みのみ取得",
        ]),
    ])


def slide_08_search_filter(prs, blank_layout):
    """検索・フィルター仕様"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "検索・フィルター仕様", badge_text="管理画面")

    two_col_cards(slide,
        "検索フィールド（5種）",
        [
            "① 氏名検索",
            "  placeholder: 山田太郎",
            "  氏名・ふりがな部分一致",
            "",
            "② 郵便番号検索",
            "  placeholder: 123-4567",
            "  zipSearch 状態で管理",
            "",
            "③ 電話番号検索",
            "  placeholder: 090-1234-5678",
            "  phoneSearch 状態で管理",
            "",
            "④ 診察予定日",
            "  type=date / dateFilter 状態",
            "  KPIカード「本日」と連動",
            "",
            "⑤ 予約受付日",
            "  type=date / createdAtFilter 状態",
            "",
            "・「クリア」ボタン: 全フィルター初期化",
            "・「CSV出力」ボタン: フィルター適用後データ出力",
        ],
        "ステータスフィルター",
        [
            "## ボタングループ（5種）",
            "・全て（デフォルト）",
            "・未確認",
            "・確認済み",
            "・キャンセル",
            "・診察完了",
            "",
            "## スタイル",
            "・選択中: variant=primary",
            "・非選択: variant=ghost",
            "・size=sm",
            "",
            "## URL パラメータ連携",
            "・?filter=pending でヘッダーバッジ",
            "  クリック時に自動フィルター適用",
            "・useSearchParams() で管理",
            "",
            "## KPIフィルター（kpiFilter状態）",
            "・today / month / new / pending",
            "・KPIカードクリックで切り替え",
            "・再クリックで解除",
        ],
        top=Inches(1.1),
    )


def slide_09_table(prs, blank_layout):
    """予約テーブル列定義"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "予約テーブル列定義", badge_text="管理画面")

    # テーブルヘッダー
    cols = [
        ("診察予定日時", "datetime", "日時降順デフォルト / ソート可"),
        ("氏名",         "name",     "ふりがな順でソート"),
        ("初/再",        "visitType","初診 / 再診 バッジ表示"),
        ("症状",         "symptoms", "モーダルで全文表示 / クリック展開"),
        ("電話",         "phone",    "tel: リンク（電話発信）"),
        ("予約受付日",   "createdAt","受付日降順 / ソート可"),
        ("ステータス",   "status",   "<Badge> コンポーネント / 色分け"),
        ("操作",         "—",        "「詳細」ボタン → 詳細モーダル表示"),
    ]

    col_w = Inches(2.0)
    gap   = Inches(0.02)
    headers = ["列名", "sortKey", "仕様・備考"]
    header_widths = [Inches(2.2), Inches(2.0), Inches(7.8)]

    y = Inches(1.1)
    # ヘッダー行
    x = Inches(0.4)
    for i, (hdr, hw) in enumerate(zip(headers, header_widths)):
        add_shape(slide, x, y, hw, Inches(0.45), fill_color=C_NAVY)
        add_text(slide, hdr, x + Inches(0.1), y + Inches(0.07),
                 hw - Inches(0.2), Inches(0.35), font_size=13, bold=True, color=C_WHITE)
        x += hw + gap

    y += Inches(0.45)
    row_h = Inches(0.58)
    for col_name, sort_key, spec in cols:
        x = Inches(0.4)
        bg = C_BASE if cols.index((col_name, sort_key, spec)) % 2 == 0 else C_WHITE
        for val, hw in [(col_name, header_widths[0]),
                        (sort_key, header_widths[1]),
                        (spec,     header_widths[2])]:
            add_shape(slide, x, y, hw, row_h, fill_color=bg)
            add_text(slide, val, x + Inches(0.1), y + Inches(0.1),
                     hw - Inches(0.2), row_h - Inches(0.15),
                     font_size=12, color=C_DARK)
            x += hw + gap
        y += row_h

    # ステータスバッジ色凡例
    add_text(slide, "ステータスバッジ色: pending=amber / confirmed=emerald / cancelled=red / completed=sky",
             Inches(0.4), y + Inches(0.08), Inches(12.5), Inches(0.4),
             font_size=11, color=C_MUTED)


def slide_10_modal(prs, blank_layout):
    """予約詳細モーダル仕様"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "予約詳細モーダル仕様", badge_text="管理画面")

    img_path = SS_DIR / "11_reservation_detail_modal.png"
    notes = [
        "## 表示フィールド",
        "・予約番号（reservationId）",
        "・診察予定日時",
        "・氏名 / ふりがな / 年齢",
        "・生年月日 / 郵便番号",
        "・住所 / 電話番号",
        "・メール / 性別",
        "・初診/再診 / 保険証",
        "・症状 / 伝達事項",
        "・連絡方法",
        "・ステータス / キャンセル情報",
        "・予約受付日 / 診察完了日時",
        "",
        "## アクション",
        "・確認済みにする",
        "  → status: confirmed",
        "・診察完了",
        "  → completeVisit() CF 呼び出し",
        "  → visit_histories に保存",
        "・問診票PDF ダウンロード",
        "  → pdf-lib で生成",
        "・キャンセル",
        "  → キャンセル理由選択必須",
        "  → cancelReservation CF",
        "・閉じる",
    ]
    image_with_notes(slide, img_path, notes)


# ── Section 3: 診察履歴 ──────────────────────────

def slide_12_history(prs, blank_layout):
    """診察履歴ページ仕様"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "診察履歴ページ仕様（/admin/history）", badge_text="管理画面")

    img_path = SS_DIR / "12_history_page.png"
    notes = [
        "## 検索フィールド",
        "・氏名（部分一致）",
        "・郵便番号",
        "・電話番号",
        "・診察日 From / To",
        "",
        "## テーブル列",
        "・診療日時（sortable）",
        "・氏名（sortable）",
        "・初/再（visitType）",
        "・電話番号",
        "・詳細 / 訂正 ボタン",
        "",
        "## ページネーション",
        "・20件/ページ",
        "・「件」サフィックス表示",
        "",
        "## CSV出力",
        "・フィルター適用済みデータ",
        "・ファイル名: 診察履歴_{日付}",
        "",
        "## 訂正機能（APPI第34条）",
        "・correctVisitHistory CF 呼び出し",
        "・corrections サブコレクションに",
        "  immutable な訂正記録を追加",
        "・訂正理由必須入力",
    ]
    image_with_notes(slide, img_path, notes)


# ── Section 4: 設定画面 ──────────────────────────

def slide_14_basic_info(prs, blank_layout):
    """基本情報タブ"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "設定 — 基本情報タブ", badge_text="管理画面")

    img_path = SS_DIR / "13_settings_basic_info.png"
    notes = [
        "## 入力フィールド",
        "・院名 *（clinicName）必須",
        "  バリデーション: 空欄NG",
        "・電話番号 *（phone）必須",
        "  バリデーション: 空欄NG",
        "・WebサイトURL（websiteUrl）",
        "  任意 / type=url",
        "・郵便番号（zip）",
        "  → 入力後に住所自動補完",
        "  zipcloud API 使用",
        "・都道府県・市区町村・番地 *",
        "  address / 必須",
        "・建物名・号室（addressSub）",
        "  任意",
        "",
        "## Google Maps",
        "・「所在地マップ」<iframe>",
        "  clinicName をクエリに埋め込み",
        "  frame-src: google.com/maps (CSP)",
        "",
        "## 保存",
        "・画面上部「保存」ボタン共通",
        "・Firestore settings/clinic に",
        "  merge: true で保存",
    ]
    image_with_notes(slide, img_path, notes)


def slide_15_business_days(prs, blank_layout):
    """営業日設定タブ"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "設定 — 営業日設定タブ", badge_text="管理画面")

    img_path = SS_DIR / "14_settings_business_days.png"
    notes = [
        "## タイムラインバー",
        "・午前診療 / 午後診療 の",
        "  グラフィカルバー表示",
        "・各 AM/PM: ON/OFF トグル",
        "",
        "## AM/PM 時刻設定",
        "・AM: 6:00〜12:00 (30分刻み)",
        "・PM: 13:00〜22:00 (30分刻み)",
        "・開始 < 終了 バリデーション",
        "  超えると「開始が終了以降」エラー",
        "",
        "## 曜日別設定",
        "・日〜土 7曜日",
        "  各々「営業」/「休診」トグル",
        "・土日→休診 / 土日→診療 一括ボタン",
        "",
        "## 休日管理",
        "・カレンダークリックで休日追加",
        "・「祝日を自動取得」ボタン",
        "  holidays-jp API 使用",
        "・「クリア」で全休日削除",
        "・予約がある日を休日設定",
        "  → 確認ダイアログ表示",
        "",
        "## キャンセル締切",
        "・診療開始のX分前まで",
        "  cancelCutoffMinutes",
    ]
    image_with_notes(slide, img_path, notes)


def slide_16_announcement(prs, blank_layout):
    """お知らせタブ"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "設定 — お知らせタブ", badge_text="管理画面")

    img_path = SS_DIR / "15_settings_announcement.png"
    notes = [
        "## お知らせバナー",
        "・表示中 / 非表示 バッジ",
        "・種別セレクト:",
        "  お知らせ / 注意 / メンテナンス",
        "・メッセージ テキスト入力",
        "・表示開始日時（datetime-local）",
        "・表示終了日時（datetime-local）",
        "・バリデーション:",
        "  開始 > 終了 → エラートースト",
        "・「クリア」ボタン:",
        "  バナー内容を全クリア",
        "",
        "## メンテナンスモード",
        "・ON にすると管理者以外",
        "  全ページアクセス不可",
        "・メンテナンス開始日時",
        "・メンテナンス終了日時",
        "・バリデーション:",
        "  開始 > 終了 → エラートースト",
        "・警告文表示:",
        "  「管理者以外アクセス不可に",
        "  なります」（Alert コンポーネント）",
    ]
    image_with_notes(slide, img_path, notes)


def slide_17_policy(prs, blank_layout):
    """ポリシー・利用規約タブ"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "設定 — ポリシー・利用規約タブ", badge_text="管理画面")

    two_col_cards(slide,
        "利用規約タブ（terms）",
        [
            "## フィールド",
            "・利用規約テキスト (Textarea)",
            "  管理者が初回ログイン時に同意",
            "",
            "## バージョン管理",
            "・バージョン更新ボタン",
            "  v{current} → v{next}",
            "  → 全管理者が次回ログイン時",
            "  に再同意を求められる",
            "",
            "## 自動生成",
            "・未入力時: 「自動生成」ボタン表示",
            "  院名・電話から自動生成",
            "・「基本情報から再生成」リンク",
            "  手動編集内容上書き警告あり",
            "",
            "## ポリシータブ（policy）",
            "・プライバシーポリシーテキスト",
            "・要配慮個人情報の同意文言 (500文字)",
            "・データ保存期間ポリシー",
            "  （3年推奨 / 無期限〜5年選択肢）",
            "・個人情報の利用目的（APPI第17条）",
            "・リマインダーメール機能 ON/OFF",
        ],
        "患者権利行使（APPI第28〜30条）",
        [
            "## patientRights フィールド",
            "・問い合わせ先・手続き案内 (Textarea)",
            "・予約完了画面・PPに表示",
            "",
            "## 法的根拠",
            "・個人情報保護法 第28条: 開示",
            "・個人情報保護法 第29条: 訂正",
            "・個人情報保護法 第30条: 利用停止",
            "",
            "## 診察履歴の特別扱い",
            "・診察履歴はデータ保存期間対象外",
            "  医療記録として永久保持",
            "  PPに告知する旨の注意文表示",
            "",
            "## リマインダーメール設定",
            "・機能有効/無効トグル",
            "・有効時: 予約フォームに",
            "  「リマインダーメールを受け取る」",
            "  同意チェックが表示される",
            "・HMAC-SHA256 opt-out トークン",
            "  optOutReminder CF で処理",
            "",
            "## スクリーンショット参照",
            "・16_settings_terms.png",
            "・17_settings_policy.png",
        ],
        top=Inches(1.1),
    )


def slide_18_accounts(prs, blank_layout):
    """アカウントタブ"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "設定 — アカウントタブ", badge_text="管理画面")

    img_path = SS_DIR / "18_settings_accounts.png"
    notes = [
        "## 新規ユーザー作成",
        "・メールアドレス入力",
        "  type=email / 必須",
        "  isValidEmail() バリデーション",
        "・パスワード入力",
        "  <PasswordInput> / 必須",
        "  8〜128文字",
        "  大文字・小文字・数字・記号",
        "  各1文字以上（isStrongPassword）",
        "・「作成」ボタン",
        "  → createAdminUser CF 呼び出し",
        "  → mustChangePassword=true 設定",
        "  → 初回ログイン時PW変更強制",
        "・ユーザー上限: 10名まで",
        "",
        "## ユーザー一覧",
        "・メール（マスク表示 [SEC-6]）",
        "  例: ad***@example.com",
        "・作成日",
        "・削除ボタン",
        "  自分自身は削除不可",
        "  → 確認ダイアログあり",
        "  → deleteUser CF 呼び出し",
        "",
        "## パスワード変更",
        "・/admin/change-password",
        "  へのリンク表示",
    ]
    image_with_notes(slide, img_path, notes)


# ── クロージング ──────────────────────────

def slide_closing(prs, blank_layout):
    """クロージング"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.5), W, Inches(1.5), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.08), fill_color=C_GOLD)
    add_text(slide, "Online Appointment System",
             Inches(0.5), Inches(1.5), Inches(12), Inches(0.9),
             font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "管理画面 UI 仕様書  以上",
             Inches(0.5), Inches(2.5), Inches(12), Inches(0.7),
             font_size=26, color=C_BASE, align=PP_ALIGN.CENTER)
    add_text(slide, "oas.kojinius.jp",
             Inches(0.5), Inches(3.5), Inches(12), Inches(0.6),
             font_size=18, color=C_GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# メイン処理
# ─────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank_layout = prs.slide_layouts[6]

    # カバー・目次
    slide_01_cover(prs, blank_layout)
    slide_02_toc(prs, blank_layout)

    # Section 1: 認証・アカウント
    divider_slide(prs, blank_layout, "1. 認証・アカウント", [
        "ログイン画面（メール / パスワード / バリデーション / パスワードトグル）",
        "パスワード変更強制フロー（mustChangePassword / isStrongPassword）",
        "利用規約・PP 同意モーダル（ConsentModal / バージョン管理）",
    ])
    slide_04_login(prs, blank_layout)
    slide_05_change_password(prs, blank_layout)

    # Section 2: ダッシュボード
    divider_slide(prs, blank_layout, "2. ダッシュボード", [
        "KPI カード（本日 / 今月 / 新規 / 未確認）",
        "検索・フィルター（氏名 / 郵便番号 / 電話番号 / 日付 / ステータス）",
        "予約テーブル列定義（8列 / ソート / ステータスバッジ）",
        "予約詳細モーダル（表示フィールド / ステータス変更 / キャンセル / PDF）",
    ])
    slide_07_kpi(prs, blank_layout)
    slide_08_search_filter(prs, blank_layout)
    slide_09_table(prs, blank_layout)
    slide_10_modal(prs, blank_layout)

    # Section 3: 診察履歴
    divider_slide(prs, blank_layout, "3. 診察履歴", [
        "診察履歴ページ（検索 / ページネーション / CSV出力 / 訂正機能）",
        "訂正機能（APPI 第34条 / corrections サブコレクション / immutable）",
    ])
    slide_12_history(prs, blank_layout)

    # Section 4: 設定画面
    divider_slide(prs, blank_layout, "4. 設定画面", [
        "基本情報タブ（院名 / 電話 / URL / 住所 / 郵便番号自動補完 / Google Maps）",
        "営業日設定タブ（タイムラインバー / AM・PM / 休日管理 / キャンセル締切）",
        "お知らせタブ（バナー設定 / メンテナンスモード）",
        "ポリシー・利用規約タブ（PP自動生成 / 要配慮同意文言 / 患者権利窓口）",
        "アカウントタブ（管理者作成 / マスク表示 / パスワード変更）",
    ])
    slide_14_basic_info(prs, blank_layout)
    slide_15_business_days(prs, blank_layout)
    slide_16_announcement(prs, blank_layout)
    slide_17_policy(prs, blank_layout)
    slide_18_accounts(prs, blank_layout)

    # クロージング
    slide_closing(prs, blank_layout)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"生成完了: {OUTPUT_PATH}")
    print(f"スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()

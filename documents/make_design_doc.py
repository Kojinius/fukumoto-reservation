# -*- coding: utf-8 -*-
"""
Online Appointment System 詳細設計書 PPTX 生成スクリプト（v2 - React SPA 対応版）
実行: PYTHONUTF8=1 python documents/make_design_doc.py
     （fukumoto-reservation ディレクトリから）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from PIL import Image as PILImage

# ─────────────────────────────────────────────
# カラー定義（OAS ブランドカラー Bold Navy × Gold）
# ─────────────────────────────────────────────
C_NAVY    = RGBColor(0x1A, 0x2B, 0x4A)   # ダークネイビー（管理者ヘッダー）
C_GOLD    = RGBColor(0xD4, 0xAF, 0x37)   # ゴールド（アクセント）
C_BASE    = RGBColor(0xF8, 0xF5, 0xEE)   # クリーム（背景）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x2B, 0x4A)   # ネイビー（テキスト）
C_MUTED   = RGBColor(0x6B, 0x7C, 0x9E)   # ミュートブルー
C_RED     = RGBColor(0xDC, 0x26, 0x26)   # エラー・警告

# ─────────────────────────────────────────────
# パス設定
# ─────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent.parent
SS_DIR      = BASE_DIR / "documents" / "screenshots"
OUTPUT_PATH = BASE_DIR / "documents" / "specs" / "OnlineAppointSystem_詳細設計書.pptx"

W = Inches(13.33)
H = Inches(7.5)


# ─────────────────────────────────────────────
# 共通ヘルパー関数
# ─────────────────────────────────────────────

def add_shape(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, line_spacing=None, font_name="Noto Sans JP"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def header_bar(slide, title_text, badge_text=None, badge_color=None):
    if badge_color is None:
        badge_color = C_GOLD
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.9), fill_color=C_NAVY)
    add_text(slide, title_text, Inches(0.4), Inches(0.1), Inches(10), Inches(0.75),
             font_size=26, bold=True, color=C_WHITE)
    if badge_text:
        s = add_shape(slide, Inches(11.3), Inches(0.2), Inches(1.8), Inches(0.5),
                      fill_color=badge_color)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge_text
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.name = "Noto Sans JP"
        run.font.color.rgb = C_NAVY


def card(slide, x, y, w, h, bg=None):
    if bg is None:
        bg = C_BASE
    return add_shape(slide, x, y, w, h, fill_color=bg)


def add_image_fit(slide, img_path, x, y, max_w, max_h):
    try:
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        scale = min(max_w / iw, max_h / ih)
        disp_w = int(iw * scale)
        disp_h = int(ih * scale)
        off_x = (max_w - disp_w) // 2
        off_y = (max_h - disp_h) // 2
        slide.shapes.add_picture(str(img_path), x + off_x, y + off_y, disp_w, disp_h)
    except Exception:
        add_shape(slide, x, y, max_w, max_h, fill_color=C_MUTED)


def divider_slide(prs, blank_layout, section_title, items):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.2), W, Inches(1.2), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(2.0), Inches(12.33), Inches(0.06), fill_color=C_GOLD)
    add_text(slide, section_title, Inches(0.5), Inches(1.0), Inches(12), Inches(1.0),
             font_size=40, bold=True, color=C_WHITE)
    y = Inches(2.3)
    for item in items:
        add_text(slide, "  " + item, Inches(0.7), y, Inches(12), Inches(0.5),
                 font_size=16, color=C_BASE)
        y += Inches(0.5)


def three_col_cards(slide, cards_data, top=Inches(1.1)):
    col_w = Inches(3.9)
    gap   = Inches(0.3)
    card_h = Inches(5.8)
    starts = [Inches(0.4), Inches(0.4) + col_w + gap,
              Inches(0.4) + (col_w + gap) * 2]
    for idx, (title, lines) in enumerate(cards_data):
        cx = starts[idx]
        card(slide, cx, top, col_w, card_h)
        add_shape(slide, cx, top, col_w, Inches(0.55), fill_color=C_NAVY)
        add_text(slide, title, cx + Inches(0.1), top + Inches(0.05),
                 col_w - Inches(0.2), Inches(0.5),
                 font_size=15, bold=True, color=C_WHITE)
        body_text = "\n".join(lines)
        add_text(slide, body_text, cx + Inches(0.15), top + Inches(0.65),
                 col_w - Inches(0.3), card_h - Inches(0.75),
                 font_size=13, color=C_DARK, line_spacing=Pt(20))


def two_col_cards(slide, left_title, left_lines, right_title, right_lines,
                  top=Inches(1.1)):
    col_w = Inches(5.9)
    gap   = Inches(0.4)
    card_h = Inches(5.8)
    left_x  = Inches(0.5)
    right_x = left_x + col_w + gap
    for cx, title, lines in [(left_x, left_title, left_lines),
                              (right_x, right_title, right_lines)]:
        card(slide, cx, top, col_w, card_h)
        add_shape(slide, cx, top, col_w, Inches(0.55), fill_color=C_NAVY)
        add_text(slide, title, cx + Inches(0.1), top + Inches(0.05),
                 col_w - Inches(0.2), Inches(0.5),
                 font_size=15, bold=True, color=C_WHITE)
        body_text = "\n".join(lines)
        add_text(slide, body_text, cx + Inches(0.15), top + Inches(0.65),
                 col_w - Inches(0.3), card_h - Inches(0.75),
                 font_size=13, color=C_DARK, line_spacing=Pt(20))


def image_with_notes(slide, img_path, notes_lines, top=Inches(1.1)):
    img_x, img_y = Inches(0.4), top
    img_max_w, img_max_h = Inches(8.0), Inches(5.8)
    note_x = Inches(8.7)
    note_w = Inches(4.3)
    if img_path and Path(img_path).exists():
        add_image_fit(slide, img_path, img_x, img_y, img_max_w, img_max_h)
    else:
        card(slide, img_x, img_y, img_max_w, img_max_h, bg=C_MUTED)
    card(slide, note_x, top, note_w, Inches(5.8))
    body_text = "\n".join(notes_lines)
    add_text(slide, body_text, note_x + Inches(0.15), top + Inches(0.2),
             note_w - Inches(0.3), Inches(5.4),
             font_size=13, color=C_DARK, line_spacing=Pt(22))


# ─────────────────────────────────────────────
# スライド生成関数群
# ─────────────────────────────────────────────

def slide_01_cover(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.5), W, Inches(1.5), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(2.8), Inches(12.33), Inches(0.08), fill_color=C_GOLD)
    add_text(slide, "Online Appointment System",
             Inches(0.5), Inches(1.2), Inches(12), Inches(1.1),
             font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "詳細設計書",
             Inches(0.5), Inches(3.0), Inches(12), Inches(0.9),
             font_size=36, bold=True, color=C_BASE, align=PP_ALIGN.CENTER)
    add_text(slide, "2026年3月  |  React SPA v4（i18n / PWA / 問診票 / 診察履歴 / APPI対応）",
             Inches(0.5), Inches(3.9), Inches(12), Inches(0.6),
             font_size=18, color=C_BASE, align=PP_ALIGN.CENTER)


def slide_02_toc(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "目次")

    toc_items = [
        ("1",  "システム概要",               "p.4"),
        ("2",  "技術スタック",               "p.5"),
        ("3",  "ディレクトリ構成",           "p.7"),
        ("4",  "Firestore データモデル",     "p.8"),
        ("5",  "Cloud Functions 設計",       "p.9"),
        ("6",  "予約作成フロー",             "p.11"),
        ("7",  "認証・認可設計",             "p.13"),
        ("8",  "Firestore セキュリティルール", "p.14"),
        ("9",  "XSS・入力バリデーション",    "p.15"),
        ("10", "レート制限・監査ログ",        "p.16"),
        ("11", "患者向け予約フロー",          "p.18"),
        ("12", "管理者ダッシュボード",        "p.19"),
        ("13", "設定画面",                   "p.20"),
        ("14", "デプロイ構成",               "p.22"),
        ("15", "i18n 多言語対応",            "p.24"),
        ("16", "PWA 対応",                   "p.25"),
        ("17", "問診票機能",                 "p.26"),
        ("18", "診察履歴・APPI対応",         "p.27"),
        ("19", "アクセスログ・リマインダー", "p.28"),
    ]

    left_items  = toc_items[:10]
    right_items = toc_items[10:]

    for col_idx, items in enumerate([left_items, right_items]):
        cx = Inches(0.5) + Inches(6.3) * col_idx
        y  = Inches(1.1)
        for num, title, page in items:
            badge = add_shape(slide, cx, y + Inches(0.05),
                              Inches(0.5), Inches(0.4), fill_color=C_NAVY)
            tf = badge.text_frame
            p  = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = num
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.name = "Noto Sans JP"
            run.font.color.rgb = C_WHITE
            add_text(slide, title, cx + Inches(0.6), y,
                     Inches(4.5), Inches(0.5), font_size=14, color=C_DARK)
            add_text(slide, page, cx + Inches(5.3), y,
                     Inches(0.8), Inches(0.5), font_size=13, color=C_MUTED,
                     align=PP_ALIGN.RIGHT)
            add_shape(slide, cx, y + Inches(0.5), Inches(6.0), Inches(0.02), fill_color=C_MUTED)
            y += Inches(0.52)


def slide_04_overview(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "システム概要", badge_text="概要")
    three_col_cards(slide, [
        ("システム目的", [
            "・鍼灸整骨院向けオンライン予約",
            "・24時間365日受付対応",
            "・管理者はリアルタイムで",
            "　予約管理可能",
            "・患者・院側双方の",
            "　利便性向上",
        ]),
        ("主要機能", [
            "・患者向け予約フロー",
            "　（4ステップ: 日時→情報入力",
            "　→確認→完了）",
            "・予約管理ダッシュボード",
            "・設定管理（5タブ）",
            "・ユーザー管理",
            "・パスワード変更強制",
        ]),
        ("非機能要件", [
            "・セキュリティ",
            "　（12脆弱性修正済み）",
            "・CSP ヘッダー完全実装",
            "・レート制限（DoS対策）",
            "・モバイル対応",
            "・日本語フォント対応",
            "・SSL/TLS（Firebase）",
        ]),
    ])


def slide_05_tech_stack(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "技術スタック", badge_text="技術")
    three_col_cards(slide, [
        ("フロントエンド (oas-spa/)", [
            "・Vite 6",
            "・React 19",
            "・TypeScript 5",
            "・Tailwind CSS 3.4",
            "・React Router 7",
            "・Firebase v10 SDK",
            "・pdf-lib / fontkit",
            "　（PDF 生成）",
            "・clsx / tailwind-merge",
            "・react-i18next（7言語対応）",
            "・vite-plugin-pwa（PWA）",
        ]),
        ("バックエンド", [
            "・Firebase Authentication",
            "　（Email/Password）",
            "・Cloud Firestore",
            "　（asia-northeast1）",
            "・Cloud Functions",
            "　（Node.js 24）",
            "・firebase-admin v13",
            "・Resend（メール送信）",
        ]),
        ("インフラ・ツール", [
            "・Firebase Hosting",
            "　（マルチサイト）",
            "・カスタムドメイン:",
            "　oas.kojinius.jp",
            "・さくらインターネット DNS",
            "・Playwright（E2Eテスト）",
            "・python-pptx（設計書生成）",
        ]),
    ])


def slide_07_directory(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "ディレクトリ構成", badge_text="構成")
    card(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.8))

    tree_text = (
        "fukumoto-reservation/\n"
        "├── oas-spa/                    # React SPA（患者向け + 管理者）\n"
        "│   ├── src/\n"
        "│   │   ├── pages/\n"
        "│   │   │   ├── booking/        # Index / DateTimeSelect / PatientForm\n"
        "│   │   │   │                   #   / Confirm / Complete\n"
        "│   │   │   ├── admin/          # Dashboard / Settings / ChangePassword / History\n"
        "│   │   │   ├── questionnaire/  # /questionnaire（問診票）\n"
        "│   │   │   ├── Cancel.tsx      # /cancel\n"
        "│   │   │   ├── Login.tsx       # /login\n"
        "│   │   │   └── PrivacyPolicy.tsx\n"
        "│   │   ├── components/         # ui / shared / layout\n"
        "│   │   ├── contexts/           # Auth / Clinic / Theme / Toast\n"
        "│   │   ├── hooks/              # useAuth / useAdmin / useReservation etc.\n"
        "│   │   ├── i18n/               # react-i18next 設定 / 7言語 / 5ネームスペース\n"
        "│   │   │   └── locales/{ja,en,zh,ko,vi,pt,es}/\n"
        "│   │   ├── types/              # TypeScript 型定義（questionnaire.ts 追加）\n"
        "│   │   └── utils/              # cn / date / validation / security / zip\n"
        "│   │                           #   questionnairePdf.ts（pdf-lib PDF生成）\n"
        "│   └── dist/                   # ビルド出力\n"
        "├── functions/index.js          # Cloud Functions (10+ functions)\n"
        "├── firestore.rules             # セキュリティルール\n"
        "└── firebase.json               # Hosting マルチサイト設定"
    )
    add_text(slide, tree_text, Inches(0.8), Inches(1.3), Inches(11.8), Inches(5.5),
             font_size=13, color=C_DARK, font_name="Courier New", line_spacing=Pt(21))


def slide_08_datamodel(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "Firestore データモデル", badge_text="DB")
    two_col_cards(slide,
        "コレクション一覧",
        [
            "reservations",
            "　id, name, furigana, date, time",
            "　phone, status, createdAt",
            "　email?, zip?, address?",
            "　symptoms?, visitType?,",
            "　insurance?, gender?, birthdate?",
            "",
            "slots",
            "　slotId (date_time 形式)",
            "　reservationId（参照）",
            "",
            "settings / clinic",
            "　clinicName, address, phone",
            "　businessHours, holidays",
            "　announcement",
            "　cancelCutoffMinutes",
            "",
            "users / {uid}",
            "　mustChangePassword: boolean",
            "",
            "audit_logs",
            "　event, data, severity, timestamp",
            "",
            "questionnaires / {reservationId}",
            "　18フィールド / 問診票データ",
            "",
            "visit_histories / {id}",
            "　（immutable）診察完了記録",
            "　→ corrections/ サブコレクション",
            "",
            "access_logs",
            "　uid, action, targetId, timestamp",
        ],
        "status・追加コレクション詳細",
        [
            "reservations.status",
            "・pending    → 未確認",
            "・confirmed  → 確認済み",
            "・cancelled  → キャンセル済",
            "・completed  → 診察完了",
            "",
            "visit_histories フィールド",
            "・completedAt, completedBy",
            "・（予約データをコピー保存）",
            "・immutable rules:",
            "  update/delete 管理者も不可",
            "  correctVisitHistory CF のみ",
            "  corrections サブコレクションに追記",
            "",
            "audit_logs イベント",
            "・reservation.created",
            "・reservation.slot_taken",
            "・reservation.cancelled",
            "・user.created",
            "・user.deleted",
            "・rate_limit.exceeded",
            "",
            "access_logs アクション",
            "・view_reservation",
            "・download_questionnaire_pdf",
            "・complete_visit",
        ],
    )


def slide_09_functions(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "Cloud Functions 設計（10 functions）", badge_text="Functions")

    funcs = [
        ("createReservation",    "onRequest", "予約作成・メール送信・スロット管理・監査ログ。レート制限・スキーマ検証あり"),
        ("verifyReservation",    "onRequest", "電話番号照合でキャンセル前の予約照会。電話番号ミスマッチは404で列挙攻撃防止"),
        ("cancelReservation",    "onRequest", "電話番号照合後 Firestore トランザクションでキャンセル。TOCTOU 修正済み [SEC-8]"),
        ("createAdminUser",      "onRequest", "管理者作成・admin カスタムクレーム付与・初期パスワード変更フラグ設定"),
        ("listUsers",            "onRequest", "管理者一覧取得（admin クレーム必須）。メールはマスク表示 [SEC-6]"),
        ("deleteUser",           "onRequest", "ユーザー削除（admin クレーム必須）"),
        ("sendDailyReminders",   "onSchedule","翌日予約へのリマインダーメール一括送信（opt-in 患者のみ）"),
        ("completeVisit",        "onRequest", "診察完了処理: visit_histories 保存 + status=completed。admin クレーム必須"),
        ("correctVisitHistory",  "onRequest", "診察履歴訂正（APPI第34条）: corrections サブコレクションに immutable 追記"),
        ("optOutReminder",       "onRequest", "HMAC-SHA256 トークン検証後 reminderConsent=false に更新（メール配信停止）"),
    ]

    y = Inches(1.1)
    row_h = Inches(0.72)
    for func_name, kind, desc in funcs:
        add_shape(slide, Inches(0.5), y, Inches(3.2), row_h, fill_color=C_NAVY)
        add_text(slide, func_name, Inches(0.6), y + Inches(0.08),
                 Inches(3.0), Inches(0.35), font_size=13, bold=True, color=C_WHITE)
        add_text(slide, kind, Inches(0.6), y + Inches(0.42),
                 Inches(3.0), Inches(0.26), font_size=10, color=C_GOLD)
        card(slide, Inches(3.7), y, Inches(9.2), row_h)
        add_text(slide, desc, Inches(3.85), y + Inches(0.15),
                 Inches(9.0), Inches(0.5), font_size=13, color=C_DARK)
        y += row_h + Inches(0.05)


def slide_11_reservation_flow(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "予約作成フロー（createReservation）", badge_text="フロー")

    steps = [
        ("①", "レート制限チェック",       "5 req/min per IP。超過時は 429 を返す"),
        ("②", "入力バリデーション",        "必須フィールド・型・長さ・電話番号・メール形式を検証"),
        ("③", "スロット競合チェック",      "slots/{date_time} を transaction.get で確認。取得済みなら SLOT_TAKEN"),
        ("④", "reservations ドキュメント作成", "Firestore transaction 内で atomic に作成"),
        ("⑤", "スロット予約",             "slots/{date_time} に reservationId を書き込み（同一 transaction）"),
        ("⑥", "メール送信（非同期）",      "Resend API で患者・院双方に確認メール送信"),
        ("⑦", "監査ログ記録",             "reservation.created イベントを audit_logs に保存"),
    ]

    y = Inches(1.1)
    for num, title, desc in steps:
        add_shape(slide, Inches(0.5), y, Inches(0.6), Inches(0.62), fill_color=C_GOLD)
        add_text(slide, num, Inches(0.5), y + Inches(0.08),
                 Inches(0.6), Inches(0.5), font_size=16, bold=True,
                 color=C_NAVY, align=PP_ALIGN.CENTER)
        add_shape(slide, Inches(1.15), y, Inches(3.0), Inches(0.62), fill_color=C_NAVY)
        add_text(slide, title, Inches(1.25), y + Inches(0.12),
                 Inches(2.8), Inches(0.45), font_size=13, bold=True, color=C_WHITE)
        card(slide, Inches(4.2), y, Inches(8.7), Inches(0.62))
        add_text(slide, desc, Inches(4.35), y + Inches(0.14),
                 Inches(8.5), Inches(0.42), font_size=13, color=C_DARK)
        y += Inches(0.72)


def slide_13_auth(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "認証・認可設計", badge_text="Auth")

    content = (
        "## 認証方式\n"
        "・Firebase Authentication（Email / Password）\n"
        "・管理者: admin カスタムクレーム（server-side 付与のみ）\n"
        "・セッション: browserSessionPersistence（タブ閉じでセッション消滅）\n"
        "\n"
        "## 認可ルール\n"
        "・患者向けページ（/、/cancel、/privacy-policy）: 認証不要\n"
        "・管理者ページ（/admin、/admin/settings）: admin クレーム必須\n"
        "・Firestore Rules: isAdmin() カスタム関数で厳密チェック\n"
        "\n"
        "## 初期パスワード強制変更\n"
        "・createAdminUser 時に mustChangePassword=true を Firestore に保存\n"
        "・AdminLayout: mustChangePassword=true なら /admin/change-password?forced=1 に強制リダイレクト\n"
        "・パスワード変更完了後に mustChangePassword=false に更新（本人のみ許可）\n"
        "\n"
        "## パスワードリセット\n"
        "・sendPasswordResetEmail で再設定リンク送信\n"
        "・/auth-action でメール認証アクション処理"
    )
    card(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.8))
    add_text(slide, content, Inches(0.8), Inches(1.3), Inches(12.0), Inches(5.5),
             font_size=13, color=C_DARK, line_spacing=Pt(23))


def slide_14_firestore_rules(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "Firestore セキュリティルール", badge_text="Security")
    card(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.8), bg=C_DARK)

    rules_text = (
        "// デフォルト拒否（全ルートを安全側に倒す）\n"
        "match /{document=**} { allow read, write: if false; }\n"
        "\n"
        "// isAdmin(): admin カスタムクレームで判定\n"
        "function isAdmin() {\n"
        "  return request.auth != null && request.auth.token.admin == true;\n"
        "}\n"
        "\n"
        "// isValidReservation(): 必須フィールド・型・長さを検証\n"
        "match /reservations/{id} {\n"
        "  allow read, update, delete: if isAdmin();\n"
        "  allow create: if isAdmin() && isValidReservation(request.resource.data);\n"
        "  // 未認証クライアントからの直接書き込みは全て拒否 [SEC-18]\n"
        "}\n"
        "\n"
        "match /slots/{slotId} {\n"
        "  allow read: if true;   // カレンダー表示のため公開\n"
        "  allow write: if isAdmin();\n"
        "}\n"
        "\n"
        "match /users/{userId} {\n"
        "  allow read: if request.auth.uid == userId;  // 本人のみ\n"
        "  allow write: if isAdmin();\n"
        "  allow update: if request.auth.uid == userId\n"
        "    && request.resource.data.diff(resource.data)\n"
        "         .affectedKeys().hasOnly(['mustChangePassword']);\n"
        "}\n"
        "\n"
        "match /audit_logs/{logId} { allow read, create: if isAdmin(); }\n"
        "match /settings/{settingId} { allow read: if true; allow write: if isAdmin(); }"
    )
    add_text(slide, rules_text, Inches(0.8), Inches(1.25), Inches(12.0), Inches(5.5),
             font_size=11.5, color=C_BASE, font_name="Courier New", line_spacing=Pt(18))


def slide_15_xss(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "XSS・入力バリデーション", badge_text="Security")

    two_col_cards(slide,
        "クライアントサイド対策",
        [
            "## JSX 自動エスケープ",
            "・React JSX は {value} を自動エスケープ",
            "・innerHTML 直接操作を排除",
            "",
            "## 入力バリデーション (utils/validation.ts)",
            "・isStrongPassword(): 4 カテゴリ必須",
            "  (大文字・小文字・数字・記号)",
            "・isValidEmail(): RFC 準拠",
            "  連続ドット・末尾ドット拒否",
            "・toHankaku(): 全角→半角正規化",
            "  (utils/security.ts)",
            "",
            "## 郵便番号 API 検証 [SEC-12]",
            "・Array.isArray() チェック",
            "・各フィールドは string && slice(0,50)",
        ],
        "サーバーサイド・インフラ対策",
        [
            "## Cloud Functions バリデーション",
            "・validatePasswordComplexity()",
            "・validateEmail()",
            "・全 Cloud Function で実施",
            "",
            "## CSP ヘッダー [SEC-15]",
            "・firebase.json で全ルールを設定",
            "・script-src: *.firebase*, CDN のみ",
            "・connect-src: *.run.app,",
            "  zipcloud, holidays-jp",
            "・frame-src: Google Maps",
            "・X-Frame-Options: SAMEORIGIN",
            "・X-Content-Type-Options: nosniff",
            "・Strict-Transport-Security",
            "",
            "## 監査ログ [SEC-14]",
            "・auditLog() でセキュリティ",
            "  イベントを GCP Cloud Logging へ",
        ],
        top=Inches(1.1),
    )


def slide_16_rate_limit(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "レート制限・監査ログ・その他修正", badge_text="Security")

    content = (
        "## レート制限 [SEC-11]\n"
        "・全 Cloud Function に IP ベースのレート制限を実装\n"
        "・5 req/min per IP（60 秒スライディングウィンドウ・インメモリ Map）\n"
        "・超過時: 429 Too Many Requests を返し rate_limit.exceeded イベントをログ\n"
        "\n"
        "## キャンセル競合防止 TOCTOU 修正 [SEC-8]\n"
        "・cancelReservation は Firestore transaction 内で再度ステータスを確認\n"
        "・同時キャンセルリクエスト時の二重キャンセルを防止\n"
        "\n"
        "## その他セキュリティ修正\n"
        "・[SEC-3]  予約照会を verifyReservation CF 経由に移行（電話番号照合必須）\n"
        "・[SEC-4]  createAdminUser: Firestore 書き込み失敗時は Auth ユーザーもロールバック\n"
        "・[SEC-5]  キャンセルを cancelReservation CF 経由に移行（未認証 update 廃止）\n"
        "・[SEC-6]  管理者 UI でメールアドレスをマスク表示\n"
        "・[SEC-10] window グローバル関数を廃止（addEventListener に移行）\n"
        "・[SEC-13] Firebase API キー — Won't Fix（公式設計・ドメイン/API 制限済み）"
    )
    card(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.8))
    add_text(slide, content, Inches(0.8), Inches(1.3), Inches(12.0), Inches(5.5),
             font_size=13, color=C_DARK, line_spacing=Pt(22))


def slide_18_patient_flow(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "患者向け予約フロー（4ステップ）", badge_text="画面")
    notes = [
        "## Step 1: 日時選択",
        "・カレンダーで希望日を選択",
        "・利用可能な時間枠をクリック",
        "",
        "## Step 2: 情報入力",
        "・氏名・ふりがな・電話番号",
        "・郵便番号→住所自動補完",
        "・初診/再診・保険有無",
        "",
        "## Step 3: 確認",
        "・入力内容の最終確認",
        "",
        "## Step 4: 完了",
        "・予約ID 発行",
        "・確認メール自動送信",
    ]
    image_with_notes(slide, SS_DIR / "01_calendar_top.png", notes)


def slide_19_dashboard(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "管理者ダッシュボード（/admin）", badge_text="画面")
    notes = [
        "## KPI カード（4種）",
        "・本日の予約数",
        "・今月の予約数",
        "・新規患者数",
        "・未確認件数",
        "",
        "## 予約一覧",
        "・ステータス・日付・キーワード",
        "　・郵便番号・電話番号で絞り込み",
        "・行クリックで詳細モーダル",
        "・ソート可（複数列対応）",
        "",
        "## 詳細モーダル",
        "・ステータス変更",
        "・診察完了（visit_histories 保存）",
        "・問診票PDF ダウンロード（pdf-lib）",
        "・キャンセル（理由選択必須）",
        "",
        "## アクセスログ",
        "・詳細モーダル閲覧: logAccess()",
        "・問診票PDF DL: logAccess()",
        "",
        "## エクスポート",
        "・CSV ダウンロード",
    ]
    image_with_notes(slide, SS_DIR / "10_dashboard.png", notes)


def slide_20_settings(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "設定画面（/admin/settings）", badge_text="画面")
    notes = [
        "## タブ構成（6 タブ）",
        "・基本情報",
        "　院名・住所・電話番号・Google Maps",
        "・営業日設定",
        "　曜日別 ON/OFF + AM/PM 時間設定",
        "　タイムラインバー・休日管理",
        "・お知らせ",
        "　バナー設定・メンテナンスモード",
        "・利用規約",
        "　規約テキスト・バージョン管理",
        "・ポリシー",
        "　PP・要配慮同意・保存期間",
        "　患者権利窓口（APPI 第28〜30条）",
        "　リマインダーメール ON/OFF",
        "・アカウント",
        "　管理者一覧・追加・削除",
        "　パスワード変更リンク",
        "",
        "## 保存",
        "・「保存」ボタンで Firestore に",
        "　即時反映（merge: true）",
    ]
    image_with_notes(slide, SS_DIR / "13_settings_basic_info.png", notes)


def slide_24_i18n(prs, blank_layout):
    """i18n 多言語対応"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "i18n 多言語対応", badge_text="新機能")
    two_col_cards(slide,
        "構成・対応言語",
        [
            "## ライブラリ",
            "・react-i18next",
            "・i18next-browser-languagedetector",
            "",
            "## 対応言語（7言語）",
            "・ja  日本語（デフォルト）",
            "・en  英語",
            "・zh  中国語（簡体字）",
            "・ko  韓国語",
            "・vi  ベトナム語",
            "・pt  ポルトガル語",
            "・es  スペイン語",
            "",
            "## ネームスペース（5種）",
            "・common   共通",
            "・admin    管理者画面",
            "・booking  予約フロー",
            "・toast    通知",
            "・questionnaire  問診票",
            "",
            "## LanguageSwitcher",
            "・ヘッダーに variant=header で表示",
            "・患者・管理者両方のヘッダーに配置",
        ],
        "実装詳細",
        [
            "## 初期化（i18n/index.ts）",
            "・I18nextProvider でアプリをラップ",
            "・ブラウザ言語を自動検出",
            "・fallbackLng: ja",
            "",
            "## 翻訳ファイル構成",
            "src/i18n/locales/",
            "  {lang}/",
            "    common.json",
            "    admin.json",
            "    booking.json",
            "    toast.json",
            "    questionnaire.json",
            "",
            "## 使い方",
            "const { t } = useTranslation('admin');",
            "t('dashboard.kpi.todayReservations')",
            "",
            "## PDF出力での多言語対応",
            "・問診票PDFラベルは t() で取得",
            "・管理者の表示言語に応じて出力",
        ],
        top=Inches(1.1),
    )


def slide_25_pwa(prs, blank_layout):
    """PWA 対応"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "PWA 対応", badge_text="新機能")
    two_col_cards(slide,
        "技術構成",
        [
            "## ライブラリ",
            "・vite-plugin-pwa",
            "・Workbox（Service Worker 生成）",
            "",
            "## vite.config.ts 設定",
            "・VitePWA プラグイン追加",
            "・registerType: autoUpdate",
            "・workbox:",
            "  globPatterns: ['**/*.{js,css,html}']",
            "  runtimeCaching で API キャッシュ",
            "",
            "## Web App Manifest",
            "・name: OAS 予約システム",
            "・short_name: OAS",
            "・display: standalone",
            "・theme_color: #1A2B4A（ネイビー）",
            "・background_color: #F8F5EE（クリーム）",
            "・icons: 192px / 512px",
        ],
        "オフライン・インストール",
        [
            "## オフラインキャッシュ戦略",
            "・Shell (HTML/CSS/JS): precache",
            "・静的アセット: CacheFirst",
            "  最大30日 / 60エントリ",
            "・Firebase API: NetworkFirst",
            "  オフライン時はキャッシュから提供",
            "",
            "## InstallBanner コンポーネント",
            "・beforeinstallprompt イベントを",
            "  キャプチャして表示",
            "・「ホーム画面に追加」ボタン",
            "・一度閉じると localStorage に",
            "  dismissed=true を保存",
            "",
            "## Service Worker 更新",
            "・autoUpdate: バックグラウンドで",
            "  自動更新・即時適用",
            "・skipWaiting: true",
        ],
        top=Inches(1.1),
    )


def slide_26_questionnaire(prs, blank_layout):
    """問診票機能"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "問診票機能（/questionnaire）", badge_text="新機能")

    img_path = SS_DIR / "08_questionnaire.png"
    notes = [
        "## ルート",
        "・/questionnaire?id={reservationId}",
        "  予約確認メール内リンクから遷移",
        "",
        "## 18フィールド（5セクション）",
        "《主訴・症状》",
        "・主訴 / 発症時期 / 痛み程度",
        "・痛みの部位（複数選択）",
        "・痛みの性質 / 日常生活への影響",
        "《既往歴・治療状況》",
        "・既往歴 / 現在の通院",
        "・服用中の薬 / アレルギー / 手術歴",
        "《生活習慣》",
        "・睡眠 / 食事 / 運動 / ストレス",
        "《その他》",
        "・妊娠の可能性（3択）",
        "・希望する施術 / その他",
        "",
        "## データ保存",
        "・questionnaires/{reservationId}",
        "  に Firestore 保存",
        "",
        "## PDF出力",
        "・generateQuestionnairePdf()",
        "・pdf-lib + fontkit",
        "・ファイル名:",
        "  問診票_{氏名}_{日付}.pdf",
    ]
    image_with_notes(slide, img_path, notes)


def slide_27_visit_history(prs, blank_layout):
    """診察履歴・APPI対応"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "診察履歴・APPI訂正権対応", badge_text="新機能")
    two_col_cards(slide,
        "診察履歴（completeVisit CF）",
        [
            "## 診察完了フロー",
            "① 管理者が詳細モーダルで",
            "  「診察完了」ボタンをクリック",
            "② 確認ダイアログ表示",
            "③ completeVisit CF 呼び出し",
            "   （admin クレーム必須）",
            "④ visit_histories/{id} に",
            "   予約データをコピー保存",
            "⑤ reservations.status",
            "   = completed に更新",
            "",
            "## visit_histories コレクション",
            "・completedAt: 診察完了日時",
            "・completedBy: 管理者UID",
            "・予約全フィールドをコピー",
            "・immutable Firestore ルール:",
            "  管理者も update/delete 不可",
            "  （CF のみ書き込み可）",
            "",
            "## 診察履歴ページ",
            "・/admin/history",
            "・検索・ページネーション（20件）",
            "・CSV 出力",
            "・詳細モーダル",
        ],
        "APPI 訂正権（第34条）",
        [
            "## correctVisitHistory CF",
            "・admin クレーム必須",
            "・訂正理由必須入力",
            "・corrections サブコレクション",
            "  に immutable 追記",
            "  （元レコードは更新しない）",
            "",
            "## 訂正可能フィールド",
            "・氏名 / ふりがな",
            "・生年月日 / 郵便番号",
            "・住所 / 電話番号 / メール",
            "・初診/再診 / 保険証",
            "・性別",
            "",
            "## 訂正不可フィールド",
            "・症状（施術者判断記録）",
            "・診察日時 / 完了者",
            "・伝達事項",
            "",
            "## 追記（addendum）",
            "・症状等は直接訂正不可",
            "・addendum フィールドで",
            "  事実誤記の補足説明可能",
            "",
            "## 通知",
            "・患者メールアドレスがあれば",
            "  訂正完了メール送信",
            "・未登録時は手動通知を促す",
        ],
        top=Inches(1.1),
    )


def slide_28_access_log_reminder(prs, blank_layout):
    """アクセスログ・リマインダーopt-out"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "アクセスログ・リマインダーopt-out", badge_text="新機能")
    two_col_cards(slide,
        "アクセスログ（logAccess）",
        [
            "## 目的",
            "・個人情報アクセスの証跡管理",
            "・APPI 対応・内部統制",
            "",
            "## ログ記録タイミング",
            "・詳細モーダル閲覧",
            "  action: view_reservation",
            "  targetId: reservationId",
            "  data: { patientName }",
            "",
            "・問診票PDFダウンロード",
            "  action: download_questionnaire_pdf",
            "  targetId: reservationId",
            "  data: { patientName }",
            "",
            "・診察完了処理",
            "  action: complete_visit",
            "  （CF側でも記録）",
            "",
            "## コレクション",
            "access_logs",
            "・uid: 操作者UID",
            "・action: アクション種別",
            "・targetId: 対象ドキュメントID",
            "・timestamp: 操作日時",
            "・data: 付帯情報",
            "",
            "## Firestore ルール",
            "・admin のみ read/create 可",
        ],
        "リマインダーメール opt-out",
        [
            "## フロー",
            "① 予約確認メールに",
            "  opt-out URL を記載",
            "  （HMAC-SHA256 トークン付き）",
            "② 患者がリンクをクリック",
            "③ optOutReminder CF に",
            "  GET リクエスト",
            "④ トークン検証",
            "  HMAC-SHA256(reservationId,",
            "  SECRET_KEY) と照合",
            "⑤ 一致した場合:",
            "  reservations/{id}",
            "  .reminderConsent = false",
            "  に更新",
            "",
            "## セキュリティ",
            "・トークンは URL-safe Base64",
            "・SECRET_KEY は",
            "  Firebase Secret Manager で管理",
            "・トークン不一致: 403 Forbidden",
            "",
            "## opt-in 設定（Settings）",
            "・ポリシータブで機能 ON/OFF",
            "・ON 時: 予約フォームに",
            "  「リマインダーメールを受け取る」",
            "  チェックボックスを表示",
            "・患者が同意した場合のみ",
            "  sendDailyReminders で送信",
        ],
        top=Inches(1.1),
    )


def slide_22_deploy(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "デプロイ構成", badge_text="運用")

    content = (
        "## Firebase Hosting マルチサイト\n"
        "・portfolio: kojinius.jp\n"
        "・oas: oas.kojinius.jp  ← OAS SPA の配信先\n"
        "\n"
        "## デプロイコマンド（スコープ指定を厳守）\n"
        "# Firestore ルール\n"
        "firebase deploy --only firestore:rules --project project-3040e21e-879f-4c66-a7d\n"
        "\n"
        "# Cloud Functions\n"
        "firebase deploy --only functions --project project-3040e21e-879f-4c66-a7d\n"
        "\n"
        "# Hosting (OAS SPA + portfolio)\n"
        "firebase deploy --only hosting:oas,hosting:portfolio --project project-3040e21e-879f-4c66-a7d\n"
        "\n"
        "## レガシー URL リダイレクト\n"
        "・旧 MPA URL（/index.html 等）→ SPA ルート（/）に 301 リダイレクト\n"
        "・クライアントへ旧 URL を案内済みの場合も自動転送\n"
        "\n"
        "## Firebase プロジェクト\n"
        "・project-3040e21e-879f-4c66-a7d（AMS と共有）"
    )
    card(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.8))
    add_text(slide, content, Inches(0.8), Inches(1.3), Inches(12.0), Inches(5.5),
             font_size=12.5, color=C_DARK, font_name="Courier New", line_spacing=Pt(21))


def slide_23_closing(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.5), W, Inches(1.5), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.08), fill_color=C_GOLD)
    add_text(slide, "Online Appointment System",
             Inches(0.5), Inches(1.5), Inches(12), Inches(0.9),
             font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "詳細設計書  以上",
             Inches(0.5), Inches(2.5), Inches(12), Inches(0.7),
             font_size=26, color=C_BASE, align=PP_ALIGN.CENTER)
    add_text(slide, "oas.kojinius.jp",
             Inches(0.5), Inches(3.5), Inches(12), Inches(0.6),
             font_size=18, color=C_GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# メイン処理
# ─────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank_layout = prs.slide_layouts[6]

    slide_01_cover(prs, blank_layout)
    slide_02_toc(prs, blank_layout)

    divider_slide(prs, blank_layout, "1. システム概要", [
        "システム目的・主要機能・非機能要件",
        "技術スタック（React SPA v4 + Firebase + Node.js 24 + i18n/PWA）",
    ])
    slide_04_overview(prs, blank_layout)
    slide_05_tech_stack(prs, blank_layout)

    divider_slide(prs, blank_layout, "2. アーキテクチャ設計", [
        "ディレクトリ構成（oas-spa/ SPA 構造 / i18n / questionnaire追加）",
        "Firestore データモデル（8 コレクション / visit_histories / access_logs）",
        "Cloud Functions 設計（10 functions）",
        "予約作成フロー（createReservation）",
    ])
    slide_07_directory(prs, blank_layout)
    slide_08_datamodel(prs, blank_layout)
    slide_09_functions(prs, blank_layout)
    slide_11_reservation_flow(prs, blank_layout)

    divider_slide(prs, blank_layout, "3. セキュリティ設計", [
        "認証・認可設計（admin クレーム・初期パスワード強制変更）",
        "Firestore セキュリティルール（isAdmin / isValidReservation）",
        "XSS・入力バリデーション（SEC-12/15）",
        "レート制限・監査ログ・その他修正（12 脆弱性対応済み）",
    ])
    slide_13_auth(prs, blank_layout)
    slide_14_firestore_rules(prs, blank_layout)
    slide_15_xss(prs, blank_layout)
    slide_16_rate_limit(prs, blank_layout)

    divider_slide(prs, blank_layout, "4. 画面設計", [
        "患者向け予約フロー（4 ステップ）",
        "管理者ダッシュボード（KPI / 検索 / 詳細モーダル / completeVisit）",
        "設定画面（6 タブ / ポリシー・APPI対応）",
    ])
    slide_18_patient_flow(prs, blank_layout)
    slide_19_dashboard(prs, blank_layout)
    slide_20_settings(prs, blank_layout)

    divider_slide(prs, blank_layout, "5. 新機能（v4 エンハンス）", [
        "i18n 多言語対応（react-i18next / 7言語 / 5ネームスペース）",
        "PWA 対応（vite-plugin-pwa / Workbox / offline cache / InstallBanner）",
        "問診票機能（/questionnaire / 18フィールド / pdf-lib PDF出力）",
        "診察履歴・APPI訂正権（completeVisit CF / visit_histories / corrections）",
        "アクセスログ・リマインダーopt-out（logAccess / HMAC-SHA256）",
    ])
    slide_24_i18n(prs, blank_layout)
    slide_25_pwa(prs, blank_layout)
    slide_26_questionnaire(prs, blank_layout)
    slide_27_visit_history(prs, blank_layout)
    slide_28_access_log_reminder(prs, blank_layout)

    divider_slide(prs, blank_layout, "6. 運用・デプロイ", [
        "Firebase Hosting マルチサイト構成",
        "デプロイコマンド（--only スコープ厳守）",
        "レガシー URL リダイレクト",
    ])
    slide_22_deploy(prs, blank_layout)
    slide_23_closing(prs, blank_layout)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"生成完了: {OUTPUT_PATH}")
    print(f"スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()

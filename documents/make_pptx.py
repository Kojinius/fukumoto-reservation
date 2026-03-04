"""
福元鍼灸整骨院 本番運用ロードマップ
PowerPoint生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── カラー定義 ──
C_BROWN   = RGBColor(0x73, 0x57, 0x63)  # #735763
C_ORANGE  = RGBColor(0xF7, 0x93, 0x21)  # #F79321
C_BEIGE   = RGBColor(0xFC, 0xF0, 0xDE)  # #FCF0DE
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x4A, 0x30, 0x40)  # #4A3040
C_MUTED   = RGBColor(0x9B, 0x7E, 0x8E)  # #9B7E8E
C_BEIGE2  = RGBColor(0xF5, 0xE8, 0xD8)  # #F5E8D8
C_BEIGE3  = RGBColor(0xF5, 0xF0, 0xED)  # #F5F0ED
C_LIGHT   = RGBColor(0xF0, 0xEB, 0xF0)  # #F0EBF0

# ── スライドサイズ 16:9 ──
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # 完全空白


def add_shape(slide, x, y, w, h, fill_color=None, border_color=None, radius=None):
    """塗りつぶし矩形を追加"""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, v_anchor=None,
             line_spacing=None, font_name="Noto Sans JP"):
    """テキストボックスを追加"""
    from pptx.util import Pt
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    if v_anchor:
        tf.vertical_anchor = v_anchor
    # 既存段落を使用
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
        if line_spacing:
            from pptx.util import Pt as Pt2
            p.line_spacing = line_spacing
    return txBox


def add_label_badge(slide, text, x, y, w, h, bg_color, text_color, font_size=11):
    """バッジ（ラベル付き角丸風矩形）"""
    s = add_shape(slide, x, y, w, h, fill_color=bg_color)
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.name = "Noto Sans JP"
    run.font.color.rgb = text_color
    return s


def header_bar(slide, title_text, badge_text=None, badge_color=C_ORANGE):
    """共通ヘッダーバー"""
    # 背景
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.9), fill_color=C_BROWN)
    # タイトル
    add_text(slide, title_text,
             Inches(0.4), Inches(0.08), Inches(10), Inches(0.75),
             font_size=28, bold=True, color=C_WHITE)
    # バッジ
    if badge_text:
        add_label_badge(slide, badge_text,
                        Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.5),
                        badge_color, C_WHITE, font_size=12)


def card(slide, x, y, w, h, bg=C_BEIGE):
    """カードの背景矩形"""
    return add_shape(slide, x, y, w, h, fill_color=bg)


def effort_bar(slide, text, x, y, w):
    """工数バー"""
    s = add_shape(slide, x, y, w, Inches(0.42), fill_color=C_BROWN)
    tf = s.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = "Noto Sans JP"
    run.font.color.rgb = C_ORANGE


# ═══════════════════════════════════════════════
# Slide 1 : 表紙
# ═══════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_shape(s1, Inches(0), Inches(0), W, H, fill_color=C_BROWN)
# アクセントバー
add_shape(s1, Inches(1.0), Inches(2.2), Inches(0.7), Inches(0.06), fill_color=C_ORANGE)
# メインタイトル
add_text(s1, "福元鍼灸整骨院\nオンライン予約システム",
         Inches(1.0), Inches(2.3), Inches(11), Inches(2.2),
         font_size=44, bold=True, color=C_WHITE, line_spacing=Pt(54))
# サブタイトル
add_text(s1, "本番運用ロードマップ",
         Inches(1.0), Inches(4.55), Inches(9), Inches(0.6),
         font_size=24, color=C_BEIGE)
# メタ
add_text(s1, "Phase 1〜5　全工程・概算見積もり",
         Inches(1.0), Inches(5.1), Inches(9), Inches(0.5),
         font_size=16, color=C_ORANGE)


# ═══════════════════════════════════════════════
# Slide 2 : アジェンダ
# ═══════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_shape(s2, Inches(0), Inches(0), W, H, fill_color=C_BEIGE)
header_bar(s2, "アジェンダ")

col_w = Inches(6.2)
col_h = Inches(6.0)
cx1, cx2 = Inches(0.2), Inches(6.7)
cy = Inches(1.0)

# 左カラム背景
card(s2, cx1, cy, col_w, col_h, bg=C_WHITE)
# 右カラム背景
card(s2, cx2, cy, col_w, col_h, bg=C_BEIGE2)

agenda_left = [
    ("Phase 1　インフラ構築", [
        "🔥 Firebase プロジェクト構築",
        "🌐 ホスティング・ドメイン設定",
        "🗄️ データベース設計",
        "🔐 認証実装",
        "📧 メール送信 API 構築",
    ]),
    ("Phase 2　コード移行", [
        "💾 localStorage → Firestore 移行",
        "⚡ リアルタイム同期",
        "🔒 二重予約防止",
        "🛡️ 管理画面認証",
    ]),
]
agenda_right = [
    ("Phase 3　通知機能", [
        "📩 患者への確認メール",
        "🔔 管理者通知",
        "⏰ リマインダー（任意）",
    ]),
    ("Phase 4〜5　仕上げ・リリース", [
        "✏️ コンテンツ修正",
        "📅 祝日対応",
        "🔐 セキュリティルール",
        "✅ 動作確認・デプロイ",
    ]),
]

def render_agenda_col(slide, items, cx, cy):
    oy = cy + Inches(0.15)
    for title, bullets in items:
        add_text(slide, title, cx + Inches(0.2), oy, col_w - Inches(0.4), Inches(0.45),
                 font_size=17, bold=True, color=C_BROWN)
        oy += Inches(0.42)
        for b in bullets:
            add_text(slide, b, cx + Inches(0.3), oy, col_w - Inches(0.5), Inches(0.32),
                     font_size=14, color=C_DARK)
            oy += Inches(0.3)
        oy += Inches(0.18)

render_agenda_col(s2, agenda_left, cx1, cy)
render_agenda_col(s2, agenda_right, cx2, cy)


# ═══════════════════════════════════════════════
# Slide 3 : Phase 1 インフラ構築
# ═══════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_shape(s3, Inches(0), Inches(0), W, H, fill_color=C_WHITE)
header_bar(s3, "Phase 1　インフラ構築", "🔴 MUST")

ph1_cards = [
    ("🔥 Firebase\nプロジェクト構築", [
        "・プロジェクト作成",
        "・Firestore DB 作成",
        "・Firebase Hosting",
        "  （静的ファイルの配信サービス）",
        "・Firebase Auth",
        "  （認証サービス）設定",
    ], "工数：約 0.5日"),
    ("🌐 ドメイン・\nSSL 設定", [
        "・独自ドメイン取得",
        "  （例: yoyaku.fukumoto.jp）",
        "・SSL（暗号化通信）証明書",
        "  → Firebase が自動発行",
        "・DNS 設定",
        "  （ドメインの名前解決）",
    ], "工数：約 0.5日"),
    ("📧 メール送信\nAPI 構築", [
        "・Resend または SendGrid",
        "  （メール送信サービス）契約",
        "・Firebase Functions",
        "  （サーバーレス関数）で",
        "  送信ロジック実装",
        "・SPF/DKIM 認証設定",
        "  （迷惑メール対策）",
    ], "工数：約 1.5日"),
]

card_w = Inches(4.0)
card_h = Inches(6.1)
gaps = [Inches(0.2), Inches(4.5), Inches(8.8)]
card_y = Inches(1.0)

for i, (title, bullets, effort) in enumerate(ph1_cards):
    cx = gaps[i]
    card(s3, cx, card_y, card_w, card_h)
    add_text(s3, title, cx + Inches(0.2), card_y + Inches(0.2),
             card_w - Inches(0.4), Inches(0.85),
             font_size=18, bold=True, color=C_BROWN, line_spacing=Pt(24))
    by = card_y + Inches(1.1)
    for b in bullets:
        add_text(s3, b, cx + Inches(0.2), by, card_w - Inches(0.4), Inches(0.3),
                 font_size=13, color=C_DARK)
        by += Inches(0.28)
    effort_bar(s3, effort, cx + Inches(0.2), card_y + card_h - Inches(0.55),
               card_w - Inches(0.4))


# ═══════════════════════════════════════════════
# Slide 4 : Phase 2 コード移行（Before/After）
# ═══════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_shape(s4, Inches(0), Inches(0), W, H, fill_color=C_WHITE)
header_bar(s4, "Phase 2　コード移行", "🔴 MUST")

# Before カード
card(s4, Inches(0.3), Inches(1.05), Inches(5.5), Inches(6.0), bg=C_BEIGE)
add_label_badge(s4, "Before　現在の構成",
                Inches(0.5), Inches(1.2), Inches(2.6), Inches(0.38),
                C_BEIGE2, C_BROWN, font_size=11)
add_text(s4, "localStorage\n（ブラウザ内ストレージ）",
         Inches(0.5), Inches(1.7), Inches(5.0), Inches(0.9),
         font_size=22, bold=True, color=C_MUTED, line_spacing=Pt(28))
before_items = [
    "❌ 他のデバイスから見えない",
    "❌ ブラウザを消すとデータ消失",
    "❌ 二重予約が起きる",
    "❌ 管理画面に認証なし",
]
by = Inches(2.7)
for item in before_items:
    add_text(s4, item, Inches(0.5), by, Inches(5.0), Inches(0.35),
             font_size=15, color=C_BROWN)
    by += Inches(0.38)

# 矢印
add_text(s4, "→", Inches(6.1), Inches(3.5), Inches(0.8), Inches(0.7),
         font_size=40, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# After カード
card(s4, Inches(7.1), Inches(1.05), Inches(5.9), Inches(6.0), bg=C_BROWN)
add_label_badge(s4, "After　本番構成",
                Inches(7.3), Inches(1.2), Inches(2.3), Inches(0.38),
                RGBColor(0x5a, 0x40, 0x50), C_ORANGE, font_size=11)
add_text(s4, "Firestore\n（Firebase のクラウド DB）",
         Inches(7.3), Inches(1.7), Inches(5.4), Inches(0.9),
         font_size=22, bold=True, color=C_ORANGE, line_spacing=Pt(28))
after_items = [
    "✅ 全デバイスからリアルタイム同期",
    "✅ データは永続化・クラウド保存",
    "✅ トランザクション（排他制御）",
    "   で二重予約防止",
    "✅ Firebase Auth（認証サービス）",
    "   で管理画面保護",
]
ay = Inches(2.7)
for item in after_items:
    add_text(s4, item, Inches(7.3), ay, Inches(5.4), Inches(0.35),
             font_size=15, color=C_BEIGE)
    ay += Inches(0.35)
effort_bar(s4, "工数：約 3〜4日",
           Inches(7.3), Inches(6.2), Inches(5.4))
s_ef = s4.shapes[-1]
s_ef.fill.solid()
s_ef.fill.fore_color.rgb = C_ORANGE
for para in s_ef.text_frame.paragraphs:
    for run in para.runs:
        run.font.color.rgb = C_BROWN


# ═══════════════════════════════════════════════
# Slide 5 : Phase 3 通知機能
# ═══════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_shape(s5, Inches(0), Inches(0), W, H, fill_color=C_WHITE)
header_bar(s5, "Phase 3　通知機能", "🔴 MUST")

ph3_cards = [
    ("📩 患者への\n確認メール", [
        "予約確定と同時に自動送信",
        "",
        "・予約番号",
        "・日時・担当",
        "・キャンセル方法",
    ], "工数：約 1日", C_BEIGE, False),
    ("🔔 管理者への\n新着通知", [
        "新規予約が入ったら",
        "リアルタイムで通知",
        "",
        "・患者名・日時",
        "・症状・初再診区分",
    ], "工数：約 0.5日", C_BEIGE, False),
    ("⏰ 前日\nリマインダー", [
        "予約前日に自動送信",
        "（任意実装）",
        "",
        "・Firebase Functions",
        "  （定期実行機能）で",
        "  スケジュール処理",
    ], "工数：約 1日", C_LIGHT, True),
]

for i, (title, bullets, effort, bg, is_optional) in enumerate(ph3_cards):
    cx = Inches(0.2 + i * 4.35)
    card(s5, cx, Inches(1.0), Inches(4.1), Inches(6.1), bg=bg)
    add_text(s5, title, cx + Inches(0.2), Inches(1.15),
             Inches(3.7), Inches(0.85),
             font_size=19, bold=True, color=C_BROWN, line_spacing=Pt(26))
    by = Inches(2.1)
    for b in bullets:
        add_text(s5, b, cx + Inches(0.2), by, Inches(3.7), Inches(0.3),
                 font_size=14, color=C_DARK)
        by += Inches(0.3)
    if is_optional:
        add_label_badge(s5, "NICE TO HAVE",
                        cx + Inches(0.2), Inches(6.05), Inches(2.0), Inches(0.32),
                        C_MUTED, C_WHITE, font_size=10)
    effort_bar(s5, effort, cx + Inches(0.2), Inches(6.65),
               Inches(3.7))


# ═══════════════════════════════════════════════
# Slide 6 : Phase 4〜5 仕上げ・リリース
# ═══════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_shape(s6, Inches(0), Inches(0), W, H, fill_color=C_WHITE)
header_bar(s6, "Phase 4〜5　仕上げ・リリース", "🟡 SHOULD", badge_color=C_MUTED)

ph45_cards = [
    ("✏️ コンテンツ修正", [
        "・電話番号のプレースホルダー",
        "  （仮の値）を本番値に変更",
        "・PDF の日本語ラベル修正",
        "・日本語フォントの埋め込み",
        "・デモデータの無効化",
    ], "工数：約 0.5日", C_BEIGE),
    ("📅 祝日・法的対応", [
        "・祝日 API 連携",
        "  （国民の祝日の自動グレーアウト）",
        "・プライバシーポリシー",
        "  ページ作成",
        "・個人情報同意",
        "  チェックボックス追加",
    ], "工数：約 1日", C_BEIGE),
    ("🔐 セキュリティ・\nデプロイ", [
        "・Firestore セキュリティルール",
        "  （アクセス制御ルール）設定",
        "  患者：書き込みのみ",
        "  管理者：読み書き OK",
        "・UAT（受入テスト）実施",
        "・Firebase Hosting へデプロイ",
    ], "工数：約 1日", C_BROWN),
]

for i, (title, bullets, effort, bg) in enumerate(ph45_cards):
    cx = Inches(0.2 + i * 4.35)
    card(s6, cx, Inches(1.0), Inches(4.1), Inches(6.1), bg=bg)
    tc = C_BEIGE if bg == C_BROWN else C_BROWN
    add_text(s6, title, cx + Inches(0.2), Inches(1.15),
             Inches(3.7), Inches(0.85),
             font_size=19, bold=True, color=tc, line_spacing=Pt(26))
    bc = C_BEIGE if bg == C_BROWN else C_DARK
    by = Inches(2.1)
    for b in bullets:
        add_text(s6, b, cx + Inches(0.2), by, Inches(3.7), Inches(0.3),
                 font_size=14, color=bc)
        by += Inches(0.3)
    ebar = add_shape(s6, cx + Inches(0.2),
                     Inches(6.65), Inches(3.7), Inches(0.42),
                     fill_color=C_ORANGE if bg == C_BROWN else C_BROWN)
    tf = ebar.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = effort
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = "Noto Sans JP"
    run.font.color.rgb = C_BROWN if bg == C_BROWN else C_ORANGE


# ═══════════════════════════════════════════════
# Slide 7 : 概算見積もり
# ═══════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_shape(s7, Inches(0), Inches(0), W, H, fill_color=C_WHITE)
header_bar(s7, "概算見積もり")
add_text(s7, "月額ランニングコスト（Firebase 無料枠内で想定）",
         Inches(5.5), Inches(0.12), Inches(7.5), Inches(0.6),
         font_size=13, color=C_BEIGE)

# コストカード 4枚
cost_items = [
    ("Firebase 無料枠\n（Spark プラン）", "¥ 0 / 月", "Firestore・Auth・Hosting 含む", C_BEIGE),
    ("独自ドメイン", "¥ 1,500 / 年", "約 125円/月", C_BEIGE2),
    ("メール送信\n（Resend 無料枠）", "¥ 0 / 月", "3,000通/月まで無料", C_BEIGE),
    ("合計（初月〜）", "¥ 125 / 月〜", "規模拡大で\nFunctions 有料化あり", C_ORANGE),
]

cw = Inches(3.1)
ch = Inches(1.7)
cy_cost = Inches(1.0)
for i, (label, val, note, bg) in enumerate(cost_items):
    cx = Inches(0.2 + i * 3.27)
    card(s7, cx, cy_cost, cw, ch, bg=bg)
    tc_label = C_WHITE if bg == C_ORANGE else C_MUTED
    tc_val   = C_WHITE if bg == C_ORANGE else C_BROWN
    tc_note  = C_WHITE if bg == C_ORANGE else C_MUTED
    add_text(s7, label, cx + Inches(0.1), cy_cost + Inches(0.1),
             cw - Inches(0.2), Inches(0.5),
             font_size=12, color=tc_label, line_spacing=Pt(16))
    add_text(s7, val, cx + Inches(0.1), cy_cost + Inches(0.6),
             cw - Inches(0.2), Inches(0.55),
             font_size=22, bold=True, color=tc_val)
    add_text(s7, note, cx + Inches(0.1), cy_cost + Inches(1.2),
             cw - Inches(0.2), Inches(0.45),
             font_size=11, color=tc_note, line_spacing=Pt(14))

# 工数サマリー
card(s7, Inches(0.2), Inches(2.9), Inches(8.2), Inches(4.3), bg=C_BEIGE3)
add_text(s7, "工数サマリー（ボス1人 ＋ AI 協力）",
         Inches(0.4), Inches(3.05), Inches(7.8), Inches(0.45),
         font_size=16, bold=True, color=C_BROWN)

summary_rows = [
    ("Phase 1", "インフラ構築",                             "2.5日"),
    ("Phase 2", "コード移行（localStorage → Firestore）",   "3〜4日"),
    ("Phase 3", "通知機能（メール送信）",                    "1.5日"),
    ("Phase 4", "仕上げ（コンテンツ・祝日・法的対応）",       "1.5日"),
    ("Phase 5", "セキュリティ・UAT・デプロイ",               "1日"),
]
ry = Inches(3.6)
for phase, desc, days in summary_rows:
    add_text(s7, phase, Inches(0.4), ry, Inches(1.0), Inches(0.32),
             font_size=13, bold=True, color=C_BROWN)
    add_text(s7, desc, Inches(1.5), ry, Inches(5.5), Inches(0.32),
             font_size=13, color=C_DARK)
    add_text(s7, days, Inches(7.2), ry, Inches(1.0), Inches(0.32),
             font_size=13, bold=True, color=C_BROWN, align=PP_ALIGN.RIGHT)
    ry += Inches(0.38)

# 区切り線
add_shape(s7, Inches(0.4), ry, Inches(7.8), Inches(0.02), fill_color=C_MUTED)
add_text(s7, "合計", Inches(0.4), ry + Inches(0.1), Inches(1.0), Inches(0.4),
         font_size=14, bold=True, color=C_BROWN)
add_text(s7, "約 10〜12日", Inches(7.0), ry + Inches(0.1), Inches(1.4), Inches(0.4),
         font_size=14, bold=True, color=C_ORANGE, align=PP_ALIGN.RIGHT)

# 前提条件
card(s7, Inches(8.6), Inches(2.9), Inches(4.5), Inches(4.3), bg=C_BROWN)
add_text(s7, "💡 前提条件",
         Inches(8.8), Inches(3.05), Inches(4.1), Inches(0.45),
         font_size=16, bold=True, color=C_ORANGE)
prereqs = [
    "・Firebase の基本知識あり",
    "・コーディングは AI（ワタシ）と",
    "  ペアプログラミング",
    "・1日 ＝ 実作業 4〜6時間換算",
    "・テスト・レビュー時間含む",
]
py = Inches(3.6)
for pr in prereqs:
    add_text(s7, pr, Inches(8.8), py, Inches(4.1), Inches(0.32),
             font_size=13, color=C_BEIGE)
    py += Inches(0.34)


# ═══════════════════════════════════════════════
# Slide 8 : クロージング
# ═══════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
add_shape(s8, Inches(0), Inches(0), W, H, fill_color=C_BROWN)
# アクセントバー
add_shape(s8, Inches(1.0), Inches(2.0), Inches(0.6), Inches(0.06), fill_color=C_ORANGE)
# メインメッセージ
add_text(s8, "さあ、本番を作ろう。",
         Inches(1.0), Inches(2.1), Inches(11), Inches(1.5),
         font_size=54, bold=True, color=C_WHITE)
# サブ
add_text(s8, "Phase 1 から、一緒に進めましょう。",
         Inches(1.0), Inches(3.65), Inches(10), Inches(0.6),
         font_size=26, color=C_BEIGE)
# 区切り
add_shape(s8, Inches(1.0), Inches(4.35), Inches(1.5), Inches(0.04), fill_color=C_ORANGE)

# ステップバッジ
steps = ["インフラ構築", "コード移行", "通知機能", "仕上げ", "リリース"]
for i, step in enumerate(steps):
    sx = Inches(1.0 + i * 2.3)
    sy = Inches(4.55)
    # 番号円
    circle = add_shape(s8, sx, sy, Inches(0.55), Inches(0.55), fill_color=C_ORANGE)
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(i + 1)
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.name = "Noto Sans JP"
    run.font.color.rgb = C_WHITE
    # ラベル
    add_text(s8, step, sx - Inches(0.15), sy + Inches(0.6),
             Inches(0.85), Inches(0.35),
             font_size=12, color=C_BEIGE, align=PP_ALIGN.CENTER)

# フッター
add_text(s8, "総工数 約 10〜12日　|　月額コスト 約 ¥125〜",
         Inches(1.0), Inches(6.6), Inches(9), Inches(0.45),
         font_size=16, bold=True, color=C_ORANGE)


# ── 保存 ──
output_path = r"G:\マイドライブ\Workspace\apps\fukumoto-reservation\documents\福元鍼灸整骨院_本番運用ロードマップ.pptx"
prs.save(output_path)
print(f"保存完了: {output_path}")

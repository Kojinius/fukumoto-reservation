"""
Online Appointment System ユーザーマニュアル 管理編 PPTX 生成スクリプト
実行: PYTHONUTF8=1 python documents/make_admin_manual.py
     (fukumoto-reservation ディレクトリから実行)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

# ─── スライドサイズ ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ─── ブランドカラー（Bold Navy × Gold）──────────────────
C_NAVY    = RGBColor(0x1A, 0x2B, 0x4A)
C_GOLD    = RGBColor(0xD4, 0xAF, 0x37)
C_BASE    = RGBColor(0xF8, 0xF5, 0xEE)   # クリーム背景
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x2B, 0x4A)
C_MUTED   = RGBColor(0x6B, 0x7C, 0x9E)

FONT_PRIMARY = "Noto Sans JP"

# ─── パス設定 ─────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent.parent
DOC_DIR     = BASE_DIR / "documents"
SS_DIR      = DOC_DIR / "screenshots"
OUTPUT_DIR  = DOC_DIR / "manuals"
OUTPUT_FILE = OUTPUT_DIR / "OAS_ユーザーマニュアル_管理編.pptx"


# ─── 部品関数（make_user_manual.py から継承）──────────────

def add_shape(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, line_spacing=None, font_name=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    fn = font_name or FONT_PRIMARY
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = fn
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox


def header_bar(slide, title_text, badge_text=None, badge_color=None):
    if badge_color is None:
        badge_color = C_GOLD
    bar_h = Inches(0.65)
    add_shape(slide, 0, 0, W, bar_h, fill_color=C_NAVY)
    add_text(slide, title_text, Inches(0.3), Inches(0.08),
             Inches(9.5), bar_h, font_size=22, bold=True, color=C_WHITE)
    if badge_text:
        badge_w = Inches(2.4)
        badge_x = W - badge_w - Inches(0.2)
        add_shape(slide, badge_x, Inches(0.1), badge_w, Inches(0.45),
                  fill_color=badge_color)
        add_text(slide, badge_text, badge_x, Inches(0.1),
                 badge_w, Inches(0.45), font_size=13, bold=True,
                 color=C_NAVY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, bg=None):
    return add_shape(slide, x, y, w, h, fill_color=bg or C_BASE)


def add_image_fit(slide, img_path, x, y, max_w, max_h):
    if not img_path or not Path(img_path).exists():
        box = add_shape(slide, x, y, max_w, max_h,
                        fill_color=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(slide, "（スクリーンショット準備中）",
                 x, y + max_h // 2 - Inches(0.3),
                 max_w, Inches(0.6), font_size=13,
                 color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)
        return
    try:
        with PILImage.open(img_path) as img:
            iw, ih = img.size
    except Exception:
        add_shape(slide, x, y, max_w, max_h, fill_color=RGBColor(0xCC, 0xCC, 0xCC))
        return
    scale = min(max_w / iw, max_h / ih)
    draw_w = int(iw * scale)
    draw_h = int(ih * scale)
    cx = x + (max_w - draw_w) // 2
    cy = y + (max_h - draw_h) // 2
    slide.shapes.add_picture(str(img_path), cx, cy, draw_w, draw_h)


def screen_slide(prs, blank_layout, title, subtitle, desc_lines, img_path,
                 label_text="管理者向け", label_color=None, left_w=Inches(4.3)):
    if label_color is None:
        label_color = C_GOLD
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, title, badge_text=label_text, badge_color=label_color)

    bar_h = Inches(0.65)
    content_y = bar_h + Inches(0.1)
    content_h = H - content_y - Inches(0.15)

    card(slide, 0, content_y, left_w, content_h, bg=C_BASE)

    if subtitle:
        add_text(slide, subtitle, Inches(0.25), content_y + Inches(0.15),
                 left_w - Inches(0.3), Inches(0.5),
                 font_size=12, color=C_MUTED)

    text_y = content_y + (Inches(0.55) if subtitle else Inches(0.2))
    text_h = content_h - (Inches(0.55) if subtitle else Inches(0.2)) - Inches(0.1)
    txBox = slide.shapes.add_textbox(Inches(0.25), text_y, left_w - Inches(0.3), text_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in desc_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if line.startswith("## "):
            run = p.add_run()
            run.text = line[3:]
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = C_NAVY
        elif line.startswith("・"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK
        elif line.startswith("⚠"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xB4, 0x5A, 0x00)
        elif line == "":
            run = p.add_run()
            run.text = ""
            run.font.size = Pt(5)
        else:
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK

    right_x = left_w + Inches(0.25)
    right_w = W - right_x - Inches(0.2)
    right_y = content_y + Inches(0.2)
    right_h = content_h - Inches(0.4)
    add_image_fit(slide, img_path, right_x, right_y, right_w, right_h)
    return slide


def divider_slide(prs, blank_layout, section_title, items,
                  bg=None, accent=None):
    if bg is None:
        bg = C_NAVY
    if accent is None:
        accent = C_GOLD
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=bg)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=accent)
    add_text(slide, section_title, Inches(0.5), Inches(2.5),
             Inches(8), Inches(1.2),
             font_size=42, bold=True, color=C_WHITE)
    add_shape(slide, Inches(0.5), Inches(3.6), Inches(3.5), Inches(0.06),
              fill_color=accent)
    y_start = Inches(4.0)
    for i, item in enumerate(items):
        add_text(slide, f"▸  {item}", Inches(0.8),
                 y_start + Inches(0.45) * i, Inches(10), Inches(0.45),
                 font_size=16, color=RGBColor(0xFF, 0xE0, 0x80))
    return slide


def cover_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, Inches(4.5), fill_color=C_NAVY)
    add_shape(slide, 0, Inches(4.5), W, Inches(3.0), fill_color=C_BASE)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=C_GOLD)
    add_text(slide, "Online Appointment System",
             Inches(0.6), Inches(1.2), Inches(11), Inches(1.3),
             font_size=46, bold=True, color=C_WHITE)
    add_text(slide, "ユーザーマニュアル",
             Inches(0.6), Inches(2.5), Inches(8), Inches(0.9),
             font_size=28, color=C_GOLD)
    # 管理編バッジ
    add_shape(slide, Inches(0.6), Inches(3.15), Inches(2.0), Inches(0.42),
              fill_color=C_GOLD)
    add_text(slide, "管理編", Inches(0.6), Inches(3.15),
             Inches(2.0), Inches(0.42),
             font_size=16, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(0.6), Inches(3.7), Inches(4.5), Inches(0.05), fill_color=C_GOLD)
    add_text(slide, "2026年3月  |  管理者向けガイド",
             Inches(0.6), Inches(3.9), Inches(5), Inches(0.6),
             font_size=16, color=C_MUTED)
    add_text(slide, "OAS", Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.5),
             font_size=72, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    return slide


def toc_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, "目次", badge_text="Table of Contents", badge_color=C_MUTED)

    bar_h = Inches(0.65)
    col_y = bar_h + Inches(0.25)

    # 左列
    left_x = Inches(0.4)
    col_w = Inches(5.8)
    add_text(slide, "ログイン・予約管理",
             left_x, col_y, col_w, Inches(0.55),
             font_size=17, bold=True, color=C_NAVY)
    add_shape(slide, left_x, col_y + Inches(0.52),
              Inches(2.5), Inches(0.04), fill_color=C_GOLD)

    left_items = [
        "1.  ログイン手順",
        "2.  ダッシュボード概要（KPI・検索・フィルター）",
        "3.  予約詳細・ステータス変更",
        "4.  診察完了の記録",
        "5.  診察履歴ページ",
    ]
    for i, item in enumerate(left_items):
        add_text(slide, item, left_x + Inches(0.1),
                 col_y + Inches(0.75) + Inches(0.46) * i,
                 col_w - Inches(0.1), Inches(0.45), font_size=13, color=C_DARK)

    # 右列
    right_x = Inches(7.0)
    add_text(slide, "設定・その他",
             right_x, col_y, col_w, Inches(0.55),
             font_size=17, bold=True, color=C_NAVY)
    add_shape(slide, right_x, col_y + Inches(0.52),
              Inches(2.5), Inches(0.04), fill_color=C_GOLD)

    right_items = [
        "6.  基本情報の設定",
        "7.  営業日・時間設定",
        "8.  お知らせ・メンテナンス設定",
        "9.  利用規約・ポリシー設定",
        "10. アカウント管理",
        "11. パスワード変更",
        "12. モバイル対応",
    ]
    for i, item in enumerate(right_items):
        add_text(slide, item, right_x + Inches(0.1),
                 col_y + Inches(0.75) + Inches(0.44) * i,
                 col_w - Inches(0.1), Inches(0.43), font_size=13, color=C_DARK)
    return slide


def closing_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_NAVY)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=C_GOLD)
    add_text(slide, "管理者の皆様へ",
             Inches(0.6), Inches(2.0), Inches(12), Inches(0.8),
             font_size=24, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "ご利用ありがとうございます",
             Inches(0.6), Inches(2.8), Inches(12), Inches(1.4),
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(3.5), Inches(4.1), Inches(6.3), Inches(0.06), fill_color=C_GOLD)
    add_text(slide, "Online Appointment System — 管理編",
             Inches(0.6), Inches(4.35), Inches(12), Inches(0.7),
             font_size=20, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "https://oas.kojinius.jp/login",
             Inches(0.6), Inches(5.1), Inches(12), Inches(0.6),
             font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    add_text(slide, "ご不明点は担当者まで",
             Inches(0.6), Inches(5.85), Inches(12), Inches(0.5),
             font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
    return slide


# ─── メイン処理 ───────────────────────────────────────────

def build_pptx():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank_layout = prs.slide_layouts[6]

    # ── 1. カバー ──────────────────────────────────────────
    cover_slide(prs, blank_layout)

    # ── 2. 目次 ───────────────────────────────────────────
    toc_slide(prs, blank_layout)

    # ══════════════════════════════════════════════════════
    # セクション 1：ログイン・初回設定
    # ══════════════════════════════════════════════════════
    divider_slide(prs, blank_layout,
                  "1. ログイン・初回設定",
                  [
                      "管理画面へのアクセス方法",
                      "初回ログイン時のパスワード変更（強制）",
                      "パスワードリセット手順",
                  ])

    # ── 4. ログイン手順 ───────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="ログイン手順",
        subtitle="https://oas.kojinius.jp/login",
        desc_lines=[
            "## ログイン方法",
            "・メールアドレスとパスワードを入力",
            "・「ログイン」ボタンをクリック",
            "",
            "## 初回ログイン",
            "・初回は仮パスワードを使用",
            "・ログイン後、パスワード変更画面に",
            "　自動リダイレクト（変更必須）",
            "",
            "## パスワードを忘れた場合",
            "・「パスワードを忘れた方」リンクから",
            "　リセットメールを送信",
            "",
            "⚠ 管理者 URL は患者向け URL とは異なります",
        ],
        img_path=SS_DIR / "09_login_page.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ══════════════════════════════════════════════════════
    # セクション 2：予約管理
    # ══════════════════════════════════════════════════════
    divider_slide(prs, blank_layout,
                  "2. 予約管理",
                  [
                      "ダッシュボード — KPI・検索・フィルター",
                      "予約詳細・ステータス変更",
                      "診察完了の記録",
                  ])

    # ── 6. ダッシュボード概要 ─────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="予約ダッシュボード",
        subtitle="/admin — 予約の一括管理",
        desc_lines=[
            "## KPI カード（上部）",
            "・本日の予約 / 今月の予約",
            "・新規患者 / 未確認",
            "・クリックで一覧を絞り込み",
            "",
            "## 検索・フィルター",
            "・氏名・郵便番号・電話番号で検索",
            "・ステータス別タブ（全て / 未確認 /",
            "　確認済み / キャンセル / 診察完了）",
            "・予約受付日・診察予定日で絞り込み",
            "",
            "## CSV エクスポート",
            "・「CSV出力」ボタンでダウンロード",
        ],
        img_path=SS_DIR / "10_dashboard.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 7. 予約詳細・ステータス変更 ──────────────────────
    screen_slide(
        prs, blank_layout,
        title="予約詳細・ステータス変更",
        subtitle="一覧の行をクリックしてモーダルを開く",
        desc_lines=[
            "## 表示される情報",
            "・予約番号・診察予定日時",
            "・氏名・ふりがな・生年月日",
            "・住所・電話番号・メール",
            "・性別・初診/再診・保険証",
            "・症状・伝達事項・連絡方法",
            "",
            "## ステータス操作",
            "・「確認済みにする」— 未確認 → 確認済み",
            "・「診察完了」— 診察履歴に保存",
            "・「キャンセル」— 理由選択して実行",
            "",
            "⚠ キャンセル時は患者へ通知メールが",
            "　自動送信されます",
        ],
        img_path=SS_DIR / "11_reservation_detail_modal.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 8. 診察完了の記録 ────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="診察完了の記録",
        subtitle="ダッシュボードから診察完了をマーク",
        desc_lines=[
            "## 診察完了の手順",
            "・予約詳細モーダルを開く",
            "・「診察完了」ボタンをクリック",
            "・確認ダイアログで「完了にする」",
            "",
            "## 完了後の動作",
            "・診察履歴ページへ自動移行",
            "・予約ステータスが「診察完了」に変更",
            "・完了日時・操作者が記録される",
            "",
            "⚠ 診察完了は取り消しできません",
            "　誤操作の場合は履歴ページで訂正",
        ],
        img_path=SS_DIR / "10_dashboard.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ══════════════════════════════════════════════════════
    # セクション 3：診察履歴
    # ══════════════════════════════════════════════════════
    divider_slide(prs, blank_layout,
                  "3. 診察履歴",
                  [
                      "診察完了した患者の記録を閲覧",
                      "氏名・電話番号・診察日での検索",
                      "記録の訂正（個人情報保護法 第34条）",
                      "CSV エクスポート",
                  ])

    # ── 10. 診察履歴ページ ────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="診察履歴ページ",
        subtitle="/admin/history — 診察完了記録の管理",
        desc_lines=[
            "## 検索機能",
            "・氏名・郵便番号・電話番号",
            "・診察日 From / To で期間絞り込み",
            "・「検索クリア」で条件リセット",
            "",
            "## 一覧表示",
            "・診療日時・氏名・初/再診・電話番号",
            "・「詳細」ボタンで詳細モーダルを表示",
            "・「訂正」ボタンで記録の訂正が可能",
            "",
            "## CSV エクスポート",
            "・「CSV出力」ボタンでダウンロード",
            "",
            "⚠ 訂正は患者からの請求時のみ実施",
            "　（個人情報保護法 第34条）",
        ],
        img_path=SS_DIR / "12_history_page.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ══════════════════════════════════════════════════════
    # セクション 4：設定
    # ══════════════════════════════════════════════════════
    divider_slide(prs, blank_layout,
                  "4. 設定",
                  [
                      "基本情報（院名・電話番号・住所）",
                      "営業日設定（タイムライン・休日）",
                      "お知らせ・メンテナンスモード",
                      "利用規約・ポリシー（APPI 対応）",
                      "アカウント管理",
                  ])

    # ── 12. 基本情報 ──────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 — 基本情報",
        subtitle="/admin/settings → 「基本情報」タブ",
        desc_lines=[
            "## 設定項目",
            "・院名（必須）",
            "・電話番号（必須）",
            "・WebサイトURL",
            "・郵便番号（入力で住所自動補完）",
            "・都道府県・市区町村・番地（必須）",
            "・建物名・号室",
            "",
            "## 保存と反映",
            "・「保存」ボタンで即時反映",
            "・トップページの院情報に連動",
            "・利用規約・ポリシーの自動生成に使用",
            "",
            "⚠ 院名・電話番号は必須項目です",
        ],
        img_path=SS_DIR / "13_settings_basic_info.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 13. 営業日設定 ───────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 — 営業日設定",
        subtitle="/admin/settings → 「営業日設定」タブ",
        desc_lines=[
            "## 週間スケジュール",
            "・曜日ごとに営業/休診を切り替え",
            "・午前・午後それぞれ開始・終了時間を設定",
            "・タイムラインバーで視覚的に確認",
            "",
            "## 休日設定",
            "・カレンダーをクリックして休日に設定",
            "・「祝日を自動取得」で祝日を一括登録",
            "・「クリックで解除」で休日を解除",
            "",
            "## 予約締切設定",
            "・予約受付締切（診療開始の○分前まで）",
            "・キャンセル受付締切（同様）",
            "",
            "⚠ 予約が入っている日を休日にすると",
            "　確認ダイアログが表示されます",
        ],
        img_path=SS_DIR / "14_settings_business_days.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 14. お知らせ・メンテナンス設定 ───────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 — お知らせ・メンテナンス",
        subtitle="/admin/settings → 「お知らせ」タブ",
        desc_lines=[
            "## お知らせバナー",
            "・種別: お知らせ / 注意 / メンテナンス",
            "・メッセージ内容を自由入力",
            "・表示開始日時・終了日時を設定",
            "・「表示中」「非表示」で ON/OFF",
            "",
            "## 活用例",
            "・臨時休診・年末年始の受付時間変更",
            "・院内イベントの告知",
            "",
            "## メンテナンスモード",
            "・有効中は管理者以外アクセス不可",
            "・メンテナンス開始・終了日時を設定",
            "",
            "⚠ メンテナンス中は患者の予約不可",
            "　設定前に必ず確認してください",
        ],
        img_path=SS_DIR / "15_settings_announcement.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 15. 利用規約・ポリシー設定 ───────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 — 利用規約・ポリシー",
        subtitle="/admin/settings → 「利用規約」「ポリシー」タブ",
        desc_lines=[
            "## 利用規約タブ",
            "・管理者ログイン時の同意規約を設定",
            "・基本情報から自動生成可能",
            "・バージョン更新で全管理者が再同意",
            "",
            "## ポリシータブ",
            "・プライバシーポリシー（患者向け）",
            "・要配慮個人情報の同意文言",
            "・データ保存期間ポリシー（APPI 対応）",
            "・個人情報の利用目的（APPI 第17条）",
            "・リマインダーメール配信の同意設定",
            "・患者の権利行使案内（APPI 第28〜30条）",
            "",
            "⚠ 法的要件を満たす内容に設定してください",
            "　不明な場合は専門家に確認を",
        ],
        img_path=SS_DIR / "16_settings_terms.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 16. アカウント管理 ───────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 — アカウント管理",
        subtitle="/admin/settings → 「アカウント」タブ",
        desc_lines=[
            "## 新規ユーザー作成",
            "・メールアドレスとパスワードを入力",
            "・「作成」ボタンで管理者を追加",
            "・パスワード要件: 8〜128文字",
            "　英大文字・小文字・数字・記号を含む",
            "",
            "## ユーザー一覧",
            "・登録済み管理者の一覧表示",
            "・メールアドレスはマスク表示",
            "・作成日・操作ボタン表示",
            "・「削除」ボタンで管理者を削除",
            "",
            "⚠ 自分自身は削除できません",
            "⚠ 追加後、初回ログイン時に",
            "　パスワード変更が強制されます",
        ],
        img_path=SS_DIR / "18_settings_accounts.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ══════════════════════════════════════════════════════
    # セクション 5：その他
    # ══════════════════════════════════════════════════════
    divider_slide(prs, blank_layout,
                  "5. その他",
                  [
                      "パスワード変更",
                      "モバイル対応（スマートフォン表示）",
                  ])

    # ── 18. パスワード変更 ────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="パスワード変更",
        subtitle="/admin/change-password",
        desc_lines=[
            "## パスワード変更手順",
            "・現在のパスワードを入力",
            "・新しいパスワードを入力",
            "・確認用パスワードを再入力",
            "・「パスワードを変更」をクリック",
            "",
            "## パスワード要件",
            "・8文字以上",
            "・英大文字を含む",
            "・英小文字を含む",
            "・数字を含む",
            "・記号を含む",
            "",
            "⚠ 初回ログイン時は変更必須",
            "⚠ 変更完了まで他ページへ移動不可",
        ],
        img_path=SS_DIR / "19_change_password.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 19. モバイル対応 ──────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="モバイル対応",
        subtitle="スマートフォンからも管理可能",
        desc_lines=[
            "## スマートフォン対応",
            "・iOS・Android のブラウザで利用可能",
            "・アプリインストール不要",
            "・ナビゲーションメニューはハンバーガー式",
            "",
            "## 利用可能な機能",
            "・ダッシュボード（予約一覧・KPI）",
            "・予約詳細・ステータス変更",
            "・設定（全タブ対応）",
            "・診察履歴閲覧",
            "",
            "⚠ 操作性は PC 画面での利用を推奨",
            "　外出時の確認用途に最適です",
        ],
        img_path=SS_DIR / "m02_admin_mobile_nav.png",
        label_text="管理者向け",
        label_color=C_GOLD,
    )

    # ── 20. クロージング ──────────────────────────────────
    closing_slide(prs, blank_layout)

    prs.save(str(OUTPUT_FILE))
    print(f"✅ 生成完了: {OUTPUT_FILE}")
    print(f"   スライド枚数: {len(prs.slides)}")


if __name__ == "__main__":
    build_pptx()

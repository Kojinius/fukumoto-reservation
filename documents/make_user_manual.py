"""
Online Appointment System ユーザーマニュアル PPTX 生成スクリプト（v2 - React SPA 対応版）
実行: PYTHONUTF8=1 python documents/make_user_manual.py
     (fukumoto-reservation ディレクトリから実行)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

# ─── スライドサイズ ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ─── ブランドカラー（Bold Navy × Gold）──────────────────
C_NAVY    = RGBColor(0x1A, 0x2B, 0x4A)
C_GOLD    = RGBColor(0xD4, 0xAF, 0x37)
C_BASE    = RGBColor(0xF8, 0xF5, 0xEE)   # クリーム背景
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x2B, 0x4A)
C_MUTED   = RGBColor(0x6B, 0x7C, 0x9E)

FONT_PRIMARY = "Noto Sans JP"

# ─── パス設定 ─────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent.parent
DOC_DIR     = BASE_DIR / "documents"
SS_DIR      = DOC_DIR / "screenshots"
OUTPUT_DIR  = DOC_DIR / "manuals"
OUTPUT_FILE = OUTPUT_DIR / "OnlineAppointSystem_ユーザーマニュアル.pptx"


# ─── 部品関数 ─────────────────────────────────────────────

def add_shape(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, line_spacing=None, font_name=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    fn = font_name or FONT_PRIMARY
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = fn
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox


def header_bar(slide, title_text, badge_text=None, badge_color=None):
    if badge_color is None:
        badge_color = C_GOLD
    bar_h = Inches(0.65)
    add_shape(slide, 0, 0, W, bar_h, fill_color=C_NAVY)
    add_text(slide, title_text, Inches(0.3), Inches(0.08),
             Inches(9.5), bar_h, font_size=22, bold=True, color=C_WHITE)
    if badge_text:
        badge_w = Inches(2.4)
        badge_x = W - badge_w - Inches(0.2)
        add_shape(slide, badge_x, Inches(0.1), badge_w, Inches(0.45),
                  fill_color=badge_color)
        add_text(slide, badge_text, badge_x, Inches(0.1),
                 badge_w, Inches(0.45), font_size=13, bold=True,
                 color=C_NAVY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, bg=None):
    return add_shape(slide, x, y, w, h, fill_color=bg or C_BASE)


def add_image_fit(slide, img_path, x, y, max_w, max_h):
    if not img_path or not Path(img_path).exists():
        box = add_shape(slide, x, y, max_w, max_h,
                        fill_color=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(slide, "（スクリーンショット準備中）",
                 x, y + max_h // 2 - Inches(0.3),
                 max_w, Inches(0.6), font_size=13,
                 color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)
        return
    try:
        with PILImage.open(img_path) as img:
            iw, ih = img.size
    except Exception:
        add_shape(slide, x, y, max_w, max_h, fill_color=RGBColor(0xCC, 0xCC, 0xCC))
        return
    scale = min(max_w / iw, max_h / ih)
    draw_w = int(iw * scale)
    draw_h = int(ih * scale)
    cx = x + (max_w - draw_w) // 2
    cy = y + (max_h - draw_h) // 2
    slide.shapes.add_picture(str(img_path), cx, cy, draw_w, draw_h)


def screen_slide(prs, blank_layout, title, subtitle, desc_lines, img_path,
                 label_text="患者向け", label_color=None, left_w=Inches(4.3)):
    if label_color is None:
        label_color = C_GOLD
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, title, badge_text=label_text, badge_color=label_color)

    bar_h = Inches(0.65)
    content_y = bar_h + Inches(0.1)
    content_h = H - content_y - Inches(0.15)

    card(slide, 0, content_y, left_w, content_h, bg=C_BASE)

    if subtitle:
        add_text(slide, subtitle, Inches(0.25), content_y + Inches(0.15),
                 left_w - Inches(0.3), Inches(0.5),
                 font_size=12, color=C_MUTED)

    text_y = content_y + (Inches(0.55) if subtitle else Inches(0.2))
    text_h = content_h - (Inches(0.55) if subtitle else Inches(0.2)) - Inches(0.1)
    txBox = slide.shapes.add_textbox(Inches(0.25), text_y, left_w - Inches(0.3), text_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in desc_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if line.startswith("## "):
            run = p.add_run()
            run.text = line[3:]
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = C_NAVY
        elif line.startswith("・"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK
        elif line.startswith("⚠"):
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xB4, 0x5A, 0x00)
        elif line == "":
            run = p.add_run()
            run.text = ""
            run.font.size = Pt(5)
        else:
            run = p.add_run()
            run.text = line
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK

    right_x = left_w + Inches(0.25)
    right_w = W - right_x - Inches(0.2)
    right_y = content_y + Inches(0.2)
    right_h = content_h - Inches(0.4)
    add_image_fit(slide, img_path, right_x, right_y, right_w, right_h)
    return slide


def divider_slide(prs, blank_layout, section_title, items,
                  bg=None, accent=None):
    if bg is None:
        bg = C_NAVY
    if accent is None:
        accent = C_GOLD
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=bg)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=accent)
    add_text(slide, section_title, Inches(0.5), Inches(2.5),
             Inches(8), Inches(1.2),
             font_size=42, bold=True, color=C_WHITE)
    add_shape(slide, Inches(0.5), Inches(3.6), Inches(3.5), Inches(0.06),
              fill_color=accent)
    y_start = Inches(4.0)
    for i, item in enumerate(items):
        add_text(slide, f"▸  {item}", Inches(0.8),
                 y_start + Inches(0.45) * i, Inches(10), Inches(0.45),
                 font_size=16, color=RGBColor(0xFF, 0xE0, 0x80))
    return slide


def cover_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, Inches(4.5), fill_color=C_NAVY)
    add_shape(slide, 0, Inches(4.5), W, Inches(3.0), fill_color=C_BASE)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=C_GOLD)
    add_text(slide, "Online Appointment System",
             Inches(0.6), Inches(1.4), Inches(11), Inches(1.3),
             font_size=46, bold=True, color=C_WHITE)
    add_text(slide, "ユーザーマニュアル",
             Inches(0.6), Inches(2.8), Inches(8), Inches(0.9),
             font_size=28, color=C_GOLD)
    add_shape(slide, Inches(0.6), Inches(3.7), Inches(4.5), Inches(0.05), fill_color=C_GOLD)
    add_text(slide, "2026年3月  |  React SPA v3",
             Inches(0.6), Inches(3.9), Inches(4), Inches(0.6),
             font_size=16, color=C_MUTED)
    add_text(slide, "OAS", Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.5),
             font_size=72, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    return slide


def toc_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_WHITE)
    header_bar(slide, "目次", badge_text="Table of Contents", badge_color=C_MUTED)

    bar_h = Inches(0.65)
    col_y = bar_h + Inches(0.3)

    # 患者向け列
    left_x = Inches(0.4)
    col_w  = Inches(5.8)
    add_text(slide, "患者向け機能",
             left_x, col_y, col_w, Inches(0.55),
             font_size=17, bold=True, color=C_NAVY)
    add_shape(slide, left_x, col_y + Inches(0.52),
              Inches(2.5), Inches(0.04), fill_color=C_GOLD)

    patient_items = [
        "01  トップページ（カレンダー日時選択）",
        "02  時間枠の選択",
        "03  お客様情報入力（Step 2）",
        "04  確認画面（Step 3）",
        "05  予約完了",
        "06  予約キャンセル（/cancel）",
        "07  プライバシーポリシー",
        "08  モバイル対応",
    ]
    for i, item in enumerate(patient_items):
        add_text(slide, item, left_x + Inches(0.1),
                 col_y + Inches(0.75) + Inches(0.46) * i,
                 col_w - Inches(0.1), Inches(0.45), font_size=13, color=C_DARK)

    # 管理者向け列
    right_x = Inches(7.0)
    add_text(slide, "管理者向け機能",
             right_x, col_y, col_w, Inches(0.55),
             font_size=17, bold=True, color=C_NAVY)
    add_shape(slide, right_x, col_y + Inches(0.52),
              Inches(2.5), Inches(0.04), fill_color=C_GOLD)

    admin_items = [
        "09  管理者ログイン（/login）",
        "10  予約ダッシュボード（/admin）",
        "11  予約詳細モーダル",
        "12  設定 - 基本情報",
        "13  設定 - 営業時間",
        "14  設定 - 休日設定",
        "15  設定 - お知らせ",
        "16  設定 - アカウント・ユーザー管理",
        "17  パスワード変更（/admin/change-password）",
        "18  リマインダーメール設定",
    ]
    for i, item in enumerate(admin_items):
        add_text(slide, item, right_x + Inches(0.1),
                 col_y + Inches(0.75) + Inches(0.44) * i,
                 col_w - Inches(0.1), Inches(0.43), font_size=13, color=C_DARK)
    return slide


def closing_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, W, H, fill_color=C_NAVY)
    add_shape(slide, 0, 0, Inches(0.18), H, fill_color=C_GOLD)
    add_text(slide, "ご利用ありがとうございます",
             Inches(0.6), Inches(2.6), Inches(12), Inches(1.4),
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(3.5), Inches(4.0), Inches(6.3), Inches(0.06), fill_color=C_GOLD)
    add_text(slide, "Online Appointment System",
             Inches(0.6), Inches(4.3), Inches(12), Inches(0.7),
             font_size=20, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "https://oas.kojinius.jp",
             Inches(0.6), Inches(5.1), Inches(12), Inches(0.6),
             font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    return slide


# ─── メイン処理 ───────────────────────────────────────────

def build_pptx():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank_layout = prs.slide_layouts[6]

    # ── 1. カバー ──────────────────────────────────────────
    cover_slide(prs, blank_layout)

    # ── 2. 目次 ──────────────────────────────────────────
    toc_slide(prs, blank_layout)

    # ── 章区切り：患者向け ─────────────────────────────────
    divider_slide(prs, blank_layout,
                  "患者向け機能",
                  [
                      "カレンダーで日時を選択",
                      "お客様情報を入力（郵便番号自動補完）",
                      "確認 → 予約完了",
                      "キャンセル・プライバシーポリシー",
                      "スマートフォン対応",
                  ])

    # ── 01 トップページ ────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="トップページ（カレンダー）",
        subtitle="https://oas.kojinius.jp",
        desc_lines=[
            "## アクセス",
            "・ブラウザで URL を開くだけ",
            "・アプリインストール不要",
            "",
            "## カレンダーの見方",
            "・グレー = 予約不可（休業日・過去日）",
            "・通常色 = 予約可能な日",
            "・希望日をクリックで時間枠を表示",
        ],
        img_path=SS_DIR / "01_calendar_top.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 02 時間枠の選択 ───────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="時間枠の選択",
        subtitle="Step 1 — 日時選択",
        desc_lines=[
            "## 時間枠を選ぶ",
            "・希望日をクリック後に時間枠が表示",
            "・「予約済み」枠は選択不可",
            "・選択した枠はハイライト表示",
            "",
            "## 次のステップ",
            "・枠を選択後「次へ」ボタンが有効に",
            "・クリックで情報入力フォームへ",
        ],
        img_path=SS_DIR / "02_time_slots.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 03 情報入力フォーム ───────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="お客様情報入力（Step 2）",
        subtitle="必要事項をご入力ください",
        desc_lines=[
            "## 入力項目",
            "・お名前・ふりがな（必須）",
            "・生年月日（任意）",
            "・電話番号（必須）",
            "・メールアドレス（任意）",
            "・住所（郵便番号→自動補完）",
            "・性別・初診/再診・保険有無",
            "・症状・連絡方法（任意）",
            "",
            "## 郵便番号自動補完",
            "・7桁入力で住所を自動入力",
            "⚠ 半角数字で入力してください",
        ],
        img_path=SS_DIR / "03_form_step2.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 04 確認画面 ───────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="確認画面（Step 3）",
        subtitle="予約内容を最終確認",
        desc_lines=[
            "## 予約内容の確認",
            "・日時・お名前・連絡先を最終確認",
            "・「修正する」で前の画面に戻れます",
            "",
            "## 予約の確定",
            "・「この内容で予約する」で確定",
            "・サーバー側で空き確認後に登録",
            "",
            "⚠ 「予約済み」になった場合は",
            "　別の時間帯をお選びください",
        ],
        img_path=SS_DIR / "05_confirm_step3.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 05 予約完了 ───────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="予約完了（Step 4）",
        subtitle="予約が完了しました",
        desc_lines=[
            "## 完了後の流れ",
            "・予約 ID が画面に表示されます",
            "・メールアドレスを入力した場合は",
            "　確認メールが届きます",
            "",
            "⚠ 予約 ID はキャンセル時に必要です",
            "　スクリーンショットで保存推奨",
        ],
        img_path=None,
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 06 キャンセル ─────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="予約キャンセル（/cancel）",
        subtitle="キャンセルの手順",
        desc_lines=[
            "## キャンセル方法",
            "・予約 ID を入力",
            "・電話番号を入力（照合用）",
            "・「キャンセルする」をクリック",
            "",
            "## セキュリティ",
            "・電話番号が一致しないと",
            "　キャンセル不可（なりすまし防止）",
            "・キャンセル通知メールが届きます",
            "",
            "⚠ 前日までのキャンセルを推奨",
        ],
        img_path=SS_DIR / "06_cancel_page.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 07 プライバシーポリシー ───────────────────────────
    screen_slide(
        prs, blank_layout,
        title="プライバシーポリシー",
        subtitle="https://oas.kojinius.jp/privacy-policy",
        desc_lines=[
            "## 個人情報の取り扱い",
            "・入力情報は予約管理のみに使用",
            "・第三者への提供は行いません",
            "・お問い合わせは院内窓口まで",
        ],
        img_path=SS_DIR / "07_privacy_policy.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 08 モバイル対応 ──────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="モバイル対応",
        subtitle="スマートフォンからも利用可能",
        desc_lines=[
            "## スマートフォン対応",
            "・iOS・Android 対応",
            "・ブラウザでアクセスするだけ",
            "・アプリインストール不要",
            "",
            "## 操作性",
            "・タッチ操作に最適化",
            "・ボタン・カレンダー全対応",
        ],
        img_path=SS_DIR / "m01_calendar_mobile.png",
        label_text="患者向け",
        label_color=C_GOLD,
    )

    # ── 章区切り：管理者向け ──────────────────────────────
    divider_slide(prs, blank_layout,
                  "管理者向け機能",
                  [
                      "管理者ログイン（/login）",
                      "予約ダッシュボード・KPIカード・CSV/PDF",
                      "設定（基本情報・営業時間・休日・お知らせ・アカウント）",
                      "ユーザー管理・パスワード変更",
                      "リマインダーメール",
                  ])

    # ── 09 管理者ログイン ─────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="管理者ログイン（/login）",
        subtitle="管理画面へのアクセス",
        desc_lines=[
            "## ログイン方法",
            "・メールアドレスとパスワードを入力",
            "・「ログイン」ボタンをクリック",
            "",
            "## 初回ログイン",
            "・初回は仮パスワードを使用",
            "・ログイン後、パスワード変更画面に",
            "　自動リダイレクト（強制変更）",
            "",
            "## パスワードを忘れた場合",
            "・「パスワードを忘れた方」リンクから",
            "　リセットメールを送信",
        ],
        img_path=SS_DIR / "08_login_page.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 10 予約ダッシュボード ─────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="予約ダッシュボード（/admin）",
        subtitle="予約の一括管理",
        desc_lines=[
            "## KPI カード",
            "・本日 / 今月 / 新規 / 未確認",
            "・クリックで一覧を絞り込み",
            "",
            "## 予約一覧",
            "・ステータス・日付・キーワードで絞り込み",
            "・行クリックで詳細モーダル表示",
            "",
            "## エクスポート",
            "・CSV ダウンロード",
            "・PDF 印刷（pdf-lib 使用）",
        ],
        img_path=SS_DIR / "09_dashboard.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 11 予約詳細モーダル ──────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="予約詳細モーダル",
        subtitle="予約情報の確認と更新",
        desc_lines=[
            "## 詳細情報",
            "・予約者の全情報を確認",
            "・症状・伝達事項を確認",
            "",
            "## ステータス変更",
            "・「確認済み」に変更",
            "・「キャンセル」に変更",
            "・変更は即時 Firestore に反映",
        ],
        img_path=SS_DIR / "10_reservation_detail_modal.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 12 設定 - 基本情報 ────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 - 基本情報（/admin/settings）",
        subtitle="クリニック基本情報の管理",
        desc_lines=[
            "## 設定項目",
            "・院名（必須）",
            "・住所",
            "・電話番号（必須）",
            "・その他情報",
            "",
            "## 保存",
            "・「保存」ボタンで即時反映",
            "・トップページの院情報に連動",
        ],
        img_path=SS_DIR / "11_settings_clinic_info.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 13 設定 - 営業時間 ────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 - 営業時間",
        subtitle="予約受付時間の設定",
        desc_lines=[
            "## 設定方法",
            "・曜日ごとに ON/OFF 切り替え",
            "・開始・終了時間を入力",
            "",
            "## カレンダーへの反映",
            "・設定後、即座に予約フォームに",
            "　反映される",
            "・OFF 設定日は予約不可に",
        ],
        img_path=SS_DIR / "12_settings_business_hours.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 14 設定 - 休日 ────────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 - 休日設定",
        subtitle="休診日・祝日の登録",
        desc_lines=[
            "## 休日の登録",
            "・カレンダーから休日を選択",
            "・休日名を任意で入力可能",
            "",
            "## カレンダーへの反映",
            "・登録した日は予約不可に",
            "・グレー表示で視覚的に識別",
        ],
        img_path=SS_DIR / "13_settings_holidays.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 15 設定 - お知らせ ────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 - お知らせ",
        subtitle="患者向けバナーメッセージ",
        desc_lines=[
            "## お知らせ機能",
            "・トップページにバナー表示",
            "・テキスト内容の編集",
            "・ON/OFF 切り替え",
            "",
            "## 活用例",
            "・臨時休診のお知らせ",
            "・年末年始の受付時間変更",
            "・院内イベントの告知",
        ],
        img_path=SS_DIR / "14_settings_announcement.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 16 設定 - アカウント・ユーザー管理 ────────────────
    screen_slide(
        prs, blank_layout,
        title="設定 - アカウント・ユーザー管理",
        subtitle="管理者アカウントの管理",
        desc_lines=[
            "## ユーザー管理",
            "・管理者アカウントの一覧表示",
            "・メールアドレスはマスク表示",
            "・新規管理者の追加",
            "・不要アカウントの削除",
            "",
            "## アカウント設定",
            "・パスワード変更ページへのリンク",
            "",
            "⚠ 管理者追加後は初回ログイン時に",
            "　パスワード変更が強制されます",
        ],
        img_path=SS_DIR / "15_settings_account.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 17 パスワード変更 ─────────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="パスワード変更（/admin/change-password）",
        subtitle="管理者パスワードの変更",
        desc_lines=[
            "## パスワード変更手順",
            "・現在のパスワードを入力",
            "・新しいパスワードを入力",
            "・確認用パスワードを再入力",
            "・「変更する」をクリック",
            "",
            "## パスワード要件",
            "・8文字以上",
            "・大文字・小文字・数字・記号",
            "　各1文字以上必須",
            "",
            "⚠ 初回ログイン時は変更必須です",
        ],
        img_path=SS_DIR / "16_change_password.png",
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── 18 リマインダーメール ─────────────────────────────
    screen_slide(
        prs, blank_layout,
        title="通知メール・リマインダー",
        subtitle="自動メール配信の仕組み",
        desc_lines=[
            "## 予約確認メール（自動）",
            "・予約完了時に患者へ自動送信",
            "・院内情報の連絡先が差出人",
            "",
            "## キャンセル通知メール（自動）",
            "・キャンセル完了時に自動送信",
            "",
            "## リマインダーメール（毎日自動）",
            "・翌日の予約者へ前日夜に送信",
            "・Cloud Functions スケジュール実行",
            "・(sendDailyReminders)",
            "",
            "⚠ メール送信は Resend API を使用",
        ],
        img_path=None,
        label_text="管理者向け",
        label_color=C_NAVY,
    )

    # ── クロージング ───────────────────────────────────────
    closing_slide(prs, blank_layout)

    prs.save(str(OUTPUT_FILE))
    print(f"✅ 生成完了: {OUTPUT_FILE}")
    print(f"   スライド枚数: {len(prs.slides)}")


if __name__ == "__main__":
    build_pptx()

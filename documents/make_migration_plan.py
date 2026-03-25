# -*- coding: utf-8 -*-
"""
OAS 本番環境移行計画書 PPTX 生成スクリプト
実行: PYTHONUTF8=1 python documents/make_migration_plan.py
     （fukumoto-reservation ディレクトリから）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# ─────────────────────────────────────────────
# カラー定義（OAS ブランドカラー Bold Navy × Gold）
# ─────────────────────────────────────────────
C_NAVY    = RGBColor(0x1A, 0x2B, 0x4A)
C_GOLD    = RGBColor(0xD4, 0xAF, 0x37)
C_BASE    = RGBColor(0xF8, 0xF5, 0xEE)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED   = RGBColor(0x6B, 0x7C, 0x9E)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_GREEN   = RGBColor(0x16, 0xA3, 0x4A)
C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
C_DARK    = RGBColor(0x1A, 0x2B, 0x4A)

BASE_DIR    = Path(__file__).parent.parent
OUTPUT_PATH = BASE_DIR / "documents" / "specs" / "OAS_本番環境移行計画書.pptx"

W = Inches(13.33)
H = Inches(7.5)


# ─────────────────────────────────────────────
# ヘルパー関数
# ─────────────────────────────────────────────

def add_shape(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, line_spacing=None, font_name="Noto Sans JP"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def header_bar(slide, title_text, badge_text=None, badge_color=None):
    if badge_color is None:
        badge_color = C_GOLD
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.9), fill_color=C_NAVY)
    add_text(slide, title_text, Inches(0.4), Inches(0.1), Inches(10), Inches(0.75),
             font_size=26, bold=True, color=C_WHITE)
    if badge_text:
        s = add_shape(slide, Inches(11.3), Inches(0.2), Inches(1.8), Inches(0.5),
                      fill_color=badge_color)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge_text
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.name = "Noto Sans JP"
        run.font.color.rgb = C_NAVY


def card(slide, x, y, w, h, bg=None):
    if bg is None:
        bg = C_BASE
    return add_shape(slide, x, y, w, h, fill_color=bg)


def divider_slide(prs, blank_layout, section_no, section_title, items=None):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.2), W, Inches(1.2), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(2.0), Inches(12.33), Inches(0.06), fill_color=C_GOLD)
    add_text(slide, section_no, Inches(0.5), Inches(0.8), Inches(12), Inches(0.6),
             font_size=20, bold=False, color=C_GOLD)
    add_text(slide, section_title, Inches(0.5), Inches(1.3), Inches(12), Inches(0.9),
             font_size=40, bold=True, color=C_WHITE)
    if items:
        y = Inches(2.4)
        for item in items:
            add_text(slide, f"  {item}", Inches(0.8), y, Inches(11.5), Inches(0.45),
                     font_size=16, color=C_BASE)
            y += Inches(0.45)


def checklist_rows(slide, items, x, y, w, row_h=0.42, done_color=None, todo_color=None):
    """チェックリスト行を描画（アイコン + テキスト）"""
    if done_color is None:
        done_color = C_GREEN
    if todo_color is None:
        todo_color = C_NAVY
    for label, status, note in items:
        color = done_color if status == "done" else (C_AMBER if status == "warn" else todo_color)
        icon = "✅" if status == "done" else ("⚠️" if status == "warn" else "□")
        row_text = f"{icon}  {label}"
        add_text(slide, row_text, x, y, Inches(w - 0.1), Inches(row_h),
                 font_size=14, color=color)
        if note:
            add_text(slide, f"     → {note}", x, y + Inches(0.22), Inches(w - 0.1), Inches(0.28),
                     font_size=10, color=C_MUTED)
        y += Inches(row_h + (0.2 if note else 0))
    return y


def phase_badge(slide, x, y, label, color=None):
    if color is None:
        color = C_NAVY
    s = add_shape(slide, x, y, Inches(2.2), Inches(0.38), fill_color=color)
    tf = s.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.name = "Noto Sans JP"
    run.font.color.rgb = C_WHITE


# ─────────────────────────────────────────────
# スライド生成関数
# ─────────────────────────────────────────────

def slide_cover(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), Inches(5.5), W, Inches(2.0), fill_color=C_GOLD)
    # アクセントライン
    add_shape(slide, Inches(0.5), Inches(1.8), Inches(0.08), Inches(3.0), fill_color=C_GOLD)
    add_text(slide, "オンライン予約システム",
             Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.7),
             font_size=22, color=C_MUTED)
    add_text(slide, "本番環境移行計画書",
             Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.1),
             font_size=48, bold=True, color=C_WHITE)
    add_text(slide, "Production Environment Migration Plan  v1.0",
             Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.5),
             font_size=16, color=C_GOLD)
    add_text(slide, "対象: 開発者・納品担当者",
             Inches(0.8), Inches(5.65), Inches(6.0), Inches(0.5),
             font_size=16, bold=True, color=C_NAVY)
    add_text(slide, "2026年3月",
             Inches(10.0), Inches(5.65), Inches(3.0), Inches(0.5),
             font_size=16, color=C_NAVY, align=PP_ALIGN.RIGHT)


def slide_toc(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "目次", "移行計画書")
    sections = [
        ("1.", "現状環境の整理",         "ボス個人アカウントで運用中のサービス一覧"),
        ("2.", "必要な契約・アカウント", "クライアント側で準備すべきサービス"),
        ("3.", "予算見積もり",           "従量課金サービスのコスト試算"),
        ("4.", "推奨環境構築手順",       "Firebase プロジェクト初期設定の推奨例"),
        ("5.", "移行チェックリスト",     "Phase 1〜4 の作業手順"),
        ("6.", "リスクと対策",           "移行時のリスク一覧・回避策・ロールバック"),
        ("7.", "データ移行方針",         "Firestore / Auth / Secrets の移行手順"),
    ]
    col_w = Inches(6.0)
    for i, (no, title, desc) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.1) + row * Inches(1.4)
        card(slide, x, y, Inches(6.1), Inches(1.25))
        add_text(slide, no, x + Inches(0.15), y + Inches(0.1), Inches(0.5), Inches(0.5),
                 font_size=22, bold=True, color=C_GOLD)
        add_text(slide, title, x + Inches(0.6), y + Inches(0.1), Inches(5.2), Inches(0.45),
                 font_size=17, bold=True, color=C_DARK)
        add_text(slide, desc, x + Inches(0.6), y + Inches(0.55), Inches(5.2), Inches(0.55),
                 font_size=12, color=C_MUTED)


def slide_current_env(prs, blank_layout):
    """1. 現状環境の整理"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "現状環境の整理", "Section 1")

    add_text(slide,
             "現在、すべてのサービスは開発者（個人）のアカウントで運用されています。\n"
             "本番リリースに際して、クライアント（医院）名義のアカウントへ移管が必要です。",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             font_size=14, color=C_MUTED)

    services = [
        ("ドメイン / DNS",      "kojinius.jp（開発者個人ドメイン）",    "Cloudflare にて管理。\noas.kojinius.jp が現在の公開URL。"),
        ("Firebase / GCP",      "project-3040e21e-879f-4c66-a7d",       "開発者の Google アカウントに紐づく Firebase プロジェクト。\nAuth / Firestore / Functions / Hosting すべて同一プロジェクト。"),
        ("Resend（メール配信）", "開発者の Resend アカウント",           "予約確認メール・リマインダーメール・キャンセル通知に使用。\nAPI キーは Cloud Functions の Secret Manager に格納。"),
    ]

    y = Inches(1.75)
    colors = [C_NAVY, C_GOLD, RGBColor(0x2D, 0x6A, 0x4F)]
    for i, (svc, current, detail) in enumerate(services):
        cx = Inches(0.5)
        card(slide, cx, y, Inches(12.3), Inches(1.5))
        add_shape(slide, cx, y, Inches(0.12), Inches(1.5), fill_color=colors[i])
        add_text(slide, svc, cx + Inches(0.25), y + Inches(0.1), Inches(3.5), Inches(0.45),
                 font_size=15, bold=True, color=C_DARK)
        add_text(slide, current, cx + Inches(3.8), y + Inches(0.1), Inches(8.4), Inches(0.4),
                 font_size=13, color=C_MUTED)
        add_text(slide, detail, cx + Inches(0.25), y + Inches(0.55), Inches(11.8), Inches(0.85),
                 font_size=12, color=C_DARK)
        y += Inches(1.65)


def slide_required_contracts(prs, blank_layout):
    """2. 必要な契約・アカウント"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "必要な契約・アカウント", "Section 2")

    add_text(slide, "クライアント（医院）側で以下を準備してください。",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
             font_size=14, color=C_MUTED)

    items = [
        ("① 独自ドメイン",
         "例: fukumoto-acupuncture.jp などクリニック名のドメイン",
         "お名前.com / Cloudflare Registrar で取得。\n.jp: 約 ¥1,500〜2,000/年　　.com: 約 ¥1,000〜1,500/年",
         C_NAVY),
        ("② Google アカウント（医院専用）",
         "例: admin@gmail.com（または Google Workspace ビジネスメール）",
         "このアカウントで Firebase プロジェクトを作成。個人の Google アカウントと混在しないよう医院専用を推奨。",
         C_GOLD),
        ("③ Firebase プロジェクト（Blaze プラン）",
         "Google アカウントにクレジットカードを登録し Blaze（従量課金）プランを有効化",
         "Cloud Functions の使用には Blaze プランが必須。無料枠内で収まる見込みのため実質課金なし（後述）。",
         RGBColor(0x2D, 0x6A, 0x4F)),
        ("④ Resend アカウント",
         "https://resend.com で無料登録 → 独自ドメインを認証",
         "予約確認・リマインダー・キャンセル通知メールの送信に使用。Free プランで月 3,000 通まで無料。",
         RGBColor(0x7C, 0x3A, 0xED)),
    ]

    y = Inches(1.5)
    for title, sub, detail, color in items:
        card(slide, Inches(0.5), y, Inches(12.3), Inches(1.2))
        add_shape(slide, Inches(0.5), y, Inches(0.12), Inches(1.2), fill_color=color)
        add_text(slide, title, Inches(0.75), y + Inches(0.08), Inches(4.5), Inches(0.38),
                 font_size=15, bold=True, color=C_DARK)
        add_text(slide, sub, Inches(5.3), y + Inches(0.08), Inches(7.3), Inches(0.38),
                 font_size=12, color=C_MUTED)
        add_text(slide, detail, Inches(0.75), y + Inches(0.5), Inches(11.8), Inches(0.6),
                 font_size=12, color=C_DARK)
        y += Inches(1.32)


def slide_budget(prs, blank_layout):
    """3. 予算見積もり"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "予算見積もり", "Section 3")

    add_text(slide, "想定: 月50〜200予約の小規模クリニック　　（Firebase Blaze プランの無料枠基準）",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.38),
             font_size=13, color=C_MUTED)

    # 左列: サービス別費用
    card(slide, Inches(0.5), Inches(1.5), Inches(7.2), Inches(5.5))
    add_text(slide, "サービス別 月額費用",
             Inches(0.7), Inches(1.6), Inches(6.8), Inches(0.45),
             font_size=16, bold=True, color=C_DARK)

    cost_items = [
        ("Firebase Hosting",    "無料",          "10GB / 360MB/日（無料枠内）"),
        ("Firestore",           "無料",          "読取 50,000回/日・書込 20,000回/日（無料枠内）"),
        ("Cloud Functions",     "無料",          "呼出 200万回/月・400,000 GB-秒（無料枠内）"),
        ("Firebase Auth",       "無料",          "10,000 MAU まで無料"),
        ("Resend（Free）",       "無料",          "月 3,000 通 / 100 通/日"),
        ("ドメイン",             "約 ¥125/月",    "年契約 ¥1,500/年 を月割換算"),
        ("Secret Manager",      "ほぼ無料",      "アクセス 10,000回/月まで無料"),
    ]

    y = Inches(2.15)
    for svc, cost, note in cost_items:
        cost_color = C_GREEN if cost == "無料" else (C_AMBER if "¥" in cost else C_MUTED)
        add_text(slide, svc, Inches(0.75), y, Inches(3.0), Inches(0.34), font_size=13, color=C_DARK)
        add_text(slide, cost, Inches(3.8), y, Inches(1.2), Inches(0.34),
                 font_size=13, bold=True, color=cost_color, align=PP_ALIGN.CENTER)
        add_text(slide, note, Inches(5.0), y, Inches(2.5), Inches(0.34), font_size=10, color=C_MUTED)
        y += Inches(0.56)

    add_shape(slide, Inches(0.7), y + Inches(0.1), Inches(6.8), Inches(0.04), fill_color=C_GOLD)
    add_text(slide, "合計（月額）",
             Inches(0.75), y + Inches(0.2), Inches(3.0), Inches(0.5),
             font_size=16, bold=True, color=C_DARK)
    add_text(slide, "¥0〜¥500 程度",
             Inches(3.8), y + Inches(0.2), Inches(3.7), Inches(0.5),
             font_size=20, bold=True, color=C_NAVY)

    # 右列: スケール別試算
    card(slide, Inches(7.9), Inches(1.5), Inches(5.0), Inches(5.5))
    add_text(slide, "スケール別 月額試算",
             Inches(8.1), Inches(1.6), Inches(4.6), Inches(0.45),
             font_size=16, bold=True, color=C_DARK)

    scales = [
        ("小規模",  "月 50 予約以下",    "¥0〜¥200",   C_GREEN,  "ドメイン代のみ"),
        ("中規模",  "月 100〜300 予約", "¥0〜¥500",   C_AMBER,  "ほぼ無料枠内"),
        ("大規模",  "月 1,000 予約以上", "¥500〜2,000", C_RED,    "Functions 超過分が発生"),
    ]
    y2 = Inches(2.15)
    for size, volume, cost, color, note in scales:
        s = add_shape(slide, Inches(8.1), y2, Inches(4.6), Inches(1.45))
        add_text(slide, size, Inches(8.2), y2 + Inches(0.1), Inches(1.2), Inches(0.45),
                 font_size=18, bold=True, color=color)
        add_text(slide, volume, Inches(9.5), y2 + Inches(0.1), Inches(3.0), Inches(0.45),
                 font_size=12, color=C_MUTED)
        add_text(slide, cost, Inches(8.2), y2 + Inches(0.55), Inches(4.2), Inches(0.5),
                 font_size=22, bold=True, color=C_DARK)
        add_text(slide, f"※ {note}", Inches(8.2), y2 + Inches(1.05), Inches(4.2), Inches(0.3),
                 font_size=10, color=C_MUTED)
        y2 += Inches(1.6)


def slide_recommended_setup(prs, blank_layout):
    """4. 推奨環境構築手順"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "推奨環境構築手順", "Section 4")

    add_text(slide, "以下の手順で Firebase プロジェクトを新規作成し、本番環境を構築してください。",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.38),
             font_size=13, color=C_MUTED)

    steps = [
        ("STEP 1",
         "Google Cloud プロジェクト作成",
         "① https://console.firebase.google.com/ にアクセス\n"
         "② 「プロジェクトを追加」→ プロジェクト名を入力（例: fukumoto-oas）\n"
         "③ Google アナリティクスは任意（不要なら無効化 OK）",
         C_NAVY),
        ("STEP 2",
         "Blaze プランへのアップグレード",
         "① Firebase Console → 左下「プランをアップグレード」→ Blaze を選択\n"
         "② クレジットカードを登録（Cloud Functions 使用に必須）\n"
         "③ 予算アラートを設定（例: ¥1,000/月 で通知）",
         C_GOLD),
        ("STEP 3",
         "必要サービスの有効化",
         "Authentication（メール/パスワード認証を有効化）\n"
         "Firestore Database（本番モードで作成 → rules は移行後に適用）\n"
         "Functions（Node.js 24 / 2nd Gen）　　Hosting（2 サイト: portfolio / oas）",
         RGBColor(0x2D, 0x6A, 0x4F)),
        ("STEP 4",
         "Secret Manager の設定",
         "① Google Cloud Console → Secret Manager → 「シークレットを作成」\n"
         "② 名前: resendApiKey　　値: Resend の API キーを貼り付け\n"
         "③ Cloud Functions のサービスアカウントにアクセス権を付与",
         RGBColor(0x7C, 0x3A, 0xED)),
    ]

    y = Inches(1.5)
    for i, (step, title, detail, color) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        yy = y + row * Inches(2.3)
        card(slide, x, yy, Inches(6.2), Inches(2.15))
        add_shape(slide, x, yy, Inches(0.12), Inches(2.15), fill_color=color)
        # ステップバッジ
        s = add_shape(slide, x + Inches(0.25), yy + Inches(0.12), Inches(1.0), Inches(0.32),
                      fill_color=color)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.name = "Noto Sans JP"
        run.font.color.rgb = C_WHITE
        add_text(slide, title, x + Inches(1.35), yy + Inches(0.1), Inches(4.7), Inches(0.45),
                 font_size=15, bold=True, color=C_DARK)
        add_text(slide, detail, x + Inches(0.25), yy + Inches(0.6), Inches(5.8), Inches(1.45),
                 font_size=11, color=C_DARK)


def slide_checklist_phase1(prs, blank_layout):
    """5. 移行チェックリスト — Phase 1: 契約・準備"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "移行チェックリスト — Phase 1: 契約・準備", "Section 5")

    phase_badge(slide, Inches(0.5), Inches(1.0), "Phase 1", C_NAVY)
    add_text(slide, "本番移行前にクライアント側で準備が必要な契約・設定事項",
             Inches(2.9), Inches(1.05), Inches(9.5), Inches(0.35),
             font_size=13, color=C_MUTED)

    items_left = [
        ("独自ドメインを取得する",                          "todo", ""),
        ("医院専用の Google アカウントを作成する",           "todo", ""),
        ("Firebase プロジェクトを作成する（新規）",          "todo", "project ID をメモしておく"),
        ("Blaze プランにアップグレードする",                 "todo", "クレジットカード登録が必要"),
        ("予算アラートを設定する（推奨: ¥1,000/月）",       "todo", ""),
    ]
    items_right = [
        ("Resend アカウントを作成する",                     "todo", "https://resend.com"),
        ("Resend で独自ドメインを認証する（SPF/DKIM）",     "todo", "DNS レコード追加が必要"),
        ("Resend API キーを発行する",                       "todo", ""),
        ("Secret Manager に resendApiKey を登録する",        "todo", ""),
        ("Firebase Hosting にカスタムドメインを設定する",   "todo", ""),
    ]

    checklist_rows(slide, items_left,  Inches(0.5), Inches(1.55), 6.2)
    checklist_rows(slide, items_right, Inches(6.9), Inches(1.55), 6.2)


def slide_checklist_phase2(prs, blank_layout):
    """5. 移行チェックリスト — Phase 2: コード・設定"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "移行チェックリスト — Phase 2: コード・設定移行", "Section 5")

    phase_badge(slide, Inches(0.5), Inches(1.0), "Phase 2", C_AMBER)
    add_text(slide, "ソースコードの接続先を新プロジェクトに切り替え、Firebase 設定を適用する",
             Inches(2.9), Inches(1.05), Inches(9.5), Inches(0.35),
             font_size=13, color=C_MUTED)

    items_left = [
        ("firebase.json の project ID を更新する",            "todo", "新プロジェクト ID に変更"),
        (".firebaserc の default プロジェクトを更新する",      "todo", ""),
        ("oas-spa/src/firebase.ts の firebaseConfig を更新",  "todo", "Console > プロジェクトの設定 > アプリ"),
        ("Firebase Authentication を有効化する",              "todo", "メール/パスワード認証"),
        ("Firestore をセキュリティルールごとデプロイする",     "todo", "firebase deploy --only firestore"),
    ]
    items_right = [
        ("Firebase Hosting をデプロイする（oas）",           "todo", "firebase deploy --only hosting:oas"),
        ("Cloud Functions をデプロイする",                   "todo", "firebase deploy --only functions"),
        ("Firestore インデックスをデプロイする",              "todo", "firebase deploy --only firestore:indexes"),
        ("カスタムドメインの DNS 伝播を確認する",             "todo", "nslookup / dig で確認（最大48時間）"),
        ("HTTPS アクセスを確認する（SSL証明書）",             "todo", "Firebase Hosting が自動発行"),
    ]

    checklist_rows(slide, items_left,  Inches(0.5), Inches(1.55), 6.2)
    checklist_rows(slide, items_right, Inches(6.9), Inches(1.55), 6.2)


def slide_checklist_phase3(prs, blank_layout):
    """5. 移行チェックリスト — Phase 3: データ移行"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "移行チェックリスト — Phase 3: データ移行", "Section 5")

    phase_badge(slide, Inches(0.5), Inches(1.0), "Phase 3", RGBColor(0x2D, 0x6A, 0x4F))
    add_text(slide, "既存の Firestore / Auth データを新プロジェクトへ移管する",
             Inches(2.9), Inches(1.05), Inches(9.5), Inches(0.35),
             font_size=13, color=C_MUTED)

    # Firestore 移行手順
    card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    add_text(slide, "Firestore データ移行",
             Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.4),
             font_size=15, bold=True, color=C_DARK)
    steps_fs = [
        "① 旧プロジェクトでメンテナンスモードを有効化",
        "② GCS バケットを作成（例: gs://fukumoto-backup）",
        "③ エクスポート実行:\n   gcloud firestore export gs://[bucket]",
        "④ 新プロジェクトにインポート:\n   gcloud firestore import gs://[bucket]/[export-prefix]",
        "⑤ データ整合性を確認（件数・内容チェック）",
    ]
    y = Inches(2.1)
    for s in steps_fs:
        add_text(slide, s, Inches(0.7), y, Inches(5.4), Inches(0.65),
                 font_size=12, color=C_DARK)
        y += Inches(0.72)

    # Auth 移行手順
    card(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(2.7))
    add_text(slide, "Firebase Auth ユーザー移行",
             Inches(6.7), Inches(1.6), Inches(5.9), Inches(0.4),
             font_size=15, bold=True, color=C_DARK)
    steps_auth = [
        "① エクスポート: firebase auth:export users.json",
        "② 新プロジェクトでインポート:",
        "   firebase auth:import users.json",
        "③ UID が同一のため Firestore との紐付け維持",
        "④ 移行後にパスワード変更を促す（推奨）",
    ]
    y2 = Inches(2.1)
    for s in steps_auth:
        add_text(slide, s, Inches(6.7), y2, Inches(5.9), Inches(0.38),
                 font_size=12, color=C_DARK)
        y2 += Inches(0.42)

    # Secrets 移行
    card(slide, Inches(6.5), Inches(4.35), Inches(6.3), Inches(2.6))
    add_text(slide, "Secrets / 設定値の移行",
             Inches(6.7), Inches(4.45), Inches(5.9), Inches(0.4),
             font_size=15, bold=True, color=C_DARK)
    steps_sec = [
        "① 旧プロジェクトの Secret Manager から resendApiKey の値を取得",
        "② 新プロジェクトの Secret Manager に同名で登録",
        "③ Cloud Functions デプロイ時に新 Secret が参照されること確認",
        "④ 旧プロジェクトの Secret は移行確認後に削除",
    ]
    y3 = Inches(4.98)
    for s in steps_sec:
        add_text(slide, s, Inches(6.7), y3, Inches(5.9), Inches(0.38),
                 font_size=12, color=C_DARK)
        y3 += Inches(0.44)


def slide_checklist_phase4(prs, blank_layout):
    """5. 移行チェックリスト — Phase 4: 動作確認・本番切替"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "移行チェックリスト — Phase 4: 動作確認・本番切替", "Section 5")

    phase_badge(slide, Inches(0.5), Inches(1.0), "Phase 4", C_RED)
    add_text(slide, "本番切替前の最終確認と DNS 切替手順",
             Inches(2.9), Inches(1.05), Inches(9.5), Inches(0.35),
             font_size=13, color=C_MUTED)

    items_left = [
        ("予約フロー End-to-End テストを実施する",           "todo", "患者側: 日程選択→入力→完了"),
        ("管理者ログイン・ダッシュボード確認",               "todo", "予約一覧・詳細モーダル・CSV出力"),
        ("メール送信テスト（確認/リマインダー/キャンセル）", "todo", "実際のメールアドレスで受信確認"),
        ("問診票フォームの動作確認",                         "todo", "/questionnaire?bookingId=xxx"),
        ("診察完了・診察履歴の記録確認",                     "todo", ""),
        ("モバイル（スマホ）での表示確認",                   "todo", "患者側・管理者ナビ両方"),
    ]
    items_right = [
        ("DNS の TTL を短くする（推奨: 300秒）",             "todo", "切替後の反映を早くするため事前に設定"),
        ("旧 URL からのリダイレクト確認",                    "todo", "/login.html → /login など"),
        ("DNS CNAME を新 Firebase Hosting に切替",           "todo", "タイミングは業務時間外を推奨"),
        ("SSL 証明書の自動発行を確認（〜数時間）",           "todo", ""),
        ("切替後 24 時間の動作モニタリング",                  "todo", "Firebase Console > Crashlytics / Logging"),
        ("旧プロジェクトのリソースを停止・削除",             "todo", "移行確認後 1 週間以上経過してから"),
    ]

    checklist_rows(slide, items_left,  Inches(0.5), Inches(1.55), 6.2)
    checklist_rows(slide, items_right, Inches(6.9), Inches(1.55), 6.2)


def slide_risks(prs, blank_layout):
    """6. リスクと対策"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "リスクと対策", "Section 6")

    risks = [
        ("DNS 切替時のダウンタイム",
         "HIGH",
         "TTL を事前に短縮（300秒）。切替は予約の少ない時間帯（深夜）に実施。\n切替前にメンテナンスバナーで患者に事前告知。",
         "DNS を旧設定に戻す（TTL=300秒 なら約5分で復旧）",
         C_RED),
        ("データ移行時の不整合・欠損",
         "HIGH",
         "エクスポート前にメンテナンスモード ON（新規書き込みを防止）。\n移行後は件数・内容の整合性チェックを実施（件数比較）。",
         "旧プロジェクトのバックアップから再エクスポート",
         C_AMBER),
        ("メール送信失敗（Resend 認証未完）",
         "MED",
         "DNS 切替前に Resend のドメイン認証（SPF/DKIM/DMARC）を完全に完了させる。\nテスト送信で確認後に本番切替。",
         "一時的に旧プロジェクトの Resend アカウントを流用",
         C_AMBER),
        ("Firebase Auth UID 不一致",
         "LOW",
         "firebase auth:export/import を使用することで UID は保持される。\nFirestore との参照は UID 固定のため自動的に維持。",
         "手動マッピングは不要（export/import で解消）",
         C_GREEN),
    ]

    y = Inches(1.1)
    for title, level, mitigation, rollback, color in risks:
        card(slide, Inches(0.5), y, Inches(12.3), Inches(1.4))
        add_shape(slide, Inches(0.5), y, Inches(0.12), Inches(1.4), fill_color=color)
        # レベルバッジ
        s = add_shape(slide, Inches(0.75), y + Inches(0.08), Inches(0.9), Inches(0.3),
                      fill_color=color)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = level
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.name = "Noto Sans JP"
        run.font.color.rgb = C_WHITE
        add_text(slide, title, Inches(1.75), y + Inches(0.06), Inches(10.7), Inches(0.38),
                 font_size=15, bold=True, color=C_DARK)
        add_text(slide, f"対策: {mitigation}", Inches(0.75), y + Inches(0.52),
                 Inches(8.3), Inches(0.55), font_size=11, color=C_DARK)
        add_text(slide, f"RB: {rollback}", Inches(9.2), y + Inches(0.52),
                 Inches(3.4), Inches(0.55), font_size=10, color=C_MUTED)
        y += Inches(1.52)


def slide_data_migration(prs, blank_layout):
    """7. データ移行方針"""
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_BASE)
    header_bar(slide, "データ移行方針", "Section 7")

    # コレクション別移行方針
    cols = [
        ("reservations",      "予約データ",     "移行",   "完全移行（全フィールド）。移行前に status='cancelled' のものを除外するかクライアントと合意。"),
        ("slots",             "時間枠",         "移行",   "完全移行。ただし過去分（date < 移行日）は不要な場合は除外 OK。"),
        ("users",             "管理者アカウント", "再作成", "Auth 移行後に Firestore ドキュメントも連動。UID は保持されるため参照は維持。"),
        ("visit_histories",   "診察履歴",       "移行",   "完全移行（改ざん防止のため全件）。APPI の保存義務あり。"),
        ("questionnaires",    "問診票",         "移行",   "完全移行。bookingId で予約と紐付けられているため整合性確認。"),
        ("access_logs",       "アクセスログ",   "任意",   "監査目的のため移行推奨。ただし容量が大きい場合は移行日以降のみでも可。"),
        ("settings",          "院の設定情報",   "移行",   "clinic / terms / announcement ドキュメントを完全移行。"),
    ]

    add_text(slide, "コレクション別 移行方針",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
             font_size=15, bold=True, color=C_DARK)

    # ヘッダー行
    hy = Inches(1.5)
    for label, x, w in [("コレクション", 0.5, 2.5), ("用途", 3.1, 2.2),
                         ("方針", 5.4, 1.2), ("注記", 6.8, 6.0)]:
        add_shape(slide, Inches(x), hy, Inches(w), Inches(0.34), fill_color=C_NAVY)
        add_text(slide, label, Inches(x + 0.1), hy + Inches(0.04), Inches(w - 0.2), Inches(0.3),
                 font_size=11, bold=True, color=C_WHITE)

    y = Inches(1.84)
    for i, (col, usage, policy, note) in enumerate(cols):
        bg = C_BASE if i % 2 == 0 else C_WHITE
        policy_color = C_GREEN if policy == "移行" else (C_AMBER if policy == "任意" else C_NAVY)
        add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.62), fill_color=bg)
        add_text(slide, col, Inches(0.6), y + Inches(0.12), Inches(2.4), Inches(0.38),
                 font_size=12, bold=True, color=C_DARK)
        add_text(slide, usage, Inches(3.2), y + Inches(0.12), Inches(2.1), Inches(0.38),
                 font_size=12, color=C_DARK)
        s = add_shape(slide, Inches(5.5), y + Inches(0.12), Inches(1.0), Inches(0.32),
                      fill_color=policy_color)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = policy
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.name = "Noto Sans JP"
        run.font.color.rgb = C_WHITE
        add_text(slide, note, Inches(6.8), y + Inches(0.1), Inches(5.9), Inches(0.42),
                 font_size=10, color=C_MUTED)
        y += Inches(0.65)


def slide_closing(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, Inches(0), Inches(0), W, H, fill_color=C_NAVY)
    add_shape(slide, Inches(0), H - Inches(1.2), W, Inches(1.2), fill_color=C_GOLD)
    add_shape(slide, Inches(0.5), Inches(2.5), Inches(12.33), Inches(0.06), fill_color=C_GOLD)
    add_text(slide, "移行作業における注意事項",
             Inches(0.5), Inches(1.1), Inches(12.0), Inches(0.7),
             font_size=28, bold=True, color=C_WHITE)
    notes = [
        "移行作業は必ず予約の少ない時間帯（深夜〜早朝）に実施してください。",
        "各 Phase の作業前後に Firestore と Auth のバックアップ（エクスポート）を取得してください。",
        "DNS 切替後 48 時間は旧プロジェクトを停止せず、ロールバック可能な状態を維持してください。",
        "本資料に記載のコマンドは新プロジェクト ID に読み替えて実行してください。",
    ]
    y = Inches(2.75)
    for note in notes:
        add_text(slide, f"▶  {note}", Inches(0.7), y, Inches(11.9), Inches(0.45),
                 font_size=14, color=C_BASE)
        y += Inches(0.55)
    add_text(slide, "ご不明な点は担当エンジニアまでお問い合わせください。",
             Inches(0.7), H - Inches(1.1), Inches(11.9), Inches(0.45),
             font_size=14, bold=True, color=C_NAVY)


# ─────────────────────────────────────────────
# メイン
# ─────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # 空白レイアウト

    print("スライド生成中...")

    slide_cover(prs, blank_layout)
    print("  ✓ カバー")

    slide_toc(prs, blank_layout)
    print("  ✓ 目次")

    divider_slide(prs, blank_layout, "Section 1", "現状環境の整理",
                  ["現在のサービス構成", "ボス個人アカウントで運用中のリソース一覧"])
    slide_current_env(prs, blank_layout)
    print("  ✓ 現状環境")

    divider_slide(prs, blank_layout, "Section 2", "必要な契約・アカウント",
                  ["ドメイン / DNS", "Firebase / GCP", "Resend（メール配信）"])
    slide_required_contracts(prs, blank_layout)
    print("  ✓ 必要な契約")

    divider_slide(prs, blank_layout, "Section 3", "予算見積もり",
                  ["サービス別月額費用", "スケール別コスト試算"])
    slide_budget(prs, blank_layout)
    print("  ✓ 予算見積もり")

    divider_slide(prs, blank_layout, "Section 4", "推奨環境構築手順",
                  ["Firebase プロジェクト作成", "Blaze プラン設定", "Secret Manager 設定"])
    slide_recommended_setup(prs, blank_layout)
    print("  ✓ 推奨環境構築")

    divider_slide(prs, blank_layout, "Section 5", "移行チェックリスト",
                  ["Phase 1: 契約・準備", "Phase 2: コード・設定移行",
                   "Phase 3: データ移行", "Phase 4: 動作確認・本番切替"])
    slide_checklist_phase1(prs, blank_layout)
    slide_checklist_phase2(prs, blank_layout)
    slide_checklist_phase3(prs, blank_layout)
    slide_checklist_phase4(prs, blank_layout)
    print("  ✓ 移行チェックリスト（4 Phase）")

    divider_slide(prs, blank_layout, "Section 6", "リスクと対策",
                  ["DNS 切替ダウンタイム", "データ不整合", "メール認証未完", "Auth UID 問題"])
    slide_risks(prs, blank_layout)
    print("  ✓ リスクと対策")

    divider_slide(prs, blank_layout, "Section 7", "データ移行方針",
                  ["Firestore コレクション別移行方針", "Firebase Auth 移行", "Secrets 移行"])
    slide_data_migration(prs, blank_layout)
    print("  ✓ データ移行方針")

    slide_closing(prs, blank_layout)
    print("  ✓ クロージング")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"\n完了！ → {OUTPUT_PATH}")
    print(f"スライド数: {len(prs.slides)} 枚")


if __name__ == "__main__":
    main()
